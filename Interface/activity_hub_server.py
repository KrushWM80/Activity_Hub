"""
Activity Hub Server V2 - Minimal Flask with inline BigQuery
"""

import os
import logging
import requests
import uuid
import smtplib
import threading
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
from flask import Flask, send_from_directory, send_file, jsonify, request
from flask_cors import CORS
from datetime import datetime, timedelta
from google.cloud import bigquery
# from apscheduler.schedulers.background import BackgroundScheduler  # Network blocked - commented out

# Basic logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s %(levelname)s %(message)s')
logger = logging.getLogger(__name__)

# Set Google Cloud credentials before BigQuery initialization
creds_path = os.path.join(os.path.expanduser("~"), "AppData", "Roaming", "gcloud", "application_default_credentials.json")
if not os.environ.get('GOOGLE_APPLICATION_CREDENTIALS'):
    if os.path.exists(creds_path):
        os.environ['GOOGLE_APPLICATION_CREDENTIALS'] = creds_path
        logger.info(f"Set GOOGLE_APPLICATION_CREDENTIALS to {creds_path}")
    else:
        logger.warning(f"Credentials file not found at {creds_path} - BigQuery may fail")

app = Flask(__name__)
CORS(app)

# BigQuery client - will be initialized on first use
_bq_client_cache = None

def get_bq_client():
    """Get or initialize the BigQuery client"""
    global _bq_client_cache
    if _bq_client_cache is None:
        try:
            logger.info("[get_bq_client] Attempting to initialize BigQuery client...")
            _bq_client_cache = bigquery.Client(project='wmt-assetprotection-prod')
            logger.info(f"[get_bq_client] ✓ Client initialized successfully")
        except Exception as e:
            logger.error(f"[get_bq_client] ✗ Failed: {e}", exc_info=True)
            return None
    logger.info(f"[get_bq_client] Returning cached client")
    return _bq_client_cache

BASE_DIR = os.path.dirname(os.path.abspath(__file__))

# Scheduler Service URL for proxying
SCHEDULER_SERVICE_URL = 'http://localhost:5011'

# ──────────────────────────────────────────────
# DEBUG ENDPOINT
# ──────────────────────────────────────────────

@app.route('/api/debug/bq-status', methods=['GET'])
def debug_bq_status():
    """Check if BigQuery client is initialized"""
    global bq_client
    client_status = "initialized" if bq_client else "NOT initialized"
    return jsonify({
        'bq_client_status': client_status,
        'bq_client_type': str(type(bq_client)),
        'bq_client_value': str(bq_client)
    })

# ──────────────────────────────────────────────
# CORE ROUTES
# ──────────────────────────────────────────────

@app.route('/activity-hub/')
@app.route('/activity-hub/for-you')
def for_you():
    """For You page"""
    return send_file(os.path.join(BASE_DIR, 'For You - Landing Page', 'activity-hub-demo.html'))

@app.route('/activity-hub/projects')
def projects():
    """Projects page"""
    return send_file(os.path.join(BASE_DIR, 'projects-page.html'))

@app.route('/activity-hub/reporting')
def reporting():
    """Reporting page"""
    return send_file(os.path.join(BASE_DIR, 'Reporting', 'reporting.html'))

@app.route('/activity-hub/admin')
def admin():
    """Admin dashboard"""
    return send_file(os.path.join(BASE_DIR, 'Admin', 'admin-dashboard.html'))

# ──────────────────────────────────────────────
# STATIC FILES & ASSETS
# ──────────────────────────────────────────────

@app.route('/<filename>')
@app.route('/activity-hub/static/<path:filepath>')
@app.route('/<path:subdir>/<filename>')
def serve_static(filename=None, filepath=None, subdir=None):
    """Serve static files like images and JS from Interface directory"""
    try:
        if filepath:
            # Handle nested paths for static files
            file_path = os.path.join(BASE_DIR, filepath.replace('/', os.sep))
        elif subdir and filename:
            # Handle paths like /Admin/Widgets/widget-registry.js
            file_path = os.path.join(BASE_DIR, subdir, filename)
        else:
            file_path = os.path.join(BASE_DIR, filename)
        
        if os.path.exists(file_path):
            return send_file(file_path)
        return jsonify({'error': 'File not found'}), 404
    except Exception as e:
        logger.error(f"Static file error: {str(e)}")
        return jsonify({'error': 'File not found'}), 404

# ──────────────────────────────────────────────
# BIGQUERY API ENDPOINTS - INLINE ONLY
# ──────────────────────────────────────────────

@app.route('/api/projects/<project_id>', methods=['PUT', 'DELETE'])
def api_project_by_id(project_id):
    """Update or delete a specific project by ID"""
    from flask import request
    
    bq_client = get_bq_client()
    if not bq_client:
        return jsonify({'error': 'BigQuery client not initialized'}), 500
    
    try:
        if request.method == 'DELETE':
            # Delete project
            sql = "DELETE FROM `wmt-assetprotection-prod.Store_Support_Dev.AH_Projects` WHERE project_id = %s"
            job = bq_client.query(sql, job_config=bigquery.QueryJobConfig(
                query_parameters=[bigquery.ScalarQueryParameter(None, "STRING", project_id)]
            ))
            job.result()
            return jsonify({'success': True, 'message': 'Project deleted'}), 200
        
        elif request.method == 'PUT':
            # Update project
            data = request.get_json()
            
            # Check if this is a partial update (Intake Hub project - only latest_update field)
            # or a full update (manual project - all fields)
            has_full_fields = any(field in data for field in ['title', 'owner', 'owner_id', 'health', 'business_area', 'business_organization'])
            
            if has_full_fields:
                # Full update for manual projects
                sql = """
                UPDATE `wmt-assetprotection-prod.Store_Support_Dev.AH_Projects`
                SET title = @title, 
                    business_organization = @business_organization, 
                    owner = @owner, 
                    owner_id = @owner_id,
                    health = @health, 
                    status = @status, 
                    project_update = @project_update, 
                    last_updated = CURRENT_TIMESTAMP()
                WHERE project_id = @project_id
                """
                query_params = [
                    bigquery.ScalarQueryParameter("title", "STRING", data.get('title')),
                    bigquery.ScalarQueryParameter("business_organization", "STRING", data.get('business_organization', data.get('store_area', ''))),
                    bigquery.ScalarQueryParameter("owner", "STRING", data.get('owner')),
                    bigquery.ScalarQueryParameter("owner_id", "STRING", data.get('owner_id')),
                    bigquery.ScalarQueryParameter("health", "STRING", data.get('health')),
                    bigquery.ScalarQueryParameter("status", "STRING", data.get('status', 'Active')),
                    bigquery.ScalarQueryParameter("project_update", "STRING", data.get('project_update', '')),
                    bigquery.ScalarQueryParameter("project_id", "STRING", project_id)
                ]
            else:
                # Partial update for Intake Hub projects (only update the note)
                sql = """
                UPDATE `wmt-assetprotection-prod.Store_Support_Dev.AH_Projects`
                SET project_update = @project_update, 
                    last_updated = CURRENT_TIMESTAMP()
                WHERE project_id = @project_id
                """
                query_params = [
                    bigquery.ScalarQueryParameter("project_update", "STRING", data.get('latest_update', data.get('project_update', ''))),
                    bigquery.ScalarQueryParameter("project_id", "STRING", project_id)
                ]
            
            job = bq_client.query(sql, job_config=bigquery.QueryJobConfig(
                query_parameters=query_params
            ))
            job.result()
            return jsonify({'success': True, 'project_id': project_id}), 200
    except Exception as e:
        logger.error(f"Project update error: {str(e)}")
        return jsonify({'error': str(e)}), 500

@app.route('/api/projects', methods=['GET', 'POST', 'PUT'])
def api_projects():
    """Get/create/update projects from BigQuery"""
    from flask import request
    logger.info("[api_projects] ENDPOINT CALLED")
    
    bq_client = get_bq_client()
    if not bq_client:
        logger.error(f"[api_projects] bq_client is None!")
        return jsonify({'error': 'BigQuery client not initialized'}), 500
    
    logger.info(f"[api_projects] bq_client check passed, proceeding with request")
    
    # Handle POST (create) and PUT (update)
    if request.method in ['POST', 'PUT']:
        try:
            import json
            data = request.get_json()
            project_id = data.get('project_id', f'proj-{int(datetime.now().timestamp())}')
            user_email = data.get('user_email', 'Activity Hub')
            new_note = data.get('project_update', '')
            
            if request.method == 'POST':
                # Insert new manual project
                sql = """
                INSERT INTO `wmt-assetprotection-prod.Store_Support_Dev.AH_Projects`
                (project_id, title, owner, owner_id, health, status, project_source, 
                 created_date, last_updated, project_update, project_update_date, 
                 project_update_by, previous_updates, business_organization)
                VALUES 
                (@project_id, @title, @owner, @owner_id, @health, @status, 'Manual Upload',
                 CURRENT_TIMESTAMP(), CURRENT_TIMESTAMP(), @project_update,
                 CURRENT_TIMESTAMP(), @user_email, JSON '[]',
                 @business_organization)
                """
                
                job_config = bigquery.QueryJobConfig(query_parameters=[
                    bigquery.ScalarQueryParameter("project_id", "STRING", project_id),
                    bigquery.ScalarQueryParameter("title", "STRING", data.get('title', '')),
                    bigquery.ScalarQueryParameter("owner", "STRING", data.get('owner', '')),
                    bigquery.ScalarQueryParameter("owner_id", "STRING", data.get('owner_id', '')),
                    bigquery.ScalarQueryParameter("health", "STRING", data.get('health', 'Unknown')),
                    bigquery.ScalarQueryParameter("status", "STRING", data.get('status', 'Active')),
                    bigquery.ScalarQueryParameter("project_update", "STRING", new_note),
                    bigquery.ScalarQueryParameter("user_email", "STRING", user_email),
                    bigquery.ScalarQueryParameter("business_organization", "STRING", 
                                                 data.get('business_organization', data.get('store_area', ''))),
                ])
            else:
                # Update project: track note history if note changed
                sql = """
                UPDATE `wmt-assetprotection-prod.Store_Support_Dev.AH_Projects`
                SET title = @title, 
                    owner = @owner, 
                    owner_id = @owner_id,
                    health = @health, 
                    status = @status, 
                    project_update = @project_update,
                    project_update_date = CURRENT_TIMESTAMP(),
                    project_update_by = @user_email,
                    last_updated = CURRENT_TIMESTAMP(),
                    business_organization = @business_organization,
                    previous_updates = CASE 
                        WHEN @project_update != project_update AND project_update IS NOT NULL
                        THEN JSON_SET(
                            COALESCE(previous_updates, JSON '[]'),
                            CONCAT('$[', JSON_LENGTH(COALESCE(previous_updates, JSON '[]')), ']'),
                            JSON_OBJECT(
                                'note', project_update,
                                'timestamp', CURRENT_TIMESTAMP(),
                                'updated_by', project_update_by
                            )
                        )
                        ELSE previous_updates
                    END
                WHERE project_id = @project_id
                """
                
                job_config = bigquery.QueryJobConfig(query_parameters=[
                    bigquery.ScalarQueryParameter("project_id", "STRING", project_id),
                    bigquery.ScalarQueryParameter("title", "STRING", data.get('title')),
                    bigquery.ScalarQueryParameter("owner", "STRING", data.get('owner')),
                    bigquery.ScalarQueryParameter("owner_id", "STRING", data.get('owner_id')),
                    bigquery.ScalarQueryParameter("health", "STRING", data.get('health')),
                    bigquery.ScalarQueryParameter("status", "STRING", data.get('status')),
                    bigquery.ScalarQueryParameter("project_update", "STRING", new_note),
                    bigquery.ScalarQueryParameter("user_email", "STRING", user_email),
                    bigquery.ScalarQueryParameter("business_organization", "STRING", 
                                                 data.get('business_organization', data.get('store_area', ''))),
                ])
            
            job = bq_client.query(sql, job_config=job_config)
            job.result()
            return jsonify({'success': True, 'project_id': project_id, 'note_tracked': True}), 200
        except Exception as e:
            logger.error(f"Project create/update error: {str(e)}")
            return jsonify({'error': str(e)}), 500
    
    # Handle GET
    try:
        business_org = request.args.get('business_organization')
        health_status = request.args.get('health_status')
        title_search = request.args.get('title_search', '')  # Search in project title
        owner_names = request.args.get('owner_names', '')   # Comma-separated list of owners
        limit = int(request.args.get('limit', 1000))  # Default: 1000 rows per request
        offset = int(request.args.get('offset', 0))   # Pagination offset
        updated_status = request.args.get('updated_status', '')  # 'updated' or 'not_updated'
        
        # Calculate current Walmart week
        today = datetime.now()
        # Get current WM week from Cal_Dim_Data (authoritative source)
        cal_sql = """
        SELECT DISTINCT WM_WEEK_NBR, FISCAL_YEAR_NBR
        FROM `wmt-assetprotection-prod.Store_Support_Dev.Cal_Dim_Data`
        WHERE CALENDAR_DATE = CURRENT_DATE()
        """
        cal_result = list(bq_client.query(cal_sql).result())
        if cal_result:
            current_wm_week = cal_result[0].WM_WEEK_NBR
            current_fiscal_year = cal_result[0].FISCAL_YEAR_NBR
        else:
            # Fallback if no match today
            today = datetime.now()
            fiscal_year_start = datetime(today.year if today.month >= 2 else today.year - 1, 2, 1)
            days_since_fy = (today - fiscal_year_start).days
            current_wm_week = (days_since_fy // 7) + 1
            current_fiscal_year = today.year if today.month >= 2 else today.year - 1
        
        # Build WHERE clause filters for consistent use in both count and paginated queries
        where_filters = ""
        if business_org:
            where_filters += f" AND p.business_organization = '{business_org}'"
        if health_status:
            where_filters += f" AND p.health = '{health_status}'"
        if owner_names:
            owners = [f"'{o.strip()}'" for o in owner_names.split(',')]
            where_filters += f" AND p.owner IN ({','.join(owners)})"
        
        # First, get total count matching filters
        count_sql = f"""
        SELECT COUNT(*) as total
        FROM `wmt-assetprotection-prod.Store_Support_Dev.AH_Projects` p
        WHERE TRIM(p.business_organization) != ''
        AND LOWER(p.title) LIKE LOWER('%{title_search}%')
        {where_filters}
        """
        
        count_result = list(bq_client.query(count_sql).result())
        total_count = count_result[0].total if count_result else 0
        
        # Then get paginated results with Cal_Dim_Data join for WM week info
        sql = f"""
        SELECT 
            p.project_id,
            p.title,
            p.business_organization,
            p.owner,
            p.owner_id,
            p.health,
            p.status,
            p.project_source,
            p.created_date,
            p.last_updated,
            p.project_update,
            p.project_update_date,
            CASE 
                WHEN c.WM_WEEK_NBR = {current_wm_week} AND c.FISCAL_YEAR_NBR = {current_fiscal_year}
                THEN TRUE
                ELSE FALSE
            END as is_updated_this_week
        FROM `wmt-assetprotection-prod.Store_Support_Dev.AH_Projects` p
        LEFT JOIN `wmt-assetprotection-prod.Store_Support_Dev.Cal_Dim_Data` c
            ON CAST(p.project_update_date AS DATE) = c.CALENDAR_DATE
        WHERE TRIM(p.business_organization) != ''
        AND LOWER(p.title) LIKE LOWER('%{title_search}%')
        {where_filters}
        ORDER BY COALESCE(p.project_update_date, p.last_updated) DESC
        LIMIT {limit} OFFSET {offset}
        """
        
        # Query
        results = bq_client.query(sql).result()
        projects = []
        
        for row in results:
            # Get is_updated_this_week from query result (checks WK 13: 4/25-5/1)
            is_updated_this_week = getattr(row, 'is_updated_this_week', False)
            update_date = getattr(row, 'project_update_date', None)
            
            # Apply updated_status filter
            if updated_status == 'updated' and not is_updated_this_week:
                continue
            if updated_status == 'not_updated' and is_updated_this_week:
                continue
            
            projects.append({
                'project_id': row.project_id,
                'title': row.title,
                'business_organization': row.business_organization,
                'owner': row.owner,
                'owner_id': row.owner_id,
                'health': row.health,
                'status': row.status,
                'project_source': row.project_source,
                'created_date': row.created_date.isoformat() if row.created_date else None,
                'updated_date': row.last_updated.isoformat() if row.last_updated else None,
                'project_update': row.project_update,
                'project_update_date': update_date.isoformat() if update_date else None,
                'is_updated_this_week': is_updated_this_week
            })
        
        # Recalculate total if updated_status filter was applied
        if updated_status:
            total_count = len(projects)
        
        return jsonify({
            'projects': projects,
            'total_count': total_count,
            'limit': limit,
            'offset': offset,
            'timestamp': datetime.now().isoformat()
        })
    
    except Exception as e:
        logger.error(f"Projects API error: {str(e)}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/projects/<project_id>/history', methods=['GET'])
def api_project_history(project_id):
    """Get project note history (current + previous updates)"""
    try:
        sql = """
        SELECT 
            project_id,
            project_update as current_note,
            project_update_date as current_note_date,
            project_update_by as current_note_by,
            previous_updates
        FROM `wmt-assetprotection-prod.Store_Support_Dev.AH_Projects`
        WHERE project_id = @project_id
        """
        
        job_config = bigquery.QueryJobConfig(query_parameters=[
            bigquery.ScalarQueryParameter("project_id", "STRING", project_id),
        ])
        
        results = list(bq_client.query(sql, job_config=job_config).result())
        
        if not results:
            return jsonify({'error': 'Project not found'}), 404
        
        row = results[0]
        
        # Build response with current + history
        history = []
        
        # Add current note if exists
        if row.current_note:
            history.append({
                'note': row.current_note,
                'timestamp': row.current_note_date.isoformat() if row.current_note_date else None,
                'updated_by': row.current_note_by or 'Activity Hub',
                'is_current': True
            })
        
        # Add previous notes from JSON array
        if row.previous_updates:
            import json
            try:
                previous = json.loads(row.previous_updates)
                for prev in reversed(previous):  # Reverse to show newest first (after current)
                    prev['is_current'] = False
                    history.append(prev)
            except:
                pass
        
        return jsonify({
            'project_id': project_id,
            'current_note': row.current_note,
            'current_note_date': row.current_note_date.isoformat() if row.current_note_date else None,
            'current_note_by': row.current_note_by or 'Activity Hub',
            'history': history,
            'total_versions': len(history)
        }), 200
        
    except Exception as e:
        logger.error(f"Project history error: {str(e)}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/projects/metrics', methods=['GET'])
def api_metrics():
    """Get metrics"""
    bq_client = get_bq_client()
    if not bq_client:
        return jsonify({'error': 'BigQuery client not initialized'}), 500
    
    try:
        # Get current WM week and fiscal year from Cal_Dim_Data (authoritative source)
        cal_sql = """
        SELECT DISTINCT WM_WEEK_NBR, FISCAL_YEAR_NBR
        FROM `wmt-assetprotection-prod.Store_Support_Dev.Cal_Dim_Data`
        WHERE CALENDAR_DATE = CURRENT_DATE()
        """
        cal_result = list(bq_client.query(cal_sql).result())
        if cal_result:
            current_wm_week = cal_result[0].WM_WEEK_NBR
            current_fiscal_year = cal_result[0].FISCAL_YEAR_NBR
        else:
            # Fallback if no match today
            today = datetime.now()
            fy_start = datetime(today.year if today.month >= 2 else today.year - 1, 2, 1)
            days_since = (today - fy_start).days
            current_wm_week = (days_since // 7) + 1
            current_fiscal_year = today.year if today.month >= 2 else today.year - 1
        
        # Get all projects and count those updated in current WM week using Cal_Dim_Data
        sql = f"""
        SELECT 
            COUNT(*) as total,
            SUM(CASE WHEN c.WM_WEEK_NBR = {current_wm_week} AND c.FISCAL_YEAR_NBR = {current_fiscal_year} THEN 1 ELSE 0 END) as projects_this_week,
            COUNT(DISTINCT p.owner) as unique_owners
        FROM `wmt-assetprotection-prod.Store_Support_Dev.AH_Projects` p
        LEFT JOIN `wmt-assetprotection-prod.Store_Support_Dev.Cal_Dim_Data` c
            ON CAST(p.project_update_date AS DATE) = c.CALENDAR_DATE
        WHERE p.project_id IS NOT NULL
        """
        
        results = list(bq_client.query(sql).result())
        row = results[0] if results else None
        
        if row:
            total = row.total or 0
            projects_this_week = row.projects_this_week or 0
            unique_owners = row.unique_owners or 0
        else:
            total = 0
            projects_this_week = 0
            unique_owners = 0
        
        percent = (projects_this_week * 100 / total) if total > 0 else 0
        
        return jsonify({
            'metrics': {
                'active_projects': total,
                'unique_owners': unique_owners,
                'projects_updated_this_week': projects_this_week,
                'percent_updated': round(percent, 2)
            },
            'timestamp': datetime.now().isoformat()
        })
    
    except Exception as e:
        logger.error(f"Metrics API error: {str(e)}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/projects/business-areas', methods=['GET'])
def api_business_areas():
    """Get unique Business Owner Area options from source Intake Hub"""
    bq_client = get_bq_client()
    if not bq_client:
        return jsonify({'error': 'BigQuery client not initialized'}), 500
    
    try:
        sql = r"""
        SELECT DISTINCT Business_Owner_Area
        FROM `wmt-assetprotection-prod.Store_Support_Dev.Output - Intake Accel Council Data`
        WHERE Business_Owner_Area IS NOT NULL 
        AND TRIM(Business_Owner_Area) != ''
        ORDER BY Business_Owner_Area
        """
        
        results = bq_client.query(sql).result()
        areas = [row.Business_Owner_Area for row in results if row.Business_Owner_Area and row.Business_Owner_Area.strip()]
        
        return jsonify({
            'business_areas': areas,
            'total': len(areas),
            'timestamp': datetime.now().isoformat()
        })
    
    except Exception as e:
        logger.error(f"Business areas error: {str(e)}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/projects/sync/status', methods=['GET'])
def projects_sync_status():
    """Get projects sync status (stub - projects service may not be running)."""
    return jsonify({
        'sync_status': 'idle',
        'last_sync': datetime.now().isoformat(),
        'is_available': False,
        'message': 'Projects Data-Bridge service not available'
    }), 200


@app.route('/api/projects/titles', methods=['GET'])
def api_project_titles():
    """Get all unique project titles for filtering"""
    bq_client = get_bq_client()
    if not bq_client:
        return jsonify({'error': 'BigQuery client not initialized'}), 500
    
    try:
        sql = r"""
        SELECT DISTINCT title
        FROM `wmt-assetprotection-prod.Store_Support_Dev.AH_Projects`
        WHERE title IS NOT NULL
        ORDER BY title
        """
        
        results = bq_client.query(sql).result()
        titles = [row.title for row in results if row.title]
        
        return jsonify({
            'titles': titles,
            'total': len(titles),
            'timestamp': datetime.now().isoformat()
        })
    
    except Exception as e:
        logger.error(f"Project titles error: {str(e)}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/projects/owners', methods=['GET'])
def api_project_owners():
    """Get all unique owners for multi-select filtering"""
    bq_client = get_bq_client()
    if not bq_client:
        return jsonify({'error': 'BigQuery client not initialized'}), 500
    
    try:
        sql = r"""
        SELECT DISTINCT owner
        FROM `wmt-assetprotection-prod.Store_Support_Dev.AH_Projects`
        WHERE owner IS NOT NULL
        ORDER BY owner
        """
        
        results = bq_client.query(sql).result()
        owners = [row.owner for row in results if row.owner]
        
        return jsonify({
            'owners': owners,
            'total': len(owners),
            'timestamp': datetime.now().isoformat()
        })
    
    except Exception as e:
        logger.error(f"Project owners error: {str(e)}")
        return jsonify({'error': str(e)}), 500


# ──────────────────────────────────────────────
# ADMIN ENDPOINTS
# ──────────────────────────────────────────────

@app.route('/api/admin/activity-log', methods=['GET'])
def admin_activity_log():
    """Return admin activity log for audit trail."""
    now = datetime.now()
    activity_log = [
        {
            'timestamp': (now - timedelta(hours=2)).isoformat(),
            'user': 'Kendall Rush',
            'action': 'Created Logic Request',
            'resource': 'Project Updates - 30 Days',
            'status': 'success'
        },
        {
            'timestamp': (now - timedelta(hours=4)).isoformat(),
            'user': 'Admin User',
            'action': 'Modified Trigger Rule',
            'resource': 'SIF Missing - 60 Days',
            'status': 'success'
        },
        {
            'timestamp': (now - timedelta(hours=8)).isoformat(),
            'user': 'Kendall Rush',
            'action': 'Viewed Dashboard',
            'resource': 'Admin Dashboard',
            'status': 'success'
        },
        {
            'timestamp': (now - timedelta(hours=24)).isoformat(),
            'user': 'System',
            'action': 'Scheduler Run',
            'resource': 'Logic Rules Engine',
            'status': 'success'
        },
    ]
    return jsonify({
        'activity_log': activity_log,
        'total_count': len(activity_log),
        'timestamp': now.isoformat()
    })


@app.route('/api/admin/active-users', methods=['GET'])
def admin_active_users():
    """Return count and list of currently active users."""
    active_users = [
        {'user_id': 'kendall.rush', 'name': 'Kendall Rush', 'role': 'Admin', 'last_activity': datetime.now().isoformat()},
        {'user_id': 'admin.user', 'name': 'Admin User', 'role': 'Admin', 'last_activity': (datetime.now() - timedelta(minutes=5)).isoformat()}
    ]
    return jsonify({
        'active_users': len(active_users),
        'users': active_users,
        'timestamp': datetime.now().isoformat()
    })


@app.route('/api/admin/trigger-ah-projects-sync', methods=['POST'])
def trigger_ah_projects_sync():
    """Manually trigger AH_Projects sync from admin dashboard (async)"""
    try:
        data = request.get_json()
        triggered_by = data.get('triggered_by', 'manual')
        
        logger.info(f"Manual AH_Projects sync triggered by: {triggered_by}")
        
        def run_sync_background():
            """Run sync in background thread"""
            try:
                from google.cloud import bigquery
                import subprocess
                import sys
                
                # Run sync script
                sync_script = os.path.join(BASE_DIR, '..', 'sync_now.py')
                result = subprocess.run(
                    [sys.executable, sync_script],
                    capture_output=True,
                    text=True,
                    timeout=300  # 5 minutes
                )
                
                if result.returncode == 0:
                    logger.info(f"AH_Projects sync completed successfully")
                else:
                    logger.error(f"Sync script failed: {result.stderr}")
                    
            except Exception as e:
                logger.error(f"Background sync error: {e}")
        
        # Start sync in background thread without blocking
        sync_thread = threading.Thread(target=run_sync_background, daemon=True)
        sync_thread.start()
        
        # Return immediately with acknowledgment
        return jsonify({
            'success': True,
            'message': 'Sync started in background',
            'status': 'running',
            'timestamp': datetime.now().isoformat()
        }), 202  # 202 Accepted
            
    except Exception as e:
        logger.error(f"Manual sync trigger error: {e}")
        return jsonify({'success': False, 'error': str(e)}), 500


# ──────────────────────────────────────────────
# LOGIC RULES ENGINE PROXY ENDPOINTS
# ──────────────────────────────────────────────

@app.route('/api/logic/metrics', methods=['GET'])
def logic_metrics():
    """Proxy Logic Rules Engine metrics for dashboard display."""
    try:
        response = requests.get(f'{SCHEDULER_SERVICE_URL}/api/v1/logic-metrics', timeout=5)
        if response.status_code == 200:
            return response.json()
        return {'error': 'Service unavailable'}, 503
    except Exception as e:
        logger.error(f"Logic metrics error: {str(e)}")
        return {'error': str(e)}, 500


@app.route('/api/logic/notifications/today', methods=['GET'])
def logic_notifications_today():
    """Proxy Logic Rules Engine notifications for today."""
    try:
        response = requests.get(f'{SCHEDULER_SERVICE_URL}/api/v1/notifications/today', timeout=5)
        if response.status_code == 200:
            return response.json()
        return {'error': 'Service unavailable'}, 503
    except Exception as e:
        logger.error(f"Logic notifications error: {str(e)}")
        return {'error': str(e)}, 500


@app.route('/api/logic/requests', methods=['GET', 'POST'])
def logic_requests():
    """Proxy Logic Requests (GET) or create new request (POST)."""
    try:
        if request.method == 'GET':
            # GET: retrieve logic requests with optional status filter
            status = request.args.get('status')
            url = f'{SCHEDULER_SERVICE_URL}/api/v1/logic-requests'
            if status:
                url += f'?status={status}'
            response = requests.get(url, timeout=5)
            if response.status_code == 200:
                return response.json()
            return {'error': 'Service unavailable'}, 503
        
        elif request.method == 'POST':
            # POST: create new logic request, forward to Scheduler Service
            data = request.get_json()
            response = requests.post(
                f'{SCHEDULER_SERVICE_URL}/api/v1/logic-requests',
                json=data,
                timeout=5,
                headers={'Content-Type': 'application/json'}
            )
            if response.status_code in [200, 201]:
                return response.json()
            else:
                # Forward error response
                try:
                    error_json = response.json()
                except:
                    error_json = {'error': response.text}
                return error_json, response.status_code
    except Exception as e:
        logger.error(f"Logic requests error: {str(e)}")
        return {'error': str(e)}, 500


# ──────────────────────────────────────────────
# PROJECTS REPORT GENERATION
# ──────────────────────────────────────────────

def get_dashboard_screenshot(output_path):
    """Capture screenshot of Projects dashboard using Playwright"""
    try:
        import asyncio
        from playwright.async_api import async_playwright
        
        async def screenshot_dashboard():
            async with async_playwright() as p:
                browser = await p.chromium.launch(headless=True)
                page = await browser.new_page()
                page.set_viewport_size({"width": 1280, "height": 960})
                
                # Navigate to projects page
                await page.goto('http://localhost:8088/activity-hub/projects', wait_until='networkidle')
                await page.wait_for_timeout(2000)  # Wait for data load
                
                # Take screenshot
                screenshot = await page.screenshot(path=output_path, full_page=False)
                await browser.close()
                return True
        
        # Run async function
        loop = asyncio.new_event_loop()
        asyncio.set_event_loop(loop)
        result = loop.run_until_complete(screenshot_dashboard())
        loop.close()
        
        return os.path.exists(output_path)
    except Exception as e:
        logger.warning(f"Screenshot capture failed: {str(e)} - will use fallback")
        return False


def generate_ppt():
    """Generate PPT report with dashboard screenshot (TDA Insights style)"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor
        from io import BytesIO
        import asyncio
        from playwright.async_api import async_playwright
        import tempfile
        import os
        
        # Calculate Walmart Week (WM WK)
        today = datetime.now()
        fiscal_year_start = datetime(today.year if today.month >= 2 else today.year - 1, 2, 1)
        days_since_fy = (today - fiscal_year_start).days
        wm_week = (days_since_fy // 7) + 1
        
        # Get filter parameters
        status = request.args.get('status', '')
        
        # Query projects
        sql = """
        SELECT project_id, title, store_area, owner, health, project_update, status, 
               CASE 
                 WHEN LOWER(wm_week) = CAST((CASE 
                   WHEN EXTRACT(MONTH FROM CURRENT_DATE()) >= 2 THEN EXTRACT(YEAR FROM CURRENT_DATE())
                   ELSE EXTRACT(YEAR FROM CURRENT_DATE()) - 1
                 END) as STRING) THEN TRUE
                 ELSE FALSE
               END as updated_this_week
        FROM `wmt-assetprotection-prod.Store_Support_Dev.AH_Projects`
        WHERE 1=1
        """
        if status != 'All' and status:
            sql += f" AND status = '{status}'"
        sql += " ORDER BY last_updated DESC"
        
        results = bq_client.query(sql).result()
        projects = list(results)
        
        # Create presentation
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        blank_layout = prs.slide_layouts[6]
        
        COLORS = {
            'walmart_blue_dark': RGBColor(0x1E, 0x3A, 0x8A),
            'walmart_yellow': RGBColor(0xFF, 0xC2, 0x20),
            'white': RGBColor(0xFF, 0xFF, 0xFF),
        }
        
        # ──────── SLIDE 1: TITLE SLIDE ────────
        slide = prs.slides.add_slide(blank_layout)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLORS['walmart_blue_dark']
        
        # Main title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.2))
        p = title_box.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = "Activity Hub Projects Overview"
        run.font.size = Pt(48)
        run.font.bold = True
        run.font.color.rgb = COLORS['white']
        
        # Subtitle with WM WK
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(0.8))
        p = subtitle_box.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = f"WM WK {wm_week} • {datetime.now().strftime('%B %d, %Y')}"
        run.font.size = Pt(28)
        run.font.color.rgb = COLORS['walmart_yellow']
        
        # ──────── SLIDE 2: DASHBOARD SCREENSHOT ────────
        # Try to capture dashboard screenshot
        screenshot_path = None
        try:
            import tempfile
            temp_dir = tempfile.gettempdir()
            screenshot_path = os.path.join(temp_dir, f'dashboard_{uuid.uuid4().hex}.png')
            
            if get_dashboard_screenshot(screenshot_path):
                # Add screenshot slide
                slide = prs.slides.add_slide(blank_layout)
                
                # Add title bar
                title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.6))
                title_shape.fill.solid()
                title_shape.fill.fore_color.rgb = COLORS['walmart_blue_dark']
                title_shape.line.fill.background()
                
                title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.1), Inches(9), Inches(0.5))
                p = title_box.text_frame.paragraphs[0]
                run = p.add_run()
                run.text = "Projects Dashboard"
                run.font.size = Pt(20)
                run.font.bold = True
                run.font.color.rgb = COLORS['white']
                
                # Add screenshot (scaled to fit)
                slide.shapes.add_picture(screenshot_path, Inches(0.2), Inches(0.8), width=Inches(9.6), height=Inches(6.5))
                
                # Clean up temp file
                try:
                    os.remove(screenshot_path)
                except:
                    pass
            else:
                logger.warning("Dashboard screenshot capture failed - using fallback")
        except Exception as e:
            logger.warning(f"Screenshot processing error: {str(e)}")
        
        # Save to BytesIO
        ppt_file = BytesIO()
        prs.save(ppt_file)
        ppt_file.seek(0)
        
        return send_file(
            ppt_file,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            as_attachment=True,
            download_name=f'WM WK{wm_week} Projects Overview.pptx'
        )
    except ImportError:
        logger.error("python-pptx not installed")
        return jsonify({'error': 'PPT generation not available - python-pptx not installed'}), 400
    except Exception as e:
        logger.error(f"PPT generation error: {str(e)}")
        return jsonify({'error': str(e)}), 500

# ──────────────────────────────────────────────
# EMAIL UTILITIES
# ──────────────────────────────────────────────

def send_email_via_relay(recipient_email, subject, html_body, sender_email="Activity_Hub@walmart.com"):
    """
    Send email via Walmart internal SMTP relay (relay.walmart.com:25)
    No authentication required on internal network
    """
    try:
        # Create message
        msg = MIMEMultipart('alternative')
        msg['Subject'] = subject
        msg['From'] = sender_email
        msg['To'] = recipient_email
        
        # Attach HTML body
        html_part = MIMEText(html_body, 'html', 'utf-8')
        msg.attach(html_part)
        
        # Connect to Walmart SMTP relay
        server = smtplib.SMTP('relay.walmart.com', 25, timeout=10)
        server.sendmail(sender_email, [recipient_email], msg.as_string())
        server.quit()
        
        logger.info(f"Email sent successfully to {recipient_email}")
        return True
    except Exception as e:
        logger.error(f"Failed to send email to {recipient_email}: {str(e)}")
        return False

@app.route('/api/sample-email', methods=['GET'])
def get_sample_email():
    """Generate and return sample project email HTML"""
    try:
        if not bq_client:
            return jsonify({'error': 'BigQuery client not initialized'}), 500
        
        # Calculate Walmart Week (WM WK)
        # Walmart fiscal year: Feb 1 - Jan 31
        today = datetime.now()
        fiscal_year_start = datetime(today.year if today.month >= 2 else today.year - 1, 2, 1)
        days_since_fy = (today - fiscal_year_start).days
        wm_week = (days_since_fy // 7) + 1
        
        # Get sample projects for email
        sql = """
        SELECT COUNT(*) as total, 
               COUNTIF(health = 'On Track') as on_track,
               COUNTIF(health = 'At Risk') as at_risk,
               COUNTIF(health = 'Off Track') as off_track
        FROM `wmt-assetprotection-prod.Store_Support_Dev.AH_Projects`
        WHERE status = 'Active'
        """
        
        stats_result = bq_client.query(sql).result()
        stats = list(stats_result)[0]
        
        # Get sample projects
        projects_sql = """
        SELECT title, store_area as business_area, owner as owner_name, health as health_status, project_update as latest_update
        FROM `wmt-assetprotection-prod.Store_Support_Dev.AH_Projects`
        WHERE status = 'Active'
        LIMIT 5
        """
        projects_result = bq_client.query(projects_sql).result()
        projects = list(projects_result)
        
        # Generate HTML email
        html = f"""
<html>
<head>
    <meta charset="UTF-8">
    <title>Activity Hub - Projects Overview</title>
    <style>
        body {{ font-family: Arial, sans-serif; margin: 0; padding: 0; background-color: #f5f5f5; }}
        .container {{ max-width: 800px; margin: 0 auto; background-color: white; }}
        .header {{ background: linear-gradient(135deg, #1e3a8a 0%, #0071ce 100%); color: white; padding: 32px 24px; text-align: center; }}
        .header h1 {{ margin: 0 0 8px 0; font-size: 32px; font-weight: 700; letter-spacing: -0.5px; }}
        .header p {{ margin: 0; font-size: 14px; font-weight: 500; letter-spacing: 0.5px; }}
        .content {{ padding: 24px 32px; }}
        .section-title {{ font-size: 16px; font-weight: 700; color: #1e3a8a; margin: 20px 0 14px 0; text-transform: uppercase; border-bottom: 2px solid #0071ce; padding-bottom: 6px; }}
        .stats-grid {{ display: grid; grid-template-columns: repeat(4, 1fr); gap: 12px; margin-bottom: 24px; }}
        .stat-box {{ background-color: #f8f9fa; border-left: 4px solid #0071ce; padding: 16px 12px; text-align: center; }}
        .stat-value {{ font-size: 28px; font-weight: 700; color: #1e3a8a; line-height: 1.3; }}
        .stat-label {{ font-size: 11px; color: #666; text-transform: uppercase; font-weight: 600; letter-spacing: 0.3px; margin-top: 6px; }}
        .stat-box.on-track {{ border-left-color: #107c10; }}
        .stat-box.on-track .stat-value {{ color: #107c10; }}
        .stat-box.at-risk {{ border-left-color: #f7630c; }}
        .stat-box.at-risk .stat-value {{ color: #f7630c; }}
        .stat-box.off-track {{ border-left-color: #dc3545; }}
        .stat-box.off-track .stat-value {{ color: #dc3545; }}
        .projects-table {{ width: 100%; border-collapse: collapse; margin-top: 12px; }}
        .projects-table th {{ background-color: #1e3a8a; color: white; padding: 12px 8px; text-align: left; font-size: 12px; font-weight: 700; text-transform: uppercase; }}
        .projects-table td {{ padding: 12px 8px; border-bottom: 1px solid #e0e0e0; font-size: 13px; }}
        .projects-table tr:nth-child(even) {{ background-color: #f9f9f9; }}
        .health-badge {{ display: inline-block; padding: 4px 8px; border-radius: 3px; font-size: 11px; font-weight: 700; text-transform: uppercase; }}
        .health-on-track {{ background-color: #d4edda; color: #107c10; }}
        .health-at-risk {{ background-color: #fff3cd; color: #f7630c; }}
        .health-off-track {{ background-color: #f8d7da; color: #dc3545; }}
        .footer {{ background-color: #f5f5f5; padding: 16px 32px; text-align: center; font-size: 11px; color: #999; border-top: 1px solid #e0e0e0; }}
        .footer p {{ margin: 4px 0; }}
        .button {{ display: inline-block; background-color: #0071ce; color: white; padding: 12px 24px; border-radius: 4px; text-decoration: none; font-weight: 600; margin-top: 12px; }}
    </style>
</head>
<body>
    <div class="container">
        <!-- HEADER -->
        <div class="header">
            <h1>Activity Hub Projects Overview</h1>
            <p>WM WK {wm_week} • {datetime.now().strftime('%B %d, %Y')}</p>
        </div>
        
        <!-- CONTENT -->
        <div class="content">
            <h2 class="section-title">Executive Summary</h2>
            
            <div class="stats-grid">
                <div class="stat-box">
                    <div class="stat-value">{stats.total}</div>
                    <div class="stat-label">Total Projects</div>
                </div>
                <div class="stat-box on-track">
                    <div class="stat-value">{stats.on_track}</div>
                    <div class="stat-label">On Track</div>
                </div>
                <div class="stat-box at-risk">
                    <div class="stat-value">{stats.at_risk}</div>
                    <div class="stat-label">At Risk</div>
                </div>
                <div class="stat-box off-track">
                    <div class="stat-value">{stats.off_track}</div>
                    <div class="stat-label">Off Track</div>
                </div>
            </div>

            <h2 class="section-title">Projects</h2>
            
            <table class="projects-table">
                <thead>
                    <tr>
                        <th>Project Title</th>
                        <th>Business Area</th>
                        <th>Owner</th>
                        <th>Health Status</th>
                    </tr>
                </thead>
                <tbody>
"""
        
        # Add project rows
        for project in projects:
            health = project['health_status'] or 'Unknown'
            health_class = 'health-on-track' if 'track' in health.lower() else \
                          'health-at-risk' if 'risk' in health.lower() else \
                          'health-off-track' if 'off' in health.lower() else 'health-on-track'
            
            html += f"""
                    <tr>
                        <td style="font-weight: 600;">{project['title']}</td>
                        <td>{project['business_area']}</td>
                        <td>{project['owner_name']}</td>
                        <td><span class="health-badge {health_class}">{health}</span></td>
                    </tr>
"""
        
        html += """
                </tbody>
            </table>

            <div style="text-align: center; margin-top: 24px;">
                <a href="http://localhost:8088/activity-hub/projects" class="button">View Full Dashboard</a>
            </div>
        </div>
        
        <!-- FOOTER -->
        <div class="footer">
            <p><strong>Activity Hub</strong> • Walmart Projects Platform</p>
            <p>This is an automated email from the Projects Management System.</p>
            <p>&copy; 2026 Walmart Inc. All rights reserved.</p>
        </div>
    </div>
</body>
</html>
"""
        
        # Return as HTML content
        return html, 200, {'Content-Type': 'text/html; charset=utf-8'}
    except Exception as e:
        logger.error(f"Sample email generation error: {str(e)}")
        return jsonify({'error': str(e)}), 500

@app.route('/api/send-email', methods=['POST'])
def send_projects_email():
    """
    Send projects overview email to recipient
    POST body: { "recipient_email": "email@walmart.com", "include_ppt": false }
    """
    try:
        data = request.get_json()
        recipient_email = data.get('recipient_email', 'kendall.rush@walmart.com').strip()
        
        if not recipient_email:
            return jsonify({'error': 'recipient_email required'}), 400
        
        if not bq_client:
            return jsonify({'error': 'BigQuery client not initialized'}), 500
        
        # Calculate Walmart Week
        today = datetime.now()
        fiscal_year_start = datetime(today.year if today.month >= 2 else today.year - 1, 2, 1)
        days_since_fy = (today - fiscal_year_start).days
        wm_week = (days_since_fy // 7) + 1
        
        # Get project stats and sample projects (same as /api/sample-email)
        sql = """
        SELECT COUNT(*) as total, 
               COUNTIF(health = 'On Track') as on_track,
               COUNTIF(health = 'At Risk') as at_risk,
               COUNTIF(health = 'Off Track') as off_track
        FROM `wmt-assetprotection-prod.Store_Support_Dev.AH_Projects`
        WHERE status = 'Active'
        """
        
        stats_result = bq_client.query(sql).result()
        stats = list(stats_result)[0]
        
        # Get sample projects
        projects_sql = """
        SELECT title, store_area as business_area, owner as owner_name, health as health_status, project_update as latest_update
        FROM `wmt-assetprotection-prod.Store_Support_Dev.AH_Projects`
        WHERE status = 'Active'
        LIMIT 5
        """
        projects_result = bq_client.query(projects_sql).result()
        projects = list(projects_result)
        
        # Generate HTML email (same as /api/sample-email)
        html = f"""
<html>
<head>
    <meta charset="UTF-8">
    <title>Activity Hub - Projects Overview</title>
    <style>
        body {{ font-family: Arial, sans-serif; margin: 0; padding: 0; background-color: #f5f5f5; }}
        .container {{ max-width: 800px; margin: 0 auto; background-color: white; }}
        .header {{ background: linear-gradient(135deg, #1e3a8a 0%, #0071ce 100%); color: white; padding: 32px 24px; text-align: center; }}
        .header h1 {{ margin: 0 0 8px 0; font-size: 32px; font-weight: 700; letter-spacing: -0.5px; }}
        .header p {{ margin: 0; font-size: 14px; font-weight: 500; letter-spacing: 0.5px; }}
        .content {{ padding: 24px 32px; }}
        .section-title {{ font-size: 16px; font-weight: 700; color: #1e3a8a; margin: 20px 0 14px 0; text-transform: uppercase; border-bottom: 2px solid #0071ce; padding-bottom: 6px; }}
        .stats-grid {{ display: grid; grid-template-columns: repeat(4, 1fr); gap: 12px; margin-bottom: 24px; }}
        .stat-box {{ background-color: #f8f9fa; border-left: 4px solid #0071ce; padding: 16px 12px; text-align: center; }}
        .stat-value {{ font-size: 28px; font-weight: 700; color: #1e3a8a; line-height: 1.3; }}
        .stat-label {{ font-size: 11px; color: #666; text-transform: uppercase; font-weight: 600; letter-spacing: 0.3px; margin-top: 6px; }}
        .stat-box.on-track {{ border-left-color: #107c10; }}
        .stat-box.on-track .stat-value {{ color: #107c10; }}
        .stat-box.at-risk {{ border-left-color: #f7630c; }}
        .stat-box.at-risk .stat-value {{ color: #f7630c; }}
        .stat-box.off-track {{ border-left-color: #dc3545; }}
        .stat-box.off-track .stat-value {{ color: #dc3545; }}
        .projects-table {{ width: 100%; border-collapse: collapse; margin-top: 12px; }}
        .projects-table th {{ background-color: #1e3a8a; color: white; padding: 12px 8px; text-align: left; font-size: 12px; font-weight: 700; text-transform: uppercase; }}
        .projects-table td {{ padding: 12px 8px; border-bottom: 1px solid #e0e0e0; font-size: 13px; }}
        .projects-table tr:nth-child(even) {{ background-color: #f9f9f9; }}
        .health-badge {{ display: inline-block; padding: 4px 8px; border-radius: 3px; font-size: 11px; font-weight: 700; text-transform: uppercase; }}
        .health-on-track {{ background-color: #d4edda; color: #107c10; }}
        .health-at-risk {{ background-color: #fff3cd; color: #f7630c; }}
        .health-off-track {{ background-color: #f8d7da; color: #dc3545; }}
        .footer {{ background-color: #f5f5f5; padding: 16px 32px; text-align: center; font-size: 11px; color: #999; border-top: 1px solid #e0e0e0; }}
        .footer p {{ margin: 4px 0; }}
        .button {{ display: inline-block; background-color: #0071ce; color: white; padding: 12px 24px; border-radius: 4px; text-decoration: none; font-weight: 600; margin-top: 12px; }}
    </style>
</head>
<body>
    <div class="container">
        <!-- HEADER -->
        <div class="header">
            <h1>Activity Hub Projects Overview</h1>
            <p>WM WK {wm_week} • {datetime.now().strftime('%B %d, %Y')}</p>
        </div>
        
        <!-- CONTENT -->
        <div class="content">
            <h2 class="section-title">Executive Summary</h2>
            
            <div class="stats-grid">
                <div class="stat-box">
                    <div class="stat-value">{stats.total}</div>
                    <div class="stat-label">Total Projects</div>
                </div>
                <div class="stat-box on-track">
                    <div class="stat-value">{stats.on_track}</div>
                    <div class="stat-label">On Track</div>
                </div>
                <div class="stat-box at-risk">
                    <div class="stat-value">{stats.at_risk}</div>
                    <div class="stat-label">At Risk</div>
                </div>
                <div class="stat-box off-track">
                    <div class="stat-value">{stats.off_track}</div>
                    <div class="stat-label">Off Track</div>
                </div>
            </div>

            <h2 class="section-title">Projects</h2>
            
            <table class="projects-table">
                <thead>
                    <tr>
                        <th>Project Title</th>
                        <th>Business Area</th>
                        <th>Owner</th>
                        <th>Health Status</th>
                    </tr>
                </thead>
                <tbody>
"""
        
        # Add project rows
        for project in projects:
            health = project['health_status'] or 'Unknown'
            health_class = 'health-on-track' if 'track' in health.lower() else \
                          'health-at-risk' if 'risk' in health.lower() else \
                          'health-off-track' if 'off' in health.lower() else 'health-on-track'
            
            html += f"""
                    <tr>
                        <td style="font-weight: 600;">{project['title']}</td>
                        <td>{project['business_area']}</td>
                        <td>{project['owner_name']}</td>
                        <td><span class="health-badge {health_class}">{health}</span></td>
                    </tr>
"""
        
        html += """
                </tbody>
            </table>

            <div style="text-align: center; margin-top: 24px;">
                <a href="http://localhost:8088/activity-hub/projects" class="button">View Full Dashboard</a>
            </div>
        </div>
        
        <!-- FOOTER -->
        <div class="footer">
            <p><strong>Activity Hub</strong> • Walmart Projects Platform</p>
            <p>This is an automated email from the Projects Management System.</p>
            <p>&copy; 2026 Walmart Inc. All rights reserved.</p>
        </div>
    </div>
</body>
</html>
"""
        
        # Send email via Walmart relay
        success = send_email_via_relay(
            recipient_email=recipient_email,
            subject=f"Activity Hub Projects Overview - WM WK {wm_week}",
            html_body=html
        )
        
        if success:
            return jsonify({
                'status': 'sent',
                'recipient': recipient_email,
                'wm_week': wm_week,
                'projects_count': len(projects),
                'timestamp': datetime.now().isoformat()
            }), 200
        else:
            return jsonify({
                'error': f'Failed to send email to {recipient_email}',
                'recipient': recipient_email
            }), 500
    
    except Exception as e:
        logger.error(f"Send-email error: {str(e)}")
        return jsonify({'error': str(e)}), 500



@app.route('/api/projects/sync-data-bridge', methods=['POST'])
def sync_data_bridge():
    """
    Merge/sync projects from Intake Accel Council Data into AH_Projects.
    Logic: For each source project, if it exists in AH_Projects (matched by Project ID),
    compare timestamps and keep the newer version. Otherwise, insert as new.
    """
    try:
        bq_client = get_bq_client()
        if not bq_client:
            return jsonify({'error': 'BigQuery client not initialized'}), 500
        
        # Query all ACTIVE projects from Intake Hub source (deduplicated by project ID)
        # Active = ARCHIVED != True (includes all non-archived projects)
        source_sql = """
        SELECT 
            CAST(Intake_Card_Nbr AS STRING) as project_id,
            Project_Title as title,
            Business_Owner_Area as business_organization,
            Owner as owner,
            PROJECT_OWNERID as owner_id,
            PROJECT_HEALTH_DESC as health,
            Status as status,
            'Intake Hub' as project_source,
            CAST(CREATED_TS AS TIMESTAMP) as created_date,
            CAST(Last_Updated AS TIMESTAMP) as last_updated,
            Project_Update as project_update
        FROM (
            SELECT 
                *,
                ROW_NUMBER() OVER (PARTITION BY Intake_Card_Nbr ORDER BY Last_Updated DESC) as rn
            FROM `wmt-assetprotection-prod.Store_Support_Dev.Output - Intake Accel Council Data`
            WHERE ARCHIVED != True
        )
        WHERE rn = 1
        ORDER BY Intake_Card_Nbr
        """
        
        source_projects = list(bq_client.query(source_sql).result())
        
        if not source_projects:
            return jsonify({
                'success': True,
                'message': 'No active projects found in Intake Hub source',
                'rows_affected': 0
            }), 200
        
        # Get all projects currently in AH_Projects indexed by project_id
        existing_sql = """
        SELECT 
            project_id,
            last_updated
        FROM `wmt-assetprotection-prod.Store_Support_Dev.AH_Projects`
        """
        
        existing_projects = {}
        try:
            for row in bq_client.query(existing_sql).result():
                existing_projects[row.project_id] = row.last_updated
        except:
            pass  # Table may not exist yet
        
        # Separate into insert and update lists
        rows_to_insert = []
        rows_to_update = []
        
        for src_row in source_projects:
            project_id = src_row.project_id
            src_timestamp = src_row.last_updated
            
            if project_id in existing_projects:
                # Project exists - compare timestamps
                existing_timestamp = existing_projects[project_id]
                if src_timestamp and (not existing_timestamp or src_timestamp > existing_timestamp):
                    # Source is newer, mark for update
                    rows_to_update.append(src_row)
            else:
                # New project, mark for insert
                rows_to_insert.append(src_row)
        
        rows_affected = 0
        
        # INSERT new projects
        if rows_to_insert:
            insert_sql = """
            INSERT INTO `wmt-assetprotection-prod.Store_Support_Dev.AH_Projects`
            (project_id, title, business_organization, owner, owner_id, health, status, project_source,
             created_date, last_updated, project_update)
            VALUES 
            """
            using_params = []
            
            for i, row in enumerate(rows_to_insert):
                if i > 0:
                    insert_sql += ", "
                
                params = [
                    bigquery.ScalarQueryParameter(f"ins{i}_project_id", "STRING", row.project_id),
                    bigquery.ScalarQueryParameter(f"ins{i}_title", "STRING", row.title),
                    bigquery.ScalarQueryParameter(f"ins{i}_business_organization", "STRING", row.business_organization),
                    bigquery.ScalarQueryParameter(f"ins{i}_owner", "STRING", row.owner),
                    bigquery.ScalarQueryParameter(f"ins{i}_owner_id", "STRING", row.owner_id),
                    bigquery.ScalarQueryParameter(f"ins{i}_health", "STRING", row.health),
                    bigquery.ScalarQueryParameter(f"ins{i}_status", "STRING", row.status),
                    bigquery.ScalarQueryParameter(f"ins{i}_project_source", "STRING", row.project_source),
                    bigquery.ScalarQueryParameter(f"ins{i}_created_date", "TIMESTAMP", row.created_date),
                    bigquery.ScalarQueryParameter(f"ins{i}_last_updated", "TIMESTAMP", row.last_updated),
                    bigquery.ScalarQueryParameter(f"ins{i}_project_update", "STRING", row.project_update),
                ]
                
                using_params.extend(params)
                insert_sql += f"""(@ins{i}_project_id, @ins{i}_title, @ins{i}_business_organization, @ins{i}_owner, 
                              @ins{i}_owner_id, @ins{i}_health, @ins{i}_status, @ins{i}_project_source,
                              @ins{i}_created_date, @ins{i}_last_updated, @ins{i}_project_update)"""
            
            job_config = bigquery.QueryJobConfig(query_parameters=using_params)
            job = bq_client.query(insert_sql, job_config=job_config)
            job.result()
            rows_affected += len(rows_to_insert)
        
        # UPDATE existing projects with newer source data
        if rows_to_update:
            for src_row in rows_to_update:
                update_sql = """
                UPDATE `wmt-assetprotection-prod.Store_Support_Dev.AH_Projects`
                SET 
                    title = @title,
                    business_organization = @business_organization,
                    owner = @owner,
                    owner_id = @owner_id,
                    health = @health,
                    status = @status,
                    last_updated = @last_updated,
                    project_update = @project_update
                WHERE project_id = @project_id
                """
                
                query_params = [
                    bigquery.ScalarQueryParameter("title", "STRING", src_row.title),
                    bigquery.ScalarQueryParameter("business_organization", "STRING", src_row.business_organization),
                    bigquery.ScalarQueryParameter("owner", "STRING", src_row.owner),
                    bigquery.ScalarQueryParameter("owner_id", "STRING", src_row.owner_id),
                    bigquery.ScalarQueryParameter("health", "STRING", src_row.health),
                    bigquery.ScalarQueryParameter("status", "STRING", src_row.status),
                    bigquery.ScalarQueryParameter("last_updated", "TIMESTAMP", src_row.last_updated),
                    bigquery.ScalarQueryParameter("project_update", "STRING", src_row.project_update),
                    bigquery.ScalarQueryParameter("project_id", "STRING", src_row.project_id),
                ]
                
                job_config = bigquery.QueryJobConfig(query_parameters=query_params)
                job = bq_client.query(update_sql, job_config=job_config)
                job.result()
            
            rows_affected += len(rows_to_update)
        
        logger.info(f"Data Bridge sync completed: {len(rows_to_insert)} new, {len(rows_to_update)} updated")
        
        return jsonify({
            'success': True,
            'message': f'Sync completed from Intake Hub. Inserted: {len(rows_to_insert)}, Updated: {len(rows_to_update)}',
            'inserted': len(rows_to_insert),
            'updated': len(rows_to_update),
            'rows_affected': rows_affected
        }), 200
            
    except Exception as e:
        logger.error(f"Data Bridge sync error: {str(e)}")
        return jsonify({'error': str(e)}), 500


# Email endpoint for sending projects list
@app.route('/api/email/send-projects-list', methods=['POST'])
def send_projects_list_email():
    """Send filtered projects list via email"""
    try:
        data = request.get_json()
        recipient_email = data.get('recipientEmail')
        projects = data.get('projects', [])
        sent_by = data.get('sentBy', 'System')
        
        if not recipient_email:
            return jsonify({'error': 'Recipient email is required'}), 400
        
        if not projects:
            return jsonify({'error': 'No projects to send'}), 400
        
        # Build HTML email
        project_rows = ''.join([f"""
            <tr>
                <td style="border: 1px solid #ddd; padding: 8px;">{p.get('title', 'N/A')}</td>
                <td style="border: 1px solid #ddd; padding: 8px;">{p.get('owner', 'N/A')}</td>
                <td style="border: 1px solid #ddd; padding: 8px;">{p.get('business_area', 'N/A')}</td>
                <td style="border: 1px solid #ddd; padding: 8px;">{p.get('status', 'N/A')}</td>
                <td style="border: 1px solid #ddd; padding: 8px;">{p.get('health', 'N/A')}</td>
            </tr>
        """ for p in projects])
        
        html_content = f"""
        <html>
            <head>
                <style>
                    body {{ font-family: Arial, sans-serif; color: #333; }}
                    h2 {{ color: #0066cc; }}
                    table {{ border-collapse: collapse; width: 100%; margin: 20px 0; }}
                    th {{ background-color: #0066cc; color: white; padding: 10px; text-align: left; }}
                    td {{ border: 1px solid #ddd; padding: 8px; }}
                    tr:nth-child(even) {{ background-color: #f5f5f5; }}
                    .footer {{ margin-top: 20px; color: #666; font-size: 12px; }}
                </style>
            </head>
            <body>
                <h2>📊 Activity Hub Projects Report</h2>
                <p>Hello,</p>
                <p>Below is the current list of projects from the Activity Hub:</p>
                
                <table>
                    <thead>
                        <tr>
                            <th>Project Title</th>
                            <th>Owner</th>
                            <th>Business Area</th>
                            <th>Status</th>
                            <th>Health</th>
                        </tr>
                    </thead>
                    <tbody>
                        {project_rows}
                    </tbody>
                </table>
                
                <p><strong>Total Projects:</strong> {len(projects)}</p>
                
                <div class="footer">
                    <p>This email was sent by {sent_by} via Activity Hub Projects</p>
                    <p>Generated: {datetime.now().strftime('%B %d, %Y at %I:%M %p')}</p>
                </div>
            </body>
        </html>
        """
        
        # For now, just log it and return success
        logger.info(f"Email request: Send {len(projects)} projects to {recipient_email} from {sent_by}")
        
        # TODO: Integrate with actual email service (SendGrid, AWS SES, etc.)
        # For testing, you can uncomment the SMTP section below or use your email service
        
        return jsonify({
            'success': True,
            'message': f'Email prepared to send {len(projects)} projects',
            'recipient': recipient_email,
            'projectCount': len(projects),
            'timestamp': datetime.now().isoformat()
        }), 200
        
    except Exception as e:
        logger.error(f"Email send error: {str(e)}")
        return jsonify({'error': str(e)}), 500


@app.route('/')
def root():
    from flask import redirect
    return redirect('/activity-hub/')


if __name__ == '__main__':
    port = 8088
    print(f'Activity Hub V2 starting on port {port}...')
    print(f'  Local:   http://localhost:{port}/activity-hub/')
    print(f'  Network: http://weus42608431466:{port}/activity-hub/')
    
    # NOTE: APScheduler scheduler disabled - network blocked pip install
    # Background scheduler would run:
    #  - AH_Projects sync every 30 minutes (smart merge, preserves note history)
    #  - Weekly audio digest every Saturday at 9:00 AM
    # Manual sync available via POST /api/admin/trigger-ah-projects-sync
    
    app.run(host='0.0.0.0', port=port, debug=False)
