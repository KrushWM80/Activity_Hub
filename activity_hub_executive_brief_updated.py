"""
Activity Hub Executive Brief - Updated with Actual Platform Assessment Data
Based on comprehensive codebase analysis and platform assessment

This script generates a 5-slide PowerPoint presentation with accurate
Activity Hub specifications, costs, and ROI projections.

Requirements:
    pip install python-pptx

Usage:
    python activity_hub_executive_brief_updated.py

Output:
    Activity_Hub_Executive_Brief_Updated.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Walmart brand colors
WALMART_BLUE = RGBColor(0, 113, 206)  # #0071ce
WALMART_YELLOW = RGBColor(255, 194, 32)  # #ffc220
WHITE = RGBColor(255, 255, 255)
DARK_GRAY = RGBColor(50, 50, 50)
LIGHT_GRAY = RGBColor(200, 200, 200)
GREEN = RGBColor(0, 150, 0)
RED = RGBColor(200, 0, 0)


def add_title_slide(prs, title, subtitle, bottom_text):
    """Create title slide with Walmart branding."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WALMART_BLUE
    shapes = slide.shapes

    # Title
    title_box = shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.2))
    title_tf = title_box.text_frame
    title_tf.text = title
    title_tf.paragraphs[0].font.size = Pt(54)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = WHITE
    title_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = shapes.add_textbox(Inches(1), Inches(3.2), Inches(8), Inches(0.8))
    subtitle_tf = subtitle_box.text_frame
    subtitle_tf.text = subtitle
    subtitle_tf.paragraphs[0].font.size = Pt(32)
    subtitle_tf.paragraphs[0].font.color.rgb = WHITE
    subtitle_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Bottom text (investment/ROI)
    bottom_box = shapes.add_textbox(Inches(1), Inches(6), Inches(8), Inches(0.8))
    bottom_tf = bottom_box.text_frame
    bottom_tf.text = bottom_text
    bottom_tf.paragraphs[0].font.size = Pt(24)
    bottom_tf.paragraphs[0].font.bold = True
    bottom_tf.paragraphs[0].font.color.rgb = WALMART_YELLOW
    bottom_tf.paragraphs[0].alignment = PP_ALIGN.CENTER


def add_callout_box(slide, text, left, top, width, height, color):
    """Add a highlighted callout box."""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    tf = box.text_frame
    tf.text = text
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = 1  # Middle


def create_presentation():
    """Create the executive brief presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    print("Creating Activity Hub Executive Brief Presentation...")

    # Slide 1: Title Slide
    print("  Adding Slide 1: Title Slide")
    add_title_slide(
        prs,
        "Walmart Activity Hub",
        "Store Operations Command Center",
        "$227K Year 1 Investment | $27M Annual Value | 694% ROI"
    )

    # Slide 2: Platform Overview & The Problem
    print("  Adding Slide 2: Platform Overview & The Problem")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shapes = slide.shapes

    title_box = shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_tf = title_box.text_frame
    title_tf.text = "Enterprise-Scale Store Operations Platform"
    title_tf.paragraphs[0].font.size = Pt(44)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = WALMART_BLUE

    # Left: Platform Specs
    left_box = shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(2.5))
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = RGBColor(240, 248, 255)
    left_tf = left_box.text_frame
    left_tf.text = "Platform Scale & Reach"
    left_tf.paragraphs[0].font.size = Pt(22)
    left_tf.paragraphs[0].font.bold = True
    left_tf.paragraphs[0].font.color.rgb = WALMART_BLUE

    platform_items = [
        "🏪 4,700+ Walmart US stores",
        "👥 50,000+ target users",
        "🌍 5 regional deployments",
        "📊 Real-time data processing",
        "🔐 8 role-based access levels",
        "🤖 AI-powered insights (Sparky)",
        "",
        "Status: Production Ready v1.0"
    ]

    for item in platform_items:
        p = left_tf.add_paragraph()
        p.text = item
        p.font.size = Pt(14) if "Status" not in item else Pt(16)
        p.font.bold = "Status" in item
        p.font.color.rgb = GREEN if "Status" in item else DARK_GRAY

    # Right: Technology Stack
    right_box = shapes.add_textbox(Inches(5.5), Inches(1.5), Inches(4), Inches(2.5))
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = RGBColor(240, 255, 240)
    right_tf = right_box.text_frame
    right_tf.text = "Enterprise Architecture"
    right_tf.paragraphs[0].font.size = Pt(22)
    right_tf.paragraphs[0].font.bold = True
    right_tf.paragraphs[0].font.color.rgb = GREEN

    tech_items = [
        "⚛️ React 18 + TypeScript frontend",
        "🐍 FastAPI microservices backend",
        "🗄️ PostgreSQL + Redis caching",
        "🔌 WebSocket real-time updates",
        "🤖 OpenAI & ML integration",
        "🐳 Docker + Kubernetes ready",
        "",
        "Complexity Rating: HIGH"
    ]

    for item in tech_items:
        p = right_tf.add_paragraph()
        p.text = item
        p.font.size = Pt(14) if "Complexity" not in item else Pt(16)
        p.font.bold = "Complexity" in item
        p.font.color.rgb = WALMART_BLUE if "Complexity" in item else DARK_GRAY

    # Bottom: Current Problem
    problem_box = shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2.5))
    problem_tf = problem_box.text_frame
    problem_tf.text = "The Problem We're Solving:"
    problem_tf.paragraphs[0].font.size = Pt(20)
    problem_tf.paragraphs[0].font.bold = True
    problem_tf.paragraphs[0].font.color.rgb = RED

    problems = [
        "❌ Disconnected tools causing 4-6 hours/week per user inefficiency",
        "❌ No centralized visibility across 4,700+ store operations",
        "❌ Manual activity tracking costing $3.1M annually in lost productivity",
        "❌ 30-40% message overlap creating store communication confusion",
        "❌ 30% time wasted on administrative tasks across all roles"
    ]

    for problem in problems:
        p = problem_tf.add_paragraph()
        p.text = problem
        p.font.size = Pt(15)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(4)

    # Slide 3: Financial Investment Breakdown
    print("  Adding Slide 3: Financial Investment Breakdown")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shapes = slide.shapes

    title_box = shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_tf = title_box.text_frame
    title_tf.text = "Year 1 Investment Breakdown"
    title_tf.paragraphs[0].font.size = Pt(44)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = WALMART_BLUE

    # Development Costs
    dev_box = shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4.5), Inches(2.8))
    dev_box.fill.solid()
    dev_box.fill.fore_color.rgb = RGBColor(240, 248, 255)
    dev_tf = dev_box.text_frame
    dev_tf.text = "💻 Development Costs"
    dev_tf.paragraphs[0].font.size = Pt(20)
    dev_tf.paragraphs[0].font.bold = True
    dev_tf.paragraphs[0].font.color.rgb = WALMART_BLUE

    dev_items = [
        "",
        "React Frontend: $35,000",
        "FastAPI Backend: $50,000",
        "AI/ML Integration: $15,000",
        "Core Features (12): $36,000",
        "Advanced Features (4): $12,000",
        "Security (Advanced RBAC): $10,000",
        "Integration & Testing (15%): $29,700",
        "",
        "Development Total: $187,700"
    ]

    for item in dev_items:
        p = dev_tf.add_paragraph()
        p.text = item
        p.font.size = Pt(14) if "Total" not in item else Pt(16)
        p.font.bold = "Total" in item
        p.font.color.rgb = WALMART_BLUE if "Total" in item else DARK_GRAY

    # Infrastructure & Ongoing
    infra_box = shapes.add_textbox(Inches(5.5), Inches(1.3), Inches(4), Inches(2.8))
    infra_box.fill.solid()
    infra_box.fill.fore_color.rgb = RGBColor(240, 255, 240)
    infra_tf = infra_box.text_frame
    infra_tf.text = "☁️ Infrastructure (Year 1)"
    infra_tf.paragraphs[0].font.size = Pt(20)
    infra_tf.paragraphs[0].font.bold = True
    infra_tf.paragraphs[0].font.color.rgb = GREEN

    infra_items = [
        "",
        "Cloud Hosting: $25,000",
        "Database (PostgreSQL): $8,000",
        "Redis Caching: $3,000",
        "Monitoring & Logging: $4,000",
        "",
        "Infrastructure Total: $40,000",
        "",
        "YEAR 1 TOTAL: $227,700"
    ]

    for item in infra_items:
        p = infra_tf.add_paragraph()
        p.text = item
        p.font.size = Pt(14) if "YEAR 1" not in item else Pt(18)
        p.font.bold = "Total" in item or "YEAR 1" in item
        p.font.color.rgb = GREEN if "Total" in item or "YEAR 1" in item else DARK_GRAY

    # ROI Summary
    roi_box = shapes.add_textbox(Inches(0.5), Inches(4.3), Inches(9), Inches(2.5))
    roi_box.fill.solid()
    roi_box.fill.fore_color.rgb = RGBColor(255, 250, 240)
    roi_tf = roi_box.text_frame
    roi_tf.text = "💰 Return on Investment Analysis"
    roi_tf.paragraphs[0].font.size = Pt(22)
    roi_tf.paragraphs[0].font.bold = True
    roi_tf.paragraphs[0].font.color.rgb = RGBColor(200, 100, 0)

    roi_items = [
        "",
        "Year 1 Investment: $227,700 (documented, platform-assessed cost)",
        "Annual Benefits: $27,000,000 (time savings, productivity, better decisions)",
        "3-Year Total Benefits: $81,000,000",
        "3-Year Net Value: $80,772,300",
        "",
        "First-Year ROI: 11,755% | Break-Even: < 3 days"
    ]

    for item in roi_items:
        p = roi_tf.add_paragraph()
        p.text = item
        p.font.size = Pt(15) if "ROI" not in item else Pt(18)
        p.font.bold = "ROI" in item or "Year 1" in item
        p.font.color.rgb = RGBColor(200, 100, 0) if "ROI" in item else DARK_GRAY

    # Slide 4: Feature Set & Capabilities
    print("  Adding Slide 4: Feature Set & Capabilities")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shapes = slide.shapes

    title_box = shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_tf = title_box.text_frame
    title_tf.text = "Comprehensive Feature Set (15+ Major Features)"
    title_tf.paragraphs[0].font.size = Pt(40)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = WALMART_BLUE

    # Executive Dashboard Features
    exec_box = shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4.5), Inches(2.5))
    exec_tf = exec_box.text_frame
    exec_tf.text = "📊 Executive Dashboard"
    exec_tf.paragraphs[0].font.size = Pt(18)
    exec_tf.paragraphs[0].font.bold = True
    exec_tf.paragraphs[0].font.color.rgb = WALMART_BLUE

    exec_features = [
        "• 8 real-time KPIs (stores, projects, safety, satisfaction)",
        "• Regional performance (5 regions, 4,700+ stores)",
        "• Critical alerts dashboard (priority-based)",
        "• Trend analysis & forecasting",
        "• Board-ready reports & analytics",
        "• Strategic goal tracking"
    ]

    for feature in exec_features:
        p = exec_tf.add_paragraph()
        p.text = feature
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_GRAY

    # Manager Dashboard Features
    mgr_box = shapes.add_textbox(Inches(5.5), Inches(1.3), Inches(4), Inches(2.5))
    mgr_tf = mgr_box.text_frame
    mgr_tf.text = "🎯 Manager Dashboard"
    mgr_tf.paragraphs[0].font.size = Pt(18)
    mgr_tf.paragraphs[0].font.bold = True
    mgr_tf.paragraphs[0].font.color.rgb = GREEN

    mgr_features = [
        "• Store-level performance tracking",
        "• Activity management (CRUD operations)",
        "• Team coordination & scheduling",
        "• Overdue activity alerts",
        "• Completion tracking (daily/weekly)",
        "• Drill-down analytics by store"
    ]

    for feature in mgr_features:
        p = mgr_tf.add_paragraph()
        p.text = feature
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_GRAY

    # Core Platform Features
    core_box = shapes.add_textbox(Inches(0.5), Inches(4), Inches(4.5), Inches(2.7))
    core_tf = core_box.text_frame
    core_tf.text = "⚙️ Core Platform Capabilities"
    core_tf.paragraphs[0].font.size = Pt(18)
    core_tf.paragraphs[0].font.bold = True
    core_tf.paragraphs[0].font.color.rgb = WALMART_BLUE

    core_features = [
        "• Real-time notifications (WebSocket)",
        "• Advanced search & filtering",
        "• Custom reporting & data export",
        "• Multi-channel communication hub",
        "• User authentication (SSO + AD Groups)",
        "• Mobile-responsive design",
        "• Drag-and-drop widget customization"
    ]

    for feature in core_features:
        p = core_tf.add_paragraph()
        p.text = feature
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_GRAY

    # AI & Integration Features
    ai_box = shapes.add_textbox(Inches(5.5), Inches(4), Inches(4), Inches(2.7))
    ai_box.fill.solid()
    ai_box.fill.fore_color.rgb = RGBColor(255, 250, 240)
    ai_tf = ai_box.text_frame
    ai_tf.text = "🤖 AI & Integration Features"
    ai_tf.paragraphs[0].font.size = Pt(18)
    ai_tf.paragraphs[0].font.bold = True
    ai_tf.paragraphs[0].font.color.rgb = RGBColor(200, 100, 0)

    ai_features = [
        "• Sparky AI assistant (context-aware)",
        "• Sentiment analysis for communications",
        "• Predictive analytics & recommendations",
        "• Automated workflow optimization",
        "• Intake Hub API integration",
        "• WalmartOne integration",
        "• Store Operations API connector"
    ]

    for feature in ai_features:
        p = ai_tf.add_paragraph()
        p.text = feature
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_GRAY

    # Slide 5: Timeline, Benefits & Call to Action
    print("  Adding Slide 5: Timeline, Benefits & Call to Action")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shapes = slide.shapes

    title_box = shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_tf = title_box.text_frame
    title_tf.text = "Implementation & Value Realization"
    title_tf.paragraphs[0].font.size = Pt(44)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = WALMART_BLUE

    # Timeline
    timeline_box = shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4.5), Inches(2.2))
    timeline_tf = timeline_box.text_frame
    timeline_tf.text = "⏱️ Deployment Timeline"
    timeline_tf.paragraphs[0].font.size = Pt(20)
    timeline_tf.paragraphs[0].font.bold = True
    timeline_tf.paragraphs[0].font.color.rgb = WALMART_BLUE
    
    timeline_items = [
        "",
        "✅ Phase 1-3: COMPLETED",
        "  Foundation, features, AI integration",
        "",
        "✅ Phase 4: COMPLETED",
        "  Integration, testing, optimization",
        "",
        "🎯 Current Status:",
        "  Production Ready v1.0.0",
        "",
        "📅 Rollout Timeline:",
        "  Months 1-2: Pilot (1,000 users)",
        "  Months 3-6: Regional (10,000 users)",
        "  Months 7-12: Full scale (50,000)"
    ]

    for item in timeline_items:
        p = timeline_tf.add_paragraph()
        p.text = item
        p.font.size = Pt(12) if item.startswith("  ") else Pt(14)
        p.font.bold = not item.startswith("  ") and item != ""
        p.font.color.rgb = GREEN if "✅" in item or "🎯" in item else DARK_GRAY

    # Annual Benefits
    benefits_box = shapes.add_textbox(Inches(5.5), Inches(1.3), Inches(4), Inches(2.2))
    benefits_box.fill.solid()
    benefits_box.fill.fore_color.rgb = RGBColor(240, 255, 240)
    benefits_tf = benefits_box.text_frame
    benefits_tf.text = "💰 Annual Benefits ($27M)"
    benefits_tf.paragraphs[0].font.size = Pt(20)
    benefits_tf.paragraphs[0].font.bold = True
    benefits_tf.paragraphs[0].font.color.rgb = GREEN

    benefits_items = [
        "",
        "Time Savings:",
        "  $15.6M (4-6 hrs/week × 50K users)",
        "",
        "Productivity Improvement:",
        "  $6.2M (15% project delivery gain)",
        "",
        "Administrative Reduction:",
        "  $3.1M (30% task time savings)",
        "",
        "Better Collaboration:",
        "  $2.1M (40% cross-functional efficiency)",
        "",
        "Delay Cost: $2.25M per month"
    ]

    for item in benefits_items:
        p = benefits_tf.add_paragraph()
        p.text = item
        p.font.size = Pt(12) if item.startswith("  ") else Pt(14)
        p.font.bold = ":" in item and not item.startswith("  ")
        p.font.color.rgb = RED if "Delay" in item else (GREEN if ":" in item else DARK_GRAY)

    # Call to Action Box
    cta_box = shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(9), Inches(3))
    cta_box.fill.solid()
    cta_box.fill.fore_color.rgb = WALMART_BLUE
    cta_tf = cta_box.text_frame
    cta_tf.text = "Ready to Deploy?"
    cta_tf.paragraphs[0].font.size = Pt(36)
    cta_tf.paragraphs[0].font.bold = True
    cta_tf.paragraphs[0].font.color.rgb = WALMART_YELLOW
    cta_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    cta_items = [
        "",
        "Platform Status: Production Ready ✅",
        "Year 1 Investment: $227,700",
        "Annual Value: $27,000,000",
        "ROI: 11,755% | Payback: 3 days",
        "",
        "Approval Needed: Deploy to 50,000+ Walmart users",
        "",
        "Target Launch: Q1 2026 | Pilot Ready Now"
    ]

    for item in cta_items:
        p = cta_tf.add_paragraph()
        p.text = item
        p.font.size = Pt(18) if "Approval" in item or "Target" in item else Pt(20)
        p.font.bold = "Year 1" in item or "ROI" in item or "Approval" in item
        p.font.color.rgb = WALMART_YELLOW if "Year 1" in item or "ROI" in item else WHITE
        p.alignment = PP_ALIGN.CENTER

    return prs


def main():
    """Main execution function."""
    print("\n" + "="*70)
    print("Activity Hub Executive Brief Generator (Updated with Assessment Data)")
    print("="*70 + "\n")

    prs = create_presentation()

    output_file = "Activity_Hub_Executive_Brief_Updated.pptx"
    prs.save(output_file)

    print("\n" + "="*70)
    print(f"✅ Executive Brief created successfully!")
    print(f"✅ File saved as: {output_file}")
    print(f"✅ Total slides: {len(prs.slides)}")
    print("="*70 + "\n")
    print("Updated brief includes:")
    print("  • Slide 1: Title with accurate investment ($227K)")
    print("  • Slide 2: Platform scale (4,700 stores, 50K users)")
    print("  • Slide 3: Detailed financial breakdown")
    print("  • Slide 4: Complete feature inventory (15+ features)")
    print("  • Slide 5: Timeline & $27M annual benefits")
    print("\n")
    print("Key Updates from Platform Assessment:")
    print("  ✓ Accurate Year 1 cost: $227,700 (vs $55K placeholder)")
    print("  ✓ Comprehensive tech stack details")
    print("  ✓ All 15+ implemented features documented")
    print("  ✓ Production-ready status highlighted")
    print("  ✓ Real user scale: 50,000 target users")
    print("  ✓ Actual complexity rating: HIGH")
    print("\n")


if __name__ == "__main__":
    main()
