"""Compare PPT sizing across reference, email, and dashboard-generated files."""
from pptx import Presentation
import glob, os

email_dir = r'Store Support\Projects\VET_Dashboard\reports'
files_to_check = {}

# Dashboard-generated
wk09 = os.path.join(email_dir, 'VET_Executive_Report_WK09.pptx')
if os.path.exists(wk09):
    files_to_check['WK09 (dashboard)'] = wk09

# Reference
ref = os.path.join(email_dir, 'VET_Executive_Report (3).pptx')
if os.path.exists(ref):
    files_to_check['Reference (3)'] = ref

# Find email-generated from today (April 2) or latest
all_tda = sorted(glob.glob(os.path.join(email_dir, 'TDA_Report_*.pptx')), reverse=True)
for f in all_tda:
    if '20260402' in os.path.basename(f):
        files_to_check['Email 6AM (Apr 2)'] = f
        break
if all_tda and 'Email 6AM (Apr 2)' not in files_to_check:
    files_to_check['Latest TDA Report'] = all_tda[0]

for label, path in files_to_check.items():
    prs = Presentation(path)
    sw = prs.slide_width / 914400
    sh = prs.slide_height / 914400
    print(f'\n=== {label} ({os.path.basename(path)}) ===')
    print(f'Slide size: {sw:.2f}" x {sh:.2f}"')
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.shape_type == 13:  # Picture
                l = shape.left / 914400
                t = shape.top / 914400
                w = shape.width / 914400
                h = shape.height / 914400
                print(f'  Slide {i+1}: pos=({l:.2f}",{t:.2f}") size=({w:.2f}"x{h:.2f}")')
