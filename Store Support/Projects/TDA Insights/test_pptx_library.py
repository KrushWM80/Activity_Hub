"""
TEST SCRIPT: python-pptx library for mobile-compatible PPTX generation.

This is a SEPARATE test — it does NOT affect the production email.
It reads the existing screenshots from the last PPT run (or captures new ones)
and generates a proper PPTX using the python-pptx library with slide masters,
layouts, and themes — which mobile/tablet PowerPoint apps require.

Usage:
    python test_pptx_library.py
"""
import sys, os, io, tempfile
from pathlib import Path
from PIL import Image

SCRIPT_DIR = Path(__file__).parent
sys.path.insert(0, str(SCRIPT_DIR))

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("ERROR: python-pptx is not installed.")
    print("Install it with:  pip install python-pptx")
    sys.exit(1)

# Import the screenshot capture logic from the production script
from send_weekly_report import (
    fetch_data, PHASE_ORDER, _measure_row_heights, _paginate_by_height,
    _build_phase_html, _capture_html_screenshot, WEEK_NUM
)

OUTPUT_PATH = SCRIPT_DIR / f"TDA_WK{WEEK_NUM}_Report_PPTX_TEST.pptx"

# Slide dimensions (standard 10x7.5 inches)
SLIDE_WIDTH = Inches(10)
SLIDE_HEIGHT = Inches(7.5)

# Colors
WALMART_BLUE = RGBColor(0x3B, 0x82, 0xF6)
WALMART_BLUE_DARK = RGBColor(0x1E, 0x3A, 0x5F)
WALMART_YELLOW = RGBColor(0xFF, 0xC2, 0x20)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x66, 0x66, 0x66)


def build_sections(data):
    """Build ownership → phase sections (same logic as production)."""
    UNASSIGNED = 'No Selection Provided'
    sections = []

    pending_unassigned = [r for r in data if r['Phase'] == 'Pending' and r.get('TDA Ownership', UNASSIGNED) == UNASSIGNED]
    if pending_unassigned:
        sections.append(('TDA Ownership - Currently No TDA Ownership', 'Pending', pending_unassigned))

    OWNERSHIP_ORDER = ['Dallas VET', 'Intake & Test', 'Deployment', UNASSIGNED]
    all_ownerships = set(r.get('TDA Ownership', UNASSIGNED) or UNASSIGNED for r in data)
    known_set = set(OWNERSHIP_ORDER)
    unknown_ownerships = sorted(o for o in all_ownerships if o not in known_set)
    ownership_list = [o for o in OWNERSHIP_ORDER if o in all_ownerships] + unknown_ownerships
    for ownership in ownership_list:
        for phase in PHASE_ORDER:
            if phase == 'Pending' and ownership == UNASSIGNED:
                continue
            rows = [r for r in data if (r.get('TDA Ownership', UNASSIGNED) or UNASSIGNED) == ownership and r['Phase'] == phase]
            if not rows:
                continue
            label = 'TDA Ownership - Currently No TDA Ownership' if ownership == UNASSIGNED else ownership
            sections.append((label, phase, rows))

    return sections


def add_title_slide(prs):
    """Add a branded title slide."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Blue header bar
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Emu(0), Emu(0), SLIDE_WIDTH, Inches(1.2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = WALMART_BLUE
    shape.line.fill.background()

    # Header text
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = "TDA Initiatives Insights"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = "Segoe UI"
    tf.margin_left = Inches(0.5)

    # Main title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(9), Inches(1.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Initiative Status Insights"
    run.font.size = Pt(40)
    run.font.bold = True
    run.font.color.rgb = WALMART_BLUE_DARK
    run.font.name = "Segoe UI"

    # Subtitle
    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(3.9), Inches(9), Inches(0.8))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = "Store Support  |  Asset Protection"
    run2.font.size = Pt(20)
    run2.font.color.rgb = GRAY
    run2.font.name = "Segoe UI"

    # Yellow accent bar at bottom
    accent = slide.shapes.add_shape(
        1,  # RECTANGLE
        Emu(0), Inches(7.1), SLIDE_WIDTH, Inches(0.4)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = WALMART_YELLOW
    accent.line.fill.background()


def add_screenshot_slide(prs, label, png_bytes):
    """Add a slide with a screenshot image, scaled to fill width."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Get image dimensions
    img = Image.open(io.BytesIO(png_bytes))
    img_w, img_h = img.size

    # Scale to slide width, maintaining aspect ratio
    aspect = img_h / img_w
    final_w = SLIDE_WIDTH
    final_h = int(SLIDE_WIDTH * aspect)
    if final_h > SLIDE_HEIGHT:
        final_h = SLIDE_HEIGHT

    # Add the image
    img_stream = io.BytesIO(png_bytes)
    slide.shapes.add_picture(img_stream, Emu(0), Emu(0), final_w, Emu(final_h))


def main():
    print("=" * 60)
    print("TEST: python-pptx Library — Mobile-Compatible PPTX")
    print("=" * 60)

    os.environ.setdefault(
        'GOOGLE_APPLICATION_CREDENTIALS',
        r'C:\Users\krush\AppData\Roaming\gcloud\application_default_credentials.json'
    )

    # 1. Fetch data
    print("\n[1/3] Fetching data from BigQuery...")
    data = fetch_data()
    print(f"  Loaded {len(data)} projects")

    # 2. Capture screenshots (same as production)
    print("[2/3] Capturing screenshots...")
    sections = build_sections(data)
    screenshots = []

    with tempfile.TemporaryDirectory() as tmp_dir:
        for ownership_label, phase, rows in sections:
            if not rows:
                continue
            heights = _measure_row_heights(phase, rows)
            pages = _paginate_by_height(rows, heights)
            total_pages = len(pages)
            for page, page_rows in enumerate(pages):
                label = f'{ownership_label} — {phase}'
                if total_pages > 1:
                    label += f' ({page + 1}/{total_pages})'

                html = _build_phase_html(phase, page_rows, ownership=ownership_label)
                png_path = os.path.join(tmp_dir, f'slide_{len(screenshots) + 1}.png')
                _capture_html_screenshot(html, png_path)
                png_bytes = Path(png_path).read_bytes()
                screenshots.append((label, png_bytes))
                print(f"    {label} ({len(png_bytes):,} bytes)")

    # 3. Build PPTX with python-pptx
    print("[3/3] Building PPTX with python-pptx...")
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    add_title_slide(prs)

    for label, png_bytes in screenshots:
        add_screenshot_slide(prs, label, png_bytes)

    prs.save(str(OUTPUT_PATH))
    file_size = OUTPUT_PATH.stat().st_size
    print(f"\n  Saved: {OUTPUT_PATH}")
    print(f"  Size: {file_size:,} bytes ({file_size / 1024 / 1024:.1f} MB)")
    print(f"  Slides: {len(prs.slides)}")

    print("\n" + "=" * 60)
    print("TEST COMPLETE — Try opening this file on mobile/tablet!")
    print("=" * 60)


if __name__ == '__main__':
    main()
