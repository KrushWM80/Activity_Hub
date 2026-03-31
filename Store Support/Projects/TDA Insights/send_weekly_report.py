#!/usr/bin/env python3
"""
TDA Insights Weekly Email Report
Generates a PPT of all projects and sends via Outlook.
Scheduled for Thursdays at 11:00 AM Central Time.

Usage:
  python send_weekly_report.py              # Send email
  python send_weekly_report.py --preview    # Preview only (save PPT + HTML, no send)
"""

import os
import sys
import json
import io
import re
import time
import base64
import struct
import subprocess
import tempfile
from datetime import datetime
from pathlib import Path
from xml.sax.saxutils import escape
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image, ImageChops

# Configuration
RECIPIENTS = ["Kendall.rush@walmart.com", "Matthew.Farnworth@walmart.com", "Justin.Barrick@walmart.com"]
SUBJECT = "TDA Initiative Insights - Weekly Report"
DASHBOARD_URL = "http://WEUS42608431466:5000/tda-initiatives-insights"
SCRIPT_DIR = Path(__file__).parent
# Walmart fiscal week from BQ Cal_Dim_Data table
def _walmart_week():
    from google.cloud import bigquery
    client = bigquery.Client(project='wmt-assetprotection-prod')
    query = """
    SELECT WM_WEEK_NBR, FISCAL_YEAR_NBR
    FROM `wmt-assetprotection-prod.Store_Support_Dev.Cal_Dim_Data`
    WHERE CALENDAR_DATE = CURRENT_DATE()
    LIMIT 1
    """
    rows = list(client.query(query).result())
    if rows:
        return rows[0]['WM_WEEK_NBR'], rows[0]['FISCAL_YEAR_NBR']
    # Fallback: manual calc if Cal_Dim_Data row missing
    from datetime import date, timedelta
    today = date.today()
    feb1 = date(today.year, 2, 1)
    days_to_sat = (5 - feb1.weekday()) % 7
    nearest_sat = feb1 + timedelta(days=days_to_sat - 7 if days_to_sat > 3 else days_to_sat)
    if today < nearest_sat:
        feb1_prev = date(today.year - 1, 2, 1)
        days_to_sat = (5 - feb1_prev.weekday()) % 7
        nearest_sat = feb1_prev + timedelta(days=days_to_sat - 7 if days_to_sat > 3 else days_to_sat)
    return (today - nearest_sat).days // 7 + 1, today.year

WEEK_NUM, FISCAL_YEAR = _walmart_week()
PPT_OUTPUT = SCRIPT_DIR / f"TDA_WK{WEEK_NUM}_Report.pptx"
PDF_OUTPUT = SCRIPT_DIR / f"TDA_WK{WEEK_NUM}_Report.pdf"
HTML_PREVIEW = SCRIPT_DIR / "email_preview.html"

# BQ Configuration
os.environ['GOOGLE_APPLICATION_CREDENTIALS'] = os.path.join(
    os.environ.get('APPDATA', ''), 'gcloud', 'application_default_credentials.json'
)
BQ_PROJECT = 'wmt-assetprotection-prod'
BQ_DATASET = 'Store_Support_Dev'
BQ_TABLE = 'Output- TDA Report'

EXCLUDED_PHASES = {'Complete'}
PHASE_ORDER = ['Pending', 'Vet', 'Test', 'Test Markets', 'Roll/Deploy']

SPARK_LOGO = Path(__file__).parent.parent.parent / "General Setup" / "Design" / "Spark Blank.png"
EDGE_PATH = r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe"

# ── Temporary data normalization (until BQ data reflects new names) ──
_PHASE_MAP = {'POC/POT': 'Vet', 'Mkt Scale': 'Test Markets'}
_OWNERSHIP_MAP = {'Dallas POC': 'Dallas VET'}

def _normalize_phase(phase):
    return _PHASE_MAP.get(phase, phase)

def _normalize_ownership(raw):
    own = raw or 'No Selection Provided'
    if own in ('No Selection Provided', '*Select Owner'):
        return 'No Selection Provided'
    return _OWNERSHIP_MAP.get(own, own)

# Colors
COLORS = {
    'walmart_blue': '#0071CE',
    'walmart_blue_dark': '#004C91',
    'walmart_yellow': '#FFC220',
    'on_track': '#107C10',
    'at_risk': '#F7630C',
    'off_track': '#DC3545',
    'continuous': '#6F42C1',
}


def fetch_data():
    """Fetch latest data from BigQuery"""
    from google.cloud import bigquery
    client = bigquery.Client(project=BQ_PROJECT)
    query = f"""
    SELECT
        Topic as `Initiative - Project Title`,
        Health_Update as `Health Status`,
        Phase,
        SUM(CASE WHEN Phase = Facility_Phase THEN Facility ELSE 0 END) as `# of Stores`,
        Dallas_POC as `Dallas POC`,
        Intake_n_Testing as `Intake & Testing`,
        Deployment,
        MAX(Intake_Card_Nbr) as `Project ID`,
        COALESCE(TDA_Ownership, 'No Selection Provided') as `TDA Ownership`
    FROM `{BQ_PROJECT}.{BQ_DATASET}.{BQ_TABLE}`
    WHERE Topic IS NOT NULL
    GROUP BY Topic, Health_Update, Phase, Intake_n_Testing, Dallas_POC, Deployment, TDA_Ownership
    ORDER BY Topic
    """
    results = client.query(query).result()
    data = []
    for row in results:
        phase = row['Phase'] or 'Unknown'
        if phase in EXCLUDED_PHASES:
            continue
        data.append({
            'Initiative - Project Title': row['Initiative - Project Title'] or 'Unknown',
            'Health Status': row['Health Status'] or 'Unknown',
            'Phase': _normalize_phase(phase),
            '# of Stores': row['# of Stores'] or 0,
            'Dallas VET': row['Dallas POC'] or 'N/A',
            'Intake & Testing': row['Intake & Testing'] or 'N/A',
            'Deployment': row['Deployment'] or 'N/A',
            'Project ID': row['Project ID'] or 0,
            'TDA Ownership': _normalize_ownership(row['TDA Ownership']),
        })
    return data


def build_email_html(data):
    """Build professional HTML email body with summary + phase tables"""
    today = datetime.now().strftime('%B %d, %Y')

    # Summary counts
    total = len(data)
    total_stores = sum(int(r.get('# of Stores', 0) or 0) for r in data)
    health_counts = {}
    phase_counts = {}
    for r in data:
        h = r['Health Status']
        p = r['Phase']
        health_counts[h] = health_counts.get(h, 0) + 1
        phase_counts[p] = phase_counts.get(p, 0) + 1

    on_track = health_counts.get('On Track', 0)
    at_risk = health_counts.get('At Risk', 0)
    off_track = health_counts.get('Off Track', 0)
    continuous = health_counts.get('Continuous', 0)

    def badge_style(status):
        s = status.lower()
        if 'on track' in s:
            return f"background-color:rgba(16,124,16,0.15);color:{COLORS['on_track']};padding:3px 10px;border-radius:12px;font-weight:600;font-size:12px;display:inline-block;"
        elif 'at risk' in s:
            return f"background-color:rgba(247,99,12,0.15);color:{COLORS['at_risk']};padding:3px 10px;border-radius:12px;font-weight:600;font-size:12px;display:inline-block;"
        elif 'off track' in s:
            return f"background-color:rgba(220,53,69,0.15);color:{COLORS['off_track']};padding:3px 10px;border-radius:12px;font-weight:600;font-size:12px;display:inline-block;"
        elif 'continuous' in s:
            return f"background-color:rgba(111,66,193,0.15);color:{COLORS['continuous']};padding:3px 10px;border-radius:12px;font-weight:600;font-size:12px;display:inline-block;"
        return "padding:3px 10px;border-radius:12px;font-weight:600;font-size:12px;display:inline-block;"

    total_stores_display = f'{total_stores:,}'

    # Spark logo as base64 for email embedding
    spark_b64 = ''
    if SPARK_LOGO.exists():
        spark_b64 = base64.b64encode(SPARK_LOGO.read_bytes()).decode('ascii')

    spark_img = f'<img src="data:image/png;base64,{spark_b64}" width="44" height="44" alt="Spark" style="display:block;"/>' if spark_b64 else '&#10058;'

    html = f"""<!DOCTYPE html>
<html>
<head><meta charset="utf-8"></head>
<body style="font-family: 'Segoe UI', Arial, sans-serif; color: #333; max-width: 1000px; margin: 0 auto; padding: 20px;">

<!-- Header -->
<table width="100%" cellpadding="0" cellspacing="0" style="margin-bottom: 24px;">
<tr>
<td bgcolor="{COLORS['walmart_blue_dark']}" style="background-color:{COLORS['walmart_blue_dark']}; padding: 24px 30px;">
    <table><tr>
    <td style="padding-right: 15px; vertical-align: middle;">
        {spark_img}
    </td>
    <td>
        <div style="color: white; font-size: 26px; font-weight: 700;">TDA Initiative Insights</div>
        <div style="color: #cccccc; font-size: 13px; margin-top: 2px;">Weekly Report &mdash; {today}</div>
    </td>
    </tr></table>
</td>
</tr>
</table>

<!-- Dashboard Button -->
<table width="100%" cellpadding="0" cellspacing="0" style="margin-bottom: 24px;">
<tr>
<td align="center" style="padding: 16px 0;">
    <table cellpadding="0" cellspacing="0" style="border-collapse: separate;">
    <tr>
    <td align="center" bgcolor="#0052CC" style="background-color: #0052CC; border-radius: 30px; mso-padding-alt: 14px 40px;">
        <a href="{DASHBOARD_URL}" target="_blank" style="background-color: #0052CC; color: #FFFFFF; font-family: 'Segoe UI', Arial, sans-serif; font-size: 16px; font-weight: 700; text-decoration: none; padding: 14px 40px; border-radius: 30px; display: inline-block; letter-spacing: 0.5px; mso-padding-alt: 0; border: 2px solid #0052CC;">Go to Dashboard</a>
    </td>
    </tr>
    </table>
</td>
</tr>
</table>

<!-- Context -->
<table width="100%" cellpadding="0" cellspacing="0" style="margin-bottom: 24px;"><tr>
<td bgcolor="#f8f9fa" style="background-color:#f8f9fa; border: 1px solid #e5e5e5; padding: 16px 20px; font-size: 13px; line-height: 1.6; color: #444;">
    <strong style="color: {COLORS['walmart_blue_dark']};">About This Report</strong><br>
    This is your weekly TDA (Technology Deployment Acceleration) Initiative Insights report. It provides a snapshot of all active projects currently tracked through the Intake Hub, filtered to Implement, Validation, and Intake &amp; Test branches. Projects are grouped by their current phase in the deployment pipeline. Each project title links directly to its Intake Hub card for full details. A PowerPoint version is attached for offline review and sharing.
</td></tr></table>

<!-- Summary Cards -->
<table width="100%" cellpadding="0" cellspacing="0" style="margin-bottom: 24px;">
<tr>
    <td width="16%" style="padding: 0 4px;">
        <div style="background:white; border:1px solid #e5e5e5; border-left:4px solid {COLORS['walmart_yellow']}; border-radius:6px; padding:12px; text-align:center;">
            <div style="font-size:11px; color:#666; text-transform:uppercase; letter-spacing:0.5px;">Total Initiatives</div>
            <div style="font-size:24px; font-weight:700; color:{COLORS['walmart_blue_dark']};">{total}</div>
        </div>
    </td>
    <td width="16%" style="padding: 0 4px;">
        <div style="background:white; border:1px solid #e5e5e5; border-left:4px solid {COLORS['walmart_yellow']}; border-radius:6px; padding:12px; text-align:center;">
            <div style="font-size:11px; color:#666; text-transform:uppercase; letter-spacing:0.5px;">Total Stores</div>
            <div style="font-size:24px; font-weight:700; color:{COLORS['walmart_blue_dark']};">{total_stores_display}</div>
        </div>
    </td>
    <td width="16%" style="padding: 0 4px;">
        <div style="background:white; border:1px solid #e5e5e5; border-left:4px solid {COLORS['on_track']}; border-radius:6px; padding:12px; text-align:center;">
            <div style="font-size:11px; color:#666; text-transform:uppercase; letter-spacing:0.5px;">On Track</div>
            <div style="font-size:24px; font-weight:700; color:{COLORS['on_track']};">{on_track}</div>
        </div>
    </td>
    <td width="16%" style="padding: 0 4px;">
        <div style="background:white; border:1px solid #e5e5e5; border-left:4px solid {COLORS['at_risk']}; border-radius:6px; padding:12px; text-align:center;">
            <div style="font-size:11px; color:#666; text-transform:uppercase; letter-spacing:0.5px;">At Risk</div>
            <div style="font-size:24px; font-weight:700; color:{COLORS['at_risk']};">{at_risk}</div>
        </div>
    </td>
    <td width="16%" style="padding: 0 4px;">
        <div style="background:white; border:1px solid #e5e5e5; border-left:4px solid {COLORS['off_track']}; border-radius:6px; padding:12px; text-align:center;">
            <div style="font-size:11px; color:#666; text-transform:uppercase; letter-spacing:0.5px;">Off Track</div>
            <div style="font-size:24px; font-weight:700; color:{COLORS['off_track']};">{off_track}</div>
        </div>
    </td>
    <td width="16%" style="padding: 0 4px;">
        <div style="background:white; border:1px solid #e5e5e5; border-left:4px solid {COLORS['continuous']}; border-radius:6px; padding:12px; text-align:center;">
            <div style="font-size:11px; color:#666; text-transform:uppercase; letter-spacing:0.5px;">Continuous</div>
            <div style="font-size:24px; font-weight:700; color:{COLORS['continuous']};">{continuous}</div>
        </div>
    </td>
</tr>
</table>
"""

    # Phase tables — grouped by TDA Ownership → Phase
    UNASSIGNED = 'No Selection Provided'

    # Build ordered sections: Pending+unassigned first, then custom ownership order
    OWNERSHIP_ORDER = ['Dallas VET', 'Intake & Test', 'Deployment', UNASSIGNED]
    sections = []
    pending_unassigned = [r for r in data if r['Phase'] == 'Pending' and r.get('TDA Ownership', UNASSIGNED) == UNASSIGNED]
    if pending_unassigned:
        sections.append(('TDA Ownership - Currently No TDA Ownership', 'Pending', pending_unassigned, False))

    # Ordered ownership list: defined order first, then any new values alphabetically
    all_ownerships = set(r.get('TDA Ownership', UNASSIGNED) or UNASSIGNED for r in data)
    known_set = set(OWNERSHIP_ORDER)
    unknown_ownerships = sorted(o for o in all_ownerships if o not in known_set)
    ownership_list = [o for o in OWNERSHIP_ORDER if o in all_ownerships] + unknown_ownerships
    for ownership in ownership_list:
        for phase in PHASE_ORDER:
            if phase == 'Pending' and ownership == UNASSIGNED:
                continue  # already handled above
            phase_data = [r for r in data if (r.get('TDA Ownership', UNASSIGNED) or UNASSIGNED) == ownership and r['Phase'] == phase]
            if not phase_data:
                continue
            label = 'TDA Ownership - Currently No TDA Ownership' if ownership == UNASSIGNED else ownership
            sections.append((label, phase, phase_data, False))

    for ownership_label, phase, phase_data, _unused in sections:
        count = len(phase_data)
        # Ownership banner (navy)
        html += f"""
<!-- {ownership_label} / {phase} Section -->
<div style="background:#1E3A8A; color:white; padding:10px 16px; font-size:15px; font-weight:700; border-radius:6px 6px 0 0; margin-top:16px;">
    {escape(ownership_label)} &mdash; {count} Project{'s' if count != 1 else ''}
</div>
"""
        # Phase sub-banner (blue, half height)
        html += f"""<div style="background:{COLORS['walmart_blue']}; color:white; padding:5px 16px; font-size:13px; font-weight:600;">
    {escape(phase)}
</div>
"""
        html += f"""<table width="100%" cellpadding="0" cellspacing="0" style="border:1px solid #e5e5e5; border-top:none; border-collapse:collapse; margin-bottom:8px; font-size:14px;">
<thead>
<tr style="background:#f5f5f5;">
    <th style="padding:10px 12px; text-align:left; border-bottom:2px solid #ddd; font-weight:700; color:#333; font-size:14px;">Initiative - Project Title</th>
    <th style="padding:10px 12px; text-align:left; border-bottom:2px solid #ddd; font-weight:700; color:#333; font-size:14px; width:100px;">Health</th>
    <th style="padding:10px 12px; text-align:center; border-bottom:2px solid #ddd; font-weight:700; color:#333; font-size:14px; width:80px;"># of Stores</th>
</tr>
</thead>
<tbody>
"""
        for i, row in enumerate(phase_data):
            bg = '#ffffff' if i % 2 == 0 else '#fafafa'
            title = escape(row['Initiative - Project Title'])
            pid = row.get('Project ID', 0)
            link = f'https://hoops.wal-mart.com/intake-hub/projects/{pid}' if pid else '#'
            health = escape(row['Health Status'])
            stores = int(row.get('# of Stores', 0) or 0)
            stores_display = f'{stores:,}'
            bs = badge_style(row['Health Status'])

            html += f"""<tr style="background:{bg};">
    <td style="padding:10px 12px; border-bottom:1px solid #eee;"><a href="{link}" style="color:{COLORS['walmart_blue']}; text-decoration:none; font-weight:500; font-size:14px;">{title}</a></td>
    <td style="padding:10px 12px; border-bottom:1px solid #eee;"><span style="{bs}">{health}</span></td>
    <td style="padding:10px 12px; border-bottom:1px solid #eee; text-align:center; font-weight:600; color:{COLORS['walmart_blue']}; font-size:14px;">{stores_display}</td>
</tr>
"""
        html += "</tbody></table>\n"

    # Footer
    html += f"""
<div style="margin-top:24px; padding:16px; border-top:2px solid #e5e5e5; text-align:center; color:#999; font-size:11px;">
    <p>This report was automatically generated from the TDA Initiative Insights Dashboard.<br>
    PPT report is attached for offline review.</p>
    <p style="margin-top:8px;">Generated on {today} at {datetime.now().strftime('%I:%M %p')} CT</p>
</div>

</body>
</html>"""
    return html


def _build_phase_html(phase, rows, ownership=None, is_special=False):
    """Build the EXACT same HTML the dashboard builds for PPT screenshots."""
    columns = ['Initiative - Project Title', 'Health Status', 'Phase',
                '# of Stores', 'Dallas VET', 'Intake & Testing', 'Deployment']

    # Ownership banner (navy) — phase sub-banner removed to match dashboard
    ownership_label = ownership or phase
    banner = f'<div style="background:#1E3A8A;color:white;padding:15px 20px;font-size:18px;font-weight:700;text-align:center;border-radius:4px 4px 0 0;">{escape(ownership_label)}</div>'

    table = '<table style="width:100%;border-collapse:collapse;font-size:14px;font-family:-apple-system,BlinkMacSystemFont,\'Segoe UI\',\'Helvetica Neue\',Arial,sans-serif;">'
    table += '<thead style="background-color:#F5F5F5;border-bottom:2px solid #E5E5E5;"><tr>'
    for col in columns:
        table += f'<th style="padding:16px 24px;text-align:left;font-weight:600;color:#212121;white-space:nowrap;">{escape(col)}</th>'
    table += '</tr></thead><tbody>'

    for row in rows:
        table += '<tr>'
        for col in columns:
            value = str(row.get(col, '') or '')
            if col == '# of Stores':
                value = str(int(row.get(col, 0) or 0))
            cell_style = 'padding:16px 24px;border-bottom:1px solid #E5E5E5;text-align:left;'

            if col == 'Health Status':
                status = value.lower()
                bg, clr = '#f0f0f0', '#333'
                if 'on track' in status:
                    bg, clr = 'rgba(16,124,16,0.2)', '#107C10'
                elif 'at risk' in status:
                    bg, clr = 'rgba(247,99,12,0.2)', '#F7630C'
                elif 'off track' in status:
                    bg, clr = 'rgba(220,53,69,0.2)', '#DC3545'
                elif 'continuous' in status:
                    bg, clr = 'rgba(111,66,193,0.2)', '#6F42C1'
                table += f'<td style="{cell_style}"><span style="display:inline-block;padding:4px 10px;background-color:{bg};color:{clr};border-radius:12px;font-weight:600;font-size:12px;white-space:nowrap;">{escape(value)}</span></td>'
            elif col == 'Phase':
                table += f'<td style="{cell_style}"><span style="display:inline-block;padding:4px 10px;background-color:#DBEAFE;color:#1E3A8A;border-radius:4px;font-weight:500;font-size:12px;">{escape(value)}</span></td>'
            elif col == '# of Stores':
                table += f'<td style="{cell_style}font-weight:600;color:#0071CE;">{escape(value)}</td>'
            elif col == 'Initiative - Project Title':
                pid = row.get('Project ID', 0)
                if pid:
                    table += f'<td style="{cell_style}"><a href="https://hoops.wal-mart.com/intake-hub/projects/{pid}" style="color:#0071CE;text-decoration:none;font-weight:500;">{escape(value)}</a></td>'
                else:
                    table += f'<td style="{cell_style}">{escape(value)}</td>'
            else:
                table += f'<td style="{cell_style}">{escape(value)}</td>'
        table += '</tr>'
    table += '</tbody></table>'

    full_html = f'''<!DOCTYPE html><html><head><meta charset="utf-8">
<style>*{{margin:0;padding:0;box-sizing:border-box;}} body{{background:white;}}</style>
</head><body>{banner}{table}</body></html>'''
    return full_html


def _measure_row_heights(phase, rows):
    """Use Edge --dump-dom to measure row heights with same styling as dashboard's packRowsIntoPages."""
    columns = ['Initiative - Project Title', 'Health Status', 'Phase',
                '# of Stores', 'Dallas VET', 'Intake & Testing', 'Deployment']

    # Build measurement HTML that matches dashboard's measurement container
    # Dashboard uses padding:8px, border:1px solid #ddd, word-wrap:break-word — NOT the render padding
    table = '<table style="width:100%;border-collapse:collapse;font-size:14px;font-family:Segoe UI,Arial,sans-serif;line-height:1.5;">'
    table += '<thead><tr style="background-color:#f0f0f0;border-bottom:2px solid #999;height:50px;">'
    for col in columns:
        table += f'<th style="padding:8px;text-align:left;border:1px solid #ddd;">{escape(col)}</th>'
    table += '</tr></thead><tbody>'
    for row in rows:
        table += '<tr style="border:1px solid #ddd;">'
        for col in columns:
            value = str(row.get(col, '') or '')
            if col == '# of Stores':
                value = str(int(row.get(col, 0) or 0))
            style = 'padding:8px;border:1px solid #ddd;word-wrap:break-word;white-space:normal;'
            if col in ('Health Status', 'Phase', '# of Stores'):
                style += 'font-weight:600;'
            table += f'<td style="{style}">{escape(value)}</td>'
        table += '</tr>'
    table += '</tbody></table>'

    meas_html = f'''<!DOCTYPE html><html><head><meta charset="utf-8">
<style>*{{margin:0;padding:0;box-sizing:border-box;}}</style>
</head><body style="width:1280px;">
{table}
<script>
window.onload=function(){{
  var rs=document.querySelectorAll("tbody tr");
  var h=[];rs.forEach(function(r){{h.push(r.offsetHeight);}});
  document.getElementById("_rh").textContent="ROW_HEIGHTS:"+JSON.stringify(h);
}};
</script><pre id="_rh"></pre>
</body></html>'''

    with tempfile.NamedTemporaryFile(suffix='.html', delete=False, mode='w', encoding='utf-8') as f:
        f.write(meas_html)
        html_path = f.name
    try:
        cmd = [
            EDGE_PATH, '--headless', '--disable-gpu', '--no-sandbox',
            '--dump-dom', '--virtual-time-budget=5000',
            f'--window-size=1280,4000',
            f'file:///{html_path.replace(os.sep, "/")}'
        ]
        result = subprocess.run(cmd, capture_output=True, timeout=30)
        stdout = result.stdout.decode('utf-8', errors='replace')
        match = re.search(r'ROW_HEIGHTS:\[([\d,]+)\]', stdout)
        if match:
            return [int(x) for x in match.group(1).split(',')]
    finally:
        os.unlink(html_path)
    # Fallback: assume 55px per row
    return [55] * len(rows)


def _paginate_by_height(rows, heights):
    """Replicate dashboard packRowsIntoPages — split rows when cumulative height > AVAILABLE_HEIGHT."""
    AVAILABLE_HEIGHT = 800  # calibrated to match dashboard's packRowsIntoPages via Edge headless
    pages = []
    current_page = []
    current_h = 0
    for row, h in zip(rows, heights):
        if current_h + h > AVAILABLE_HEIGHT and current_page:
            pages.append(current_page)
            current_page = [row]
            current_h = h
        else:
            current_page.append(row)
            current_h += h
    if current_page:
        pages.append(current_page)
    return pages


def _capture_html_screenshot(html_content, output_png, width=1280):
    """Use Edge headless to capture HTML as a PNG screenshot."""
    with tempfile.NamedTemporaryFile(suffix='.html', delete=False, mode='w', encoding='utf-8') as f:
        f.write(html_content)
        html_path = f.name

    try:
        cmd = [
            EDGE_PATH,
            '--headless',
            '--disable-gpu',
            '--no-sandbox',
            '--hide-scrollbars',
            '--force-device-scale-factor=1.5',
            f'--screenshot={output_png}',
            f'--window-size={width},2000',
            f'file:///{html_path.replace(os.sep, "/")}'
        ]
        result = subprocess.run(cmd, capture_output=True, timeout=30)
        if not Path(output_png).exists():
            raise RuntimeError(f"Edge screenshot failed: {result.stderr.decode(errors='replace')[:200]}")

        # Auto-crop whitespace — match html2canvas element-only capture
        img = Image.open(output_png)
        bg = Image.new(img.mode, img.size, (255, 255, 255))
        diff = ImageChops.difference(img, bg)
        bbox = diff.getbbox()
        if bbox:
            # Crop exactly to content boundary — no padding on right
            cropped = img.crop((0, 0, bbox[2], bbox[3] + 4))
            # Save as JPEG for smaller file size
            jpg_path = output_png.replace('.png', '.jpg')
            cropped.convert('RGB').save(jpg_path, 'JPEG', quality=78, optimize=True)
            import shutil
            shutil.move(jpg_path, output_png)  # keep .png extension for caller
    finally:
        os.unlink(html_path)


def generate_report_pptx(data):
    """Generate PPTX using Edge headless screenshots — identical to dashboard Generate PPT."""
    print("    Using Edge headless for screenshot-based PPT (matches dashboard)...")

    # Build ownership → phase sections (same logic as dashboard & email)
    UNASSIGNED = 'No Selection Provided'
    sections = []

    # 1. Pending + unassigned first
    pending_unassigned = [r for r in data if r['Phase'] == 'Pending' and r.get('TDA Ownership', UNASSIGNED) == UNASSIGNED]
    if pending_unassigned:
        sections.append(('TDA Ownership - Currently No TDA Ownership', 'Pending', pending_unassigned))

    # 2. Custom ownership order: Dallas VET, Intake & Test, Deployment, then unassigned
    OWNERSHIP_ORDER = ['Dallas VET', 'Intake & Test', 'Deployment', UNASSIGNED]
    all_ownerships = set(r.get('TDA Ownership', UNASSIGNED) or UNASSIGNED for r in data)
    known_set = set(OWNERSHIP_ORDER)
    unknown_ownerships = sorted(o for o in all_ownerships if o not in known_set)
    ownership_list = [o for o in OWNERSHIP_ORDER if o in all_ownerships] + unknown_ownerships
    for ownership in ownership_list:
        for phase in PHASE_ORDER:
            if phase == 'Pending' and ownership == UNASSIGNED:
                continue
            rows = [r for r in data if (r.get('TDA Ownership', UNASSIGNED) or UNASSIGNED) == ownership and r['Phase'] == phase]
            if not rows:
                continue
            label = 'TDA Ownership - Currently No TDA Ownership' if ownership == UNASSIGNED else ownership
            sections.append((label, phase, rows))

    screenshots = []  # list of (label, png_bytes)

    # Combine phases within same ownership (phases stay in order, packed by height)
    from collections import OrderedDict
    ownership_groups = OrderedDict()
    for ownership_label, phase, rows in sections:
        if ownership_label not in ownership_groups:
            ownership_groups[ownership_label] = []
        ownership_groups[ownership_label].extend(rows)

    with tempfile.TemporaryDirectory() as tmp_dir:
        for ownership_label, all_rows in ownership_groups.items():
            if not all_rows:
                continue
            # Measure actual row heights via Edge (matches dashboard packRowsIntoPages)
            heights = _measure_row_heights(ownership_label, all_rows)
            pages = _paginate_by_height(all_rows, heights)
            total_pages = len(pages)
            for page, page_rows in enumerate(pages):
                label = ownership_label
                if total_pages > 1:
                    label += f' ({page + 1}/{total_pages})'

                html = _build_phase_html(None, page_rows, ownership=ownership_label)
                png_path = os.path.join(tmp_dir, f'slide_{len(screenshots) + 1}.png')
                _capture_html_screenshot(html, png_path)
                png_bytes = Path(png_path).read_bytes()
                screenshots.append((label, png_bytes))
                print(f"    Captured: {ownership_label} page {page + 1}/{total_pages} ({len(png_bytes):,} bytes)")

    # Build the PPTX with python-pptx (schema-compliant, no Repair dialog)
    SLIDE_WIDTH = Inches(10)
    SLIDE_HEIGHT = Inches(7.5)

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # ── Slide 1: Title slide (white bg, blue header, yellow accent) ──
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Blue header bar
    header_shape = slide.shapes.add_shape(1, Emu(0), Emu(0), SLIDE_WIDTH, Inches(1.2))
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = RGBColor(0x3B, 0x82, 0xF6)
    header_shape.line.fill.background()
    tf = header_shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.5)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = "TDA Initiatives Insights"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    run.font.name = "Segoe UI"

    # Main title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(9), Inches(1.8))
    p = txBox.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Initiative Status Insights"
    run.font.size = Pt(40)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1E, 0x3A, 0x5F)
    run.font.name = "Segoe UI"

    # Subtitle
    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(3.9), Inches(9), Inches(0.8))
    p2 = txBox2.text_frame.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = "Store Support  |  Asset Protection"
    run2.font.size = Pt(20)
    run2.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    run2.font.name = "Segoe UI"

    # Yellow accent bar at bottom
    accent = slide.shapes.add_shape(1, Emu(0), Inches(7.1), SLIDE_WIDTH, Inches(0.4))
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGBColor(0xFF, 0xC2, 0x20)
    accent.line.fill.background()

    # ── Screenshot slides ──
    for label, png_bytes in screenshots:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        img = Image.open(io.BytesIO(png_bytes))
        img_w, img_h = img.size
        aspect = img_h / img_w if img_w > 0 else 0.5625
        final_w = SLIDE_WIDTH
        final_h = Emu(int(SLIDE_WIDTH * aspect))
        if final_h > SLIDE_HEIGHT:
            final_h = SLIDE_HEIGHT
        slide.shapes.add_picture(io.BytesIO(png_bytes), Emu(0), Emu(0), final_w, final_h)

    pptx_buffer = io.BytesIO()
    prs.save(pptx_buffer)
    pptx_buffer.seek(0)
    return pptx_buffer.getvalue(), screenshots


def generate_report_pdf(pptx_screenshots):
    """Convert screenshot images to a multi-page PDF using Pillow."""
    if not pptx_screenshots:
        return None
    pages = []
    for label, png_bytes in pptx_screenshots:
        img = Image.open(io.BytesIO(png_bytes)).convert('RGB')
        pages.append(img)
    pdf_buffer = io.BytesIO()
    pages[0].save(pdf_buffer, 'PDF', save_all=True, append_images=pages[1:], resolution=150)
    pdf_buffer.seek(0)
    return pdf_buffer.getvalue()


def send_outlook_email(html_body, attachments, recipients):
    """Send email via Outlook with PPT attachment"""
    import win32com.client
    outlook = win32com.client.Dispatch('Outlook.Application')
    mail = outlook.CreateItem(0)  # 0 = olMailItem
    mail.Subject = SUBJECT
    mail.HTMLBody = html_body
    mail.To = '; '.join(recipients)
    for att in attachments:
        mail.Attachments.Add(str(att))
    mail.Send()
    print(f"[OK] Email sent to: {', '.join(recipients)}")


def main():
    preview_only = '--preview' in sys.argv

    print("=" * 60)
    print("TDA Initiative Insights - Weekly Email Report")
    print("=" * 60)

    # 1. Fetch data
    print("\n[1/4] Fetching data from BigQuery...")
    data = fetch_data()
    print(f"  Loaded {len(data)} projects")

    # 2. Generate PPT + PDF
    print("[2/5] Generating PPT report...")
    pptx_data, screenshots = generate_report_pptx(data)
    PPT_OUTPUT.write_bytes(pptx_data)
    print(f"  Saved: {PPT_OUTPUT}")
    print(f"  Size: {len(pptx_data):,} bytes")

    print("[3/5] Generating PDF report...")
    pdf_data = generate_report_pdf(screenshots)
    if pdf_data:
        PDF_OUTPUT.write_bytes(pdf_data)
        print(f"  Saved: {PDF_OUTPUT}")
        print(f"  Size: {len(pdf_data):,} bytes")

    # 4. Build email HTML
    print("[4/5] Building email HTML...")
    html = build_email_html(data)

    if preview_only:
        HTML_PREVIEW.write_text(html, encoding='utf-8')
        print(f"\n  PREVIEW MODE - files saved:")
        print(f"  PPT:   {PPT_OUTPUT}")
        print(f"  PDF:   {PDF_OUTPUT}")
        print(f"  HTML:  {HTML_PREVIEW}")
        print(f"\n  Open {HTML_PREVIEW} in a browser to review the email.")
        print("  Run without --preview to send.")
        return

    # 5. Send email with PPT + PDF attachments
    print(f"[5/5] Sending email to {', '.join(RECIPIENTS)}...")
    attachments = [PPT_OUTPUT]
    if PDF_OUTPUT.exists():
        attachments.append(PDF_OUTPUT)
    send_outlook_email(html, attachments, RECIPIENTS)

    print("\n" + "=" * 60)
    print("Done!")
    print("=" * 60)


if __name__ == '__main__':
    main()
