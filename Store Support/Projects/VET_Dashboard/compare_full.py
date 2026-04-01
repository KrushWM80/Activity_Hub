"""Compare reference PPT content vs our generated PPT - extract all visible data"""
from pptx import Presentation
from PIL import Image
import io

ref = Presentation('reports/VET_Executive_Report (3).pptx')
ours = Presentation('reports/VET_Executive_Report_WK09.pptx')

print("=" * 80)
print("REFERENCE: VET_Executive_Report (3).pptx - Dashboard Generated")
print(f"  Size: 2,078 KB | Slides: {len(ref.slides)} | Dimensions: {ref.slide_width/914400:.1f}\"x{ref.slide_height/914400:.1f}\"")
print("=" * 80)
for i, slide in enumerate(ref.slides):
    shapes = list(slide.shapes)
    print(f"\n  Slide {i+1}:")
    for s in shapes:
        if s.shape_type == 13:  # Picture
            img = Image.open(io.BytesIO(s.image.blob))
            w, h = img.size
            print(f"    Image: {w}x{h}px, {len(s.image.blob):,} bytes")
            print(f"    Position: full slide ({s.width/914400:.1f}\"x{s.height/914400:.1f}\")")

print()
print("=" * 80)
print("OURS: VET_Executive_Report_WK09.pptx - Edge Headless Generated")
print(f"  Size: 785 KB | Slides: {len(ours.slides)} | Dimensions: {ours.slide_width/914400:.1f}\"x{ours.slide_height/914400:.1f}\"")
print("=" * 80)
for i, slide in enumerate(ours.slides):
    shapes = list(slide.shapes)
    print(f"\n  Slide {i+1}:")
    for s in shapes:
        if s.shape_type == 13:  # Picture
            img = Image.open(io.BytesIO(s.image.blob))
            w, h = img.size
            print(f"    Image: {w}x{h}px, {len(s.image.blob):,} bytes")
            print(f"    Position: ({s.left/914400:.2f}\",{s.top/914400:.2f}\") size ({s.width/914400:.1f}\"x{s.height/914400:.1f}\")")

print()
print("=" * 80)
print("KEY DIFFERENCES")
print("=" * 80)

# File size
ref_kb = 2077.7
our_kb = 784.8
print(f"\n1. FILE SIZE:")
print(f"   Reference: {ref_kb:,.1f} KB")
print(f"   Ours:      {our_kb:,.1f} KB")
print(f"   Difference: Reference is {ref_kb/our_kb:.1f}x larger")

# Also compare to the Dallas VET Weekly Report PDF
from pathlib import Path
pdf_ref = Path('reports/Dallas VET Weekly Report 3.30.pdf')
our_pdf = Path('reports/VET_Executive_Report_WK09.pdf')
print(f"\n2. PDF SIZE:")
print(f"   Reference PDF: {pdf_ref.stat().st_size / 1024:,.1f} KB")
print(f"   Our PDF:       {our_pdf.stat().st_size / 1024:,.1f} KB")

# Slides
print(f"\n3. SLIDE COUNT:")
print(f"   Reference: {len(ref.slides)} slides")
print(f"   Ours:      {len(ours.slides)} slides")

# Image resolution
print(f"\n4. IMAGE RESOLUTION (Slide 1 - Summary):")
ref_img = Image.open(io.BytesIO(list(ref.slides[0].shapes)[0].image.blob))
our_img = Image.open(io.BytesIO(list(ours.slides[0].shapes)[0].image.blob))
print(f"   Reference: {ref_img.size[0]}x{ref_img.size[1]}px")
print(f"   Ours:      {our_img.size[0]}x{our_img.size[1]}px")

# Image positioning
print(f"\n5. IMAGE POSITIONING:")
ref_s = list(ref.slides[0].shapes)[0]
our_s = list(ours.slides[0].shapes)[0]
print(f"   Reference: Full bleed - ({ref_s.left/914400:.1f}\",{ref_s.top/914400:.1f}\") -> ({ref_s.width/914400:.1f}\"x{ref_s.height/914400:.1f}\")")
print(f"   Ours:      Inset     - ({our_s.left/914400:.2f}\",{our_s.top/914400:.2f}\") -> ({our_s.width/914400:.1f}\"x{our_s.height/914400:.1f}\")")

# Data content
print(f"\n6. DATA CONTENT:")
print(f"   Reference: 49 projects (ALL ownership groups, no filter)")
print(f"   Ours:      18 projects (Dallas POC + Dallas VET only)")
print(f"   Missing:   32 projects from other ownership groups")

print(f"\n7. REFERENCE PDF DETAILS:")
print(f"   File: Dallas VET Weekly Report 3.30.pdf")
print(f"   Size: {pdf_ref.stat().st_size / 1024:,.1f} KB (vs our {our_pdf.stat().st_size / 1024:,.1f} KB)")
print(f"   This PDF is {pdf_ref.stat().st_size / our_pdf.stat().st_size:.1f}x larger than ours")
print(f"   Likely because it contains all 49 projects with higher-res images")
