"""
TDA Dashboard PPT Generation API Service
Integrates PowerPoint generation with the Flask backend
"""

from flask import Blueprint, jsonify, request, send_file
import logging
from datetime import datetime
from pathlib import Path
import os
import base64
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches
from PIL import Image

logger = logging.getLogger(__name__)


def get_walmart_week():
    """
    Calculate current Walmart fiscal week (WK01-WK13)
    Walmart fiscal year starts Feb 1 with 13 periods of 28 days each
    """
    today = datetime.now()
    
    # Walmart fiscal year starts Feb 1
    fiscal_year_start = datetime(today.year if today.month >= 2 else today.year - 1, 2, 1)
    
    # Days since start of fiscal year
    days_elapsed = (today - fiscal_year_start).days
    
    # Each Walmart period is 28 days
    # WK01 = days 0-27, WK02 = days 28-55, etc
    week_num = (days_elapsed // 28) + 1
    
    # Ensure week is between 1-13
    week_num = max(1, min(13, week_num))
    
    return f"WK{week_num:02d}"


# Create Blueprint
ppt_bp = Blueprint('ppt', __name__, url_prefix='/api/ppt')

# Configuration
OUTPUT_DIR = Path(__file__).parent / 'reports'
OUTPUT_DIR.mkdir(exist_ok=True)

# Try to import optional dependencies
try:
    from generate_ppt import TDAPowerPointGenerator
    from google.cloud import bigquery
    bq_client = bigquery.Client()
    PPT_AVAILABLE = True
except Exception as e:
    logger.warning(f"PPT service not fully available: {e}")
    PPT_AVAILABLE = False
    bq_client = None


class ReportManager:
    """Manages PPT report generation and caching"""
    
    def __init__(self):
        if PPT_AVAILABLE:
            self.generator = TDAPowerPointGenerator(bq_client)
        else:
            self.generator = None
        self.last_report_path = None
        self.last_report_time = None
    
    def generate_report(self, phases_filter=None, force_regenerate=False):
        """
        Generate a PowerPoint report
        
        Args:
            phases_filter: List of phases to include (None = all)
            force_regenerate: Force generation even if recent report exists
        
        Returns:
            tuple: (file_path, success: bool, message: str)
        """
        try:
            # Check if we have a recent report (less than 5 minutes old)
            if (self.last_report_path and 
                self.last_report_time and
                not force_regenerate):
                age = (datetime.now() - self.last_report_time).total_seconds()
                if age < 300 and Path(self.last_report_path).exists():
                    logger.info(f"Using cached report from {age:.0f}s ago")
                    return self.last_report_path, True, f"Using cached report (generated {age:.0f}s ago)"
            
            # Fetch fresh data
            logger.info("Fetching data for report generation...")
            self.generator.fetch_data()
            
            if not self.generator.data:
                return None, False, "No data available for report generation"
            
            # Filter data if requested
            if phases_filter:
                self.generator.data = [
                    row for row in self.generator.data
                    if str(row.get('Phase', '')).lower() in [p.lower() for p in phases_filter]
                ]
                logger.info(f"Filtered data to {len(self.generator.data)} records for phases: {phases_filter}")
            
            # Generate report
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            output_file = OUTPUT_DIR / f"VET_Executive_Report_{timestamp}.pptx"
            
            logger.info(f"Generating PowerPoint report: {output_file}")
            self.generator.generate_report(str(output_file))
            
            # Cache the path
            self.last_report_path = str(output_file)
            self.last_report_time = datetime.now()
            
            logger.info(f"Report generated successfully: {output_file}")
            return str(output_file), True, f"Report generated successfully"
        
        except Exception as e:
            logger.error(f"Error generating report: {e}")
            return None, False, f"Error: {str(e)}"


# Initialize report manager
report_manager = ReportManager()


# API Routes

@ppt_bp.route('/generate', methods=['POST'])
def generate_ppt():
    """
    Generate a PowerPoint report
    
    Request body:
    {
        "phases": ["Test", "Production"],  // Optional
        "force_regenerate": false          // Optional
    }
    """
    try:
        data = request.get_json() or {}
        phases = data.get('phases')
        force_regen = data.get('force_regenerate', False)
        
        logger.info(f"PPT generation requested. Phases: {phases}, Force: {force_regen}")
        
        file_path, success, message = report_manager.generate_report(
            phases_filter=phases,
            force_regenerate=force_regen
        )
        
        if success:
            return jsonify({
                'success': True,
                'message': message,
                'file_path': str(file_path),
                'file_name': Path(file_path).name,
                'timestamp': datetime.now().isoformat()
            })
        else:
            return jsonify({
                'success': False,
                'message': message,
                'timestamp': datetime.now().isoformat()
            }), 500
    
    except Exception as e:
        logger.error(f"Error in /api/ppt/generate: {e}")
        return jsonify({
            'success': False,
            'message': f'Server error: {str(e)}',
            'timestamp': datetime.now().isoformat()
        }), 500


@ppt_bp.route('/download/<filename>', methods=['GET'])
def download_ppt(filename):
    """Download a generated PPT file"""
    try:
        # Security: validate filename to prevent directory traversal
        if '..' in filename or '/' in filename or '\\' in filename:
            return jsonify({'error': 'Invalid filename'}), 400
        
        file_path = OUTPUT_DIR / filename
        
        if not file_path.exists():
            return jsonify({'error': 'File not found'}), 404
        
        logger.info(f"Downloading report: {filename}")
        
        return send_file(
            file_path,
            as_attachment=True,
            download_name=filename,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )
    
    except Exception as e:
        logger.error(f"Error downloading file: {e}")
        return jsonify({'error': 'Download failed'}), 500


@ppt_bp.route('/list', methods=['GET'])
def list_reports():
    """List all generated reports"""
    try:
        reports = []
        
        if OUTPUT_DIR.exists():
            for file_path in sorted(OUTPUT_DIR.glob('*.pptx'), reverse=True):
                stat = file_path.stat()
                reports.append({
                    'filename': file_path.name,
                    'size': stat.st_size,
                    'created': datetime.fromtimestamp(stat.st_ctime).isoformat(),
                    'modified': datetime.fromtimestamp(stat.st_mtime).isoformat()
                })
        
        return jsonify({
            'success': True,
            'count': len(reports),
            'reports': reports,
            'timestamp': datetime.now().isoformat()
        })
    
    except Exception as e:
        logger.error(f"Error listing reports: {e}")
        return jsonify({
            'success': False,
            'error': str(e)
        }), 500


@ppt_bp.route('/delete/<filename>', methods=['DELETE'])
def delete_report(filename):
    """Delete a generated report"""
    try:
        # Security: validate filename
        if '..' in filename or '/' in filename or '\\' in filename:
            return jsonify({'error': 'Invalid filename'}), 400
        
        file_path = OUTPUT_DIR / filename
        
        if not file_path.exists():
            return jsonify({'error': 'File not found'}), 404
        
        file_path.unlink()
        logger.info(f"Deleted report: {filename}")
        
        return jsonify({
            'success': True,
            'message': f'Report deleted: {filename}',
            'timestamp': datetime.now().isoformat()
        })
    
    except Exception as e:
        logger.error(f"Error deleting report: {e}")
        return jsonify({
            'success': False,
            'error': str(e)
        }), 500


@ppt_bp.route('/generate-from-screenshots', methods=['POST'])
def generate_ppt_from_screenshots():
    """
    Generate a PowerPoint report from screenshots
    
    Request body:
    {
        "screenshots": [
            {
                "page": 1,
                "imageData": "data:image/png;base64,..."
            },
            ...
        ],
        "selectedPhases": ["Test"],
        "title": "TDA Initiatives Insights - Test"
    }
    """
    try:
        data = request.get_json() or {}
        screenshots = data.get('screenshots', [])
        selected_phases = data.get('selectedPhases', [])
        title = data.get('title', 'V.E.T. Executive Report')
        
        if not screenshots:
            return jsonify({
                'success': False,
                'error': 'No screenshots provided'
            }), 400
        
        logger.info(f"PPT generation from screenshots requested. Pages: {len(screenshots)}, Phases: {selected_phases}")
        
        # Generate PowerPoint with screenshots
        prs = Presentation()
        prs.slide_width = Inches(9.6)
        prs.slide_height = Inches(7.2)
        
        slides_added = 0
        for screenshot_data in screenshots:
            page_num = screenshot_data.get('page', 1)
            image_base64 = screenshot_data.get('imageData', '')
            
            if not image_base64:
                logger.warning(f"Skipping page {page_num} - no image data")
                continue
            
            try:
                # Extract base64 data (remove data:image/png;base64, prefix if present)
                if ',' in image_base64:
                    image_base64 = image_base64.split(',')[1]
                
                # Decode base64 to bytes
                image_bytes = base64.b64decode(image_base64)
                image_stream = BytesIO(image_bytes)
                
                # Add blank slide
                blank_slide_layout = prs.slide_layouts[6]  # Blank layout
                slide = prs.slides.add_slide(blank_slide_layout)
                
                # Add image to slide with proper aspect ratio (don't stretch)
                # Image is positioned at top-left with margins, maintaining aspect ratio
                image_stream.seek(0)
                pil_image = Image.open(image_stream)
                img_width_px, img_height_px = pil_image.size
                
                # Calculate dimensions that maintain aspect ratio
                # Slide is 9.6" x 7.2", leave 0.2" margin on top/left, 0.4" on right
                max_width_inches = 9.2  # 9.6 - 0.4 for margins
                max_height_inches = 6.8  # 7.2 - 0.2 top - 0.2 bottom
                
                aspect_ratio = img_width_px / img_height_px
                
                # Determine which dimension is limiting WITHOUT specifying both
                # Specifying only one dimension lets python-pptx auto-calculate the other
                left = Inches(0.2)
                top = Inches(0.2)
                
                if aspect_ratio > (max_width_inches / max_height_inches):
                    # Image is wider proportionally - constrain by width only
                    image_stream.seek(0)
                    slide.shapes.add_picture(image_stream, left, top, width=Inches(max_width_inches))
                else:
                    # Image is taller proportionally - constrain by height only
                    image_stream.seek(0)
                    slide.shapes.add_picture(image_stream, left, top, height=Inches(max_height_inches))
                
                slides_added += 1
                logger.info(f"Added screenshot for page {page_num} (original: {img_width_px}x{img_height_px}px)")
            except Exception as e:
                logger.error(f"Error processing screenshot for page {page_num}: {e}")
                continue
        
        if slides_added == 0:
            return jsonify({
                'success': False,
                'error': 'No screenshots could be processed'
            }), 400
        
        # Save file to disk with WK format
        wk_format = get_walmart_week()
        filename = f'VET_Executive_Report_{wk_format}.pptx'
        file_path = OUTPUT_DIR / filename
        
        # Save to file
        with open(file_path, 'wb') as f:
            prs.save(f)
        
        logger.info(f"PPT generated successfully: {filename} with {slides_added} slides")
        
        # Send file as download
        return send_file(
            str(file_path),
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            as_attachment=True,
            download_name=filename
        )
        
    except Exception as e:
        logger.error(f"Error in /api/ppt/generate-from-screenshots: {e}", exc_info=True)
        return jsonify({
            'success': False,
            'error': f'Failed to generate PPT: {str(e)}'
        }), 500


@ppt_bp.route('/status', methods=['GET'])
def ppt_status():
    """Get status of PPT generation service"""
    return jsonify({
        'success': True,
        'service': 'V.E.T. Executive Report Generator',
        'status': 'operational',
        'output_directory': str(OUTPUT_DIR),
        'reports_directory_exists': OUTPUT_DIR.exists(),
        'last_report_path': report_manager.last_report_path,
        'last_report_time': report_manager.last_report_time.isoformat() if report_manager.last_report_time else None,
        'timestamp': datetime.now().isoformat()
    })


# Export
def register_ppt_routes(app):
    """Register PPT routes with Flask app"""
    app.register_blueprint(ppt_bp)
    logger.info("PPT routes registered")
