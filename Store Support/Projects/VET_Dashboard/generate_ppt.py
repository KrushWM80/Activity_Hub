"""
TDA Initiatives PowerPoint Report Generator
Creates professional PowerPoint presentations with phase-based slides
"""

import json
import os
from datetime import datetime
from typing import List, Dict, Any
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from google.cloud import bigquery
import logging

# Configuration
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Walmart Living Design Colors
WALMART_BLUE_DARK = RGBColor(30, 58, 138)
WALMART_BLUE = RGBColor(0, 113, 206)
WALMART_BLUE_LIGHT = RGBColor(219, 234, 254)
WALMART_YELLOW = RGBColor(255, 204, 0)
WALMART_YELLOW_DARK = RGBColor(255, 165, 0)
TEXT_PRIMARY = RGBColor(33, 33, 33)
TEXT_SECONDARY = RGBColor(102, 102, 102)
SUCCESS = RGBColor(16, 124, 16)
WARNING = RGBColor(247, 99, 12)
ERROR = RGBColor(220, 53, 69)

# BigQuery Configuration
PROJECT_ID = "wmt-assetprotection-prod"
DATASET_ID = "Store_Support_Dev"
TABLE_ID = "Output- TDA Report"
FULL_TABLE_ID = f"{PROJECT_ID}.{DATASET_ID}.`{TABLE_ID}`"


class TDAPowerPointGenerator:
    """Generates PowerPoint presentations for TDA Initiatives"""
    
    def __init__(self, bigquery_client=None):
        self.client = bigquery_client or bigquery.Client(project=PROJECT_ID)
        self.presentation = None
        self.data = []
    
    def fetch_data(self) -> List[Dict[str, Any]]:
        """Fetch all TDA data from BigQuery"""
        try:
            query = f"""
            SELECT 
                Topic AS `Initiative - Project Title`,
                Health_Update AS `Health Status`,
                Phase,
                SUM(CASE 
                    WHEN Phase = Facility_Phase THEN Facility 
                    ELSE 0 
                END) AS `# of Stores`,
                Dallas_POC,
                TDA_Ownership,
                Intake_Card_Nbr AS `Project ID`,
                Intake_n_Testing AS `Intake & Testing`,
                Deployment
            FROM 
                {FULL_TABLE_ID}
            GROUP BY
                Topic,
                Health_Update,
                Phase,
                Dallas_POC,
                TDA_Ownership,
                Intake_Card_Nbr,
                Intake_n_Testing,
                Deployment
            ORDER BY 
                Phase ASC,
                Topic ASC
            """
            
            query_job = self.client.query(query)
            rows = query_job.result()
            self.data = [dict(row) for row in rows]
            
            logger.info(f"Fetched {len(self.data)} records from BigQuery")
            return self.data
            
        except Exception as e:
            logger.error(f"Error fetching data: {e}")
            return []
    
    def create_presentation(self) -> Presentation:
        """Create a new presentation"""
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        return prs
    
    def add_title_slide(self, prs: Presentation, title: str, subtitle: str):
        """Add a title slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = WALMART_BLUE_DARK
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5), Inches(9), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        title_para.font.size = Pt(54)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        title_para.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4.2), Inches(9), Inches(1)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.text = subtitle
        subtitle_para.font.size = Pt(28)
        subtitle_para.font.color.rgb = WALMART_YELLOW
        subtitle_para.alignment = PP_ALIGN.CENTER
        
        # Logo/Icon
        logo_box = slide.shapes.add_textbox(
            Inches(4.5), Inches(0.5), Inches(1), Inches(1)
        )
        logo_para = logo_box.text_frame.paragraphs[0]
        logo_para.text = "💼"
        logo_para.font.size = Pt(44)
        logo_para.alignment = PP_ALIGN.CENTER
    
    def add_phase_summary_slide(self, prs: Presentation, phase: str, phase_data: List[Dict], slide_index: int = 0):
        """Add a summary slide for a phase (splits into multiple slides if needed)"""
        # Calculate items per slide (dynamically based on available space)
        # Available space: 7.5" - 4.0" (table_top) = 3.5" for content
        # With header/footer: ~3.0" for items
        # Using 0.3" per row with 0.05" spacing = 0.35" per item
        # This gives us ~8-9 items per slide comfortably
        ITEMS_PER_SLIDE = 8
        
        # Determine which items to show on this slide
        start_idx = slide_index * ITEMS_PER_SLIDE
        end_idx = min(start_idx + ITEMS_PER_SLIDE, len(phase_data))
        initiatives_to_show = phase_data[start_idx:end_idx]
        
        if not initiatives_to_show:
            return
        
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)
        
        # Add border at top
        top_border = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0), Inches(0), Inches(10), Inches(0.15)
        )
        top_border.fill.solid()
        top_border.fill.fore_color.rgb = WALMART_BLUE
        top_border.line.fill.background()
        
        # Phase title (with slide number if multi-slide phase)
        total_slides = (len(phase_data) + ITEMS_PER_SLIDE - 1) // ITEMS_PER_SLIDE
        title_text = f"Phase: {phase}"
        if total_slides > 1:
            title_text += f" ({slide_index + 1}/{total_slides})"
        
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(9), Inches(0.6)
        )
        title_para = title_box.text_frame.paragraphs[0]
        title_para.text = title_text
        title_para.font.size = Pt(40)
        title_para.font.bold = True
        title_para.font.color.rgb = WALMART_BLUE_DARK
        
        # Only show summary stats on first slide of phase
        if slide_index == 0:
            # Summary stats
            total_projects = len(phase_data)
            total_stores = sum([int(row.get('# of Stores', 0) or 0) for row in phase_data])
            
            on_track = sum(1 for row in phase_data if 'on track' in str(row.get('Health Status', '')).lower())
            at_risk = sum(1 for row in phase_data if 'at risk' in str(row.get('Health Status', '')).lower())
            off_track = sum(1 for row in phase_data if 'off track' in str(row.get('Health Status', '')).lower())
            
            # Stats boxes
            stat_y = 1.2
            stats = [
                (f"Total Initiatives", total_projects, WALMART_BLUE),
                (f"On Track", on_track, SUCCESS),
                (f"At Risk", at_risk, WARNING),
                (f"Off Track", off_track, ERROR),
            ]
            
            for idx, (label, value, color) in enumerate(stats):
                x_pos = 0.5 + (idx % 5) * 1.9
                if idx >= 5:
                    stat_y = 2.8
                    x_pos = 0.5 + (idx - 5) * 1.9
                
                # Box
                stat_shape = slide.shapes.add_shape(
                    1,  # Rectangle
                    Inches(x_pos), Inches(stat_y), Inches(1.7), Inches(0.8)
                )
                stat_shape.fill.solid()
                stat_shape.fill.fore_color.rgb = RGBColor(245, 245, 245)
                stat_shape.line.color.rgb = color
                stat_shape.line.width = Pt(2)
                
                # Value
                val_box = slide.shapes.add_textbox(
                    Inches(x_pos + 0.1), Inches(stat_y + 0.1), Inches(1.5), Inches(0.35)
                )
                val_para = val_box.text_frame.paragraphs[0]
                val_para.text = str(value)
                val_para.font.size = Pt(28)
                val_para.font.bold = True
                val_para.font.color.rgb = color
                val_para.alignment = PP_ALIGN.CENTER
                
                # Label
                lbl_box = slide.shapes.add_textbox(
                    Inches(x_pos + 0.1), Inches(stat_y + 0.45), Inches(1.5), Inches(0.3)
                )
                lbl_para = lbl_box.text_frame.paragraphs[0]
                lbl_para.text = label
                lbl_para.font.size = Pt(9)
                lbl_para.font.color.rgb = TEXT_SECONDARY
                lbl_para.alignment = PP_ALIGN.CENTER
            
            table_top = 4.0
        else:
            # For continuation slides, start higher
            table_top = 1.2
        
        # Table header
        slide.shapes.add_textbox(
            Inches(0.5), Inches(table_top - 0.4), Inches(9), Inches(0.3)
        ).text_frame.paragraphs[0].text = "Key Initiatives"
        
        # Render initiatives with fixed spacing (no stretching)
        ROW_HEIGHT = 0.35  # Fixed height per row
        for idx, initiative in enumerate(initiatives_to_show):
            y_pos = table_top + (idx * ROW_HEIGHT)
            
            # Initiative name
            name_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(y_pos), Inches(5), Inches(0.3)
            )
            name_para = name_box.text_frame.paragraphs[0]
            name_para.text = str(initiative.get('Initiative - Project Title', 'Unknown'))[:50]
            name_para.font.size = Pt(10)
            name_para.font.color.rgb = TEXT_PRIMARY
            
            # Health status
            status = str(initiative.get('Health Status', 'Unknown'))
            status_color = SUCCESS
            if 'at risk' in status.lower():
                status_color = WARNING
            elif 'off track' in status.lower():
                status_color = ERROR
            
            status_box = slide.shapes.add_textbox(
                Inches(5.7), Inches(y_pos), Inches(1.5), Inches(0.3)
            )
            status_para = status_box.text_frame.paragraphs[0]
            status_para.text = status
            status_para.font.size = Pt(9)
            status_para.font.color.rgb = status_color
            status_para.font.bold = True
            
            # # of stores
            stores_box = slide.shapes.add_textbox(
                Inches(7.4), Inches(y_pos), Inches(1), Inches(0.3)
            )
            stores_para = stores_box.text_frame.paragraphs[0]
            stores_para.text = str(initiative.get('# of Stores', '0'))
            stores_para.font.size = Pt(9)
            stores_para.font.color.rgb = WALMART_BLUE
            stores_para.font.bold = True
            
            # POC
            poc_box = slide.shapes.add_textbox(
                Inches(8.6), Inches(y_pos), Inches(1.2), Inches(0.3)
            )
            poc_para = poc_box.text_frame.paragraphs[0]
            poc_para.text = str(initiative.get('Dallas POC', ''))[:15]
            poc_para.font.size = Pt(8)
            poc_para.font.color.rgb = TEXT_SECONDARY
        
        # Footer
        footer_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(7), Inches(9), Inches(0.4)
        )
        footer_para = footer_box.text_frame.paragraphs[0]
        footer_para.text = f"Dallas Team Report | Data Source: wmt-assetprotection-prod.Store_Support_Dev | Generated: {datetime.now().strftime('%B %d, %Y')}"
        footer_para.font.size = Pt(8)
        footer_para.font.color.rgb = TEXT_SECONDARY
        footer_para.alignment = PP_ALIGN.CENTER
    
    def generate_report(self, output_path: str = None) -> str:
        """Generate complete PowerPoint report"""
        
        if not self.data:
            logger.error("No data available. Run fetch_data() first.")
            return None
        
        # Create presentation
        self.presentation = self.create_presentation()
        
        # Get WM Week from data (if available)
        wm_week = "Current"
        if self.data and 'WM Week' in self.data[0]:
            wm_weeks = set(str(row.get('WM Week', 'Current')) for row in self.data if row.get('WM Week'))
            if wm_weeks:
                wm_week = ', '.join(sorted(wm_weeks))
        
        # Add title slide
        self.add_title_slide(
            self.presentation,
            "Dallas Team Report",
            f"Walmart Enterprise Transformation - {wm_week}"
        )
        
        # Get unique phases
        phases = sorted(set(str(row.get('Phase', 'Unknown')) for row in self.data))
        
        # Add slides for each phase (split into multiple slides if needed)
        for phase in phases:
            phase_data = [row for row in self.data if str(row.get('Phase', '')) == phase]
            
            # Calculate number of slides needed for this phase (8 items per slide)
            ITEMS_PER_SLIDE = 8
            num_slides = (len(phase_data) + ITEMS_PER_SLIDE - 1) // ITEMS_PER_SLIDE
            
            # Add slide for each chunk of the phase
            for slide_idx in range(num_slides):
                self.add_phase_summary_slide(self.presentation, phase, phase_data, slide_idx)
        
        # Add overview slide
        self.add_overview_slide()
        
        # Save presentation
        if not output_path:
            output_path = f"Dallas_Team_Report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        
        self.presentation.save(output_path)
        logger.info(f"PowerPoint report generated: {output_path}")
        
        return output_path
    
    def add_overview_slide(self):
        """Add an overview slide with all projects"""
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[6])
        
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(9), Inches(0.6)
        )
        title_para = title_box.text_frame.paragraphs[0]
        title_para.text = "Complete Initiative Summary"
        title_para.font.size = Pt(40)
        title_para.font.bold = True
        title_para.font.color.rgb = WALMART_BLUE_DARK
        
        # Overall stats
        total_projects = len(self.data)
        total_stores = sum([int(row.get('# of Stores', 0) or 0) for row in self.data])
        on_track = sum(1 for row in self.data if 'on track' in str(row.get('Health Status', '')).lower())
        at_risk = sum(1 for row in self.data if 'at risk' in str(row.get('Health Status', '')).lower())
        off_track = sum(1 for row in self.data if 'off track' in str(row.get('Health Status', '')).lower())
        
        stats_text = f"""
Total Initiatives: {total_projects}
On Track: {on_track}
At Risk: {at_risk}
Off Track: {off_track}

Key Insights:
• {(on_track/total_projects*100):.1f}% of initiatives are on track
• {(at_risk/total_projects*100):.1f}% of initiatives are at risk
• {(off_track/total_projects*100):.1f}% of initiatives are off track
        """
        
        stats_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.2), Inches(9), Inches(5.8)
        )
        stats_frame = stats_box.text_frame
        stats_frame.word_wrap = True
        stats_para = stats_frame.paragraphs[0]
        stats_para.text = stats_text
        stats_para.font.size = Pt(14)
        stats_para.font.color.rgb = TEXT_PRIMARY
        stats_para.line_spacing = 1.8
        
        # Footer
        footer_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(7), Inches(9), Inches(0.4)
        )
        footer_para = footer_box.text_frame.paragraphs[0]
        footer_para.text = f"Report Generated: {datetime.now().strftime('%B %d, %Y at %I:%M %p')}"
        footer_para.font.size = Pt(8)
        footer_para.font.color.rgb = TEXT_SECONDARY
        footer_para.alignment = PP_ALIGN.CENTER


def main():
    """Main execution"""
    try:
        # Initialize generator
        generator = TDAPowerPointGenerator()
        
        # Fetch data
        logger.info("Fetching TDA data from BigQuery...")
        generator.fetch_data()
        
        if not generator.data:
            logger.error("No data retrieved. Check BigQuery connection and table.")
            return
        
        # Generate report
        logger.info("Generating PowerPoint report...")
        output_file = generator.generate_report()
        
        if output_file:
            logger.info(f"✅ Report generated successfully: {output_file}")
        else:
            logger.error("Failed to generate report")
    
    except Exception as e:
        logger.error(f"Error: {e}")
        import traceback
        traceback.print_exc()


if __name__ == "__main__":
    main()
