"""
V.E.T. Dashboard PPT Generator
Creates PowerPoint reports directly from dashboard API data
"""

import logging
from datetime import datetime, timedelta
from pathlib import Path
from io import BytesIO
import json
import urllib.request
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from xml.sax.saxutils import escape

logger = logging.getLogger(__name__)

# Configuration
DASHBOARD_API = "http://127.0.0.1:5001"
OUTPUT_DIR = Path(__file__).parent / 'reports'
OUTPUT_DIR.mkdir(exist_ok=True)

COLORS = {
    'walmart_blue_dark': RGBColor(0x1E, 0x3A, 0x8A),
    'walmart_blue': RGBColor(0x00, 0x71, 0xCE),
    'walmart_yellow': RGBColor(0xFF, 0xC2, 0x20),
    'on_track': RGBColor(0x10, 0x7C, 0x10),
    'at_risk': RGBColor(0xF7, 0x63, 0x0C),
    'off_track': RGBColor(0xDC, 0x35, 0x45),
}


def get_walmart_week():
    """Calculate current Walmart Week"""
    today = datetime.now().date()
    if today.month >= 2:
        year_start = datetime(today.year, 2, 1).date()
    else:
        year_start = datetime(today.year - 1, 2, 1).date()
    
    days_until_saturday = (5 - year_start.weekday()) % 7
    first_saturday = year_start + timedelta(days=days_until_saturday)
    days_diff = (today - first_saturday).days
    weeks = max(1, (days_diff // 7) + 1)
    
    return f"WK{weeks:02d}"


def fetch_dashboard_data():
    """Fetch summary and all projects from dashboard API"""
    try:
        # Get summary
        summary_url = f"{DASHBOARD_API}/api/summary"
        with urllib.request.urlopen(summary_url, timeout=5) as response:
            summary_data = json.loads(response.read().decode())
        summary = summary_data.get('summary', {})
        
        # Get all projects
        data_url = f"{DASHBOARD_API}/api/data?phase=All"
        with urllib.request.urlopen(data_url, timeout=5) as response:
            data_response = json.loads(response.read().decode())
        projects = data_response.get('data', [])
        
        return summary, projects
    except Exception as e:
        logger.error(f"Error fetching dashboard data: {e}")
        return None, None


def generate_vet_pptx():
    """
    Generate VET Executive Report PowerPoint from dashboard data
    Returns: Path to generated PPTX file
    """
    try:
        logger.info("Fetching dashboard data for PPT generation...")
        summary, projects = fetch_dashboard_data()
        
        if not summary or not projects:
            logger.error("Failed to fetch dashboard data")
            return None
        
        logger.info(f"Generating PPT with {len(projects)} projects")
        
        # Create presentation
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        blank_layout = prs.slide_layouts[6]
        
        # ── Slide 1: Title Slide ──
        slide = prs.slides.add_slide(blank_layout)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLORS['walmart_blue_dark']
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = "V.E.T. Executive Report"
        run.font.size = Pt(44)
        run.font.bold = True
        run.font.color.rgb = RGBColor(255, 255, 255)
        
        # Subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(0.8))
        p = subtitle_box.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = "Walmart Enterprise Transformation"
        run.font.size = Pt(20)
        run.font.color.rgb = RGBColor(255, 255, 255)
        
        # Yellow accent
        accent = slide.shapes.add_shape(1, 0, Inches(7.0), prs.slide_width, Inches(0.5))
        accent.fill.solid()
        accent.fill.fore_color.rgb = COLORS['walmart_yellow']
        accent.line.fill.background()
        
        # ── Slide 2: Executive Summary with Stats ──
        slide = prs.slides.add_slide(blank_layout)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(245, 245, 245)
        
        # Header
        header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
        p = header_box.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = "Dashboard Summary"
        run.font.size = Pt(28)
        run.font.bold = True
        run.font.color.rgb = COLORS['walmart_blue_dark']
        
        # Stats in grid layout
        stats = {
            'Total Projects': summary.get('total_projects', 0),
            'Total Stores': summary.get('total_stores', 0),
            'On Track': summary.get('by_health_status', {}).get('On Track', 0),
            'At Risk': summary.get('by_health_status', {}).get('At Risk', 0),
            'Off Track': summary.get('by_health_status', {}).get('Off Track', 0),
        }
        
        stat_names = list(stats.keys())
        for idx, (name, value) in enumerate(stats.items()):
            col = idx % 5
            row = 2
            
            x = Inches(0.5 + col * 1.9)
            y = Inches(row)
            
            # Box
            box = slide.shapes.add_shape(1, x, y, Inches(1.7), Inches(1.2))
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(255, 255, 255)
            box.line.color.rgb = COLORS['walmart_blue']
            box.line.width = Pt(2)
            
            # Label
            label_box = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.1), Inches(1.5), Inches(0.4))
            p = label_box.text_frame.paragraphs[0]
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = COLORS['walmart_blue_dark']
            p.text = name
            
            # Value
            val_box = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.5), Inches(1.5), Inches(0.6))
            p = val_box.text_frame.paragraphs[0]
            p.font.size = Pt(20)
            p.font.bold = True
            run = p.add_run()
            run.text = f"{value:,}" if isinstance(value, int) else str(value)
            
            if name == 'On Track':
                run.font.color.rgb = COLORS['on_track']
            elif name == 'At Risk':
                run.font.color.rgb = COLORS['at_risk']
            elif name == 'Off Track':
                run.font.color.rgb = COLORS['off_track']
            else:
                run.font.color.rgb = COLORS['walmart_blue']
        
        # ── Slide 3+: Phase Tables ──
        # Group by phase
        phases_dict = {}
        for proj in projects:
            phase = proj.get('Phase', 'Unknown')
            if phase not in phases_dict:
                phases_dict[phase] = []
            phases_dict[phase].append(proj)
        
        phase_order = ['Pending', 'POC/POT', 'Test', 'Mkt Scale', 'Roll/Deploy']
        
        for phase in phase_order:
            if phase not in phases_dict or not phases_dict[phase]:
                continue
            
            phase_projects = phases_dict[phase]
            
            slide = prs.slides.add_slide(blank_layout)
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(255, 255, 255)
            
            # Phase header
            header_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9.4), Inches(0.5))
            p = header_box.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = f"{phase} - {len(phase_projects)} Projects"
            run.font.size = Pt(20)
            run.font.bold = True
            run.font.color.rgb = RGBColor(255, 255, 255)
            
            header_shape = slide.shapes.add_shape(1, Inches(0.3), Inches(0.2), Inches(9.4), Inches(0.5))
            header_shape.fill.solid()
            header_shape.fill.fore_color.rgb = COLORS['walmart_blue_dark']
            header_shape.line.fill.background()
            slide.shapes._spTree.remove(header_shape._element)
            slide.shapes._spTree.insert(2, header_shape._element)
            
            # Simple table with key columns
            from pptx.util import Pt as PtSize
            from pptx.enum.text import MSO_ANCHOR
            
            # Add table rows
            row_height = Inches(0.35)
            rows_per_slide = 17
            
            col_widths = [Inches(3.5), Inches(1.0), Inches(1.0), Inches(1.5), Inches(1.0), Inches(1.4)]
            columns = ['Initiative', 'Health', 'Phase', 'WM Week', 'Stores', 'Notes']
            
            y_pos = Inches(0.85)
            
            # Header row
            x_pos = Inches(0.3)
            for col_idx, col_name in enumerate(columns):
                col_box = slide.shapes.add_textbox(x_pos, y_pos, col_widths[col_idx], row_height)
                col_box.text_frame.margin_bottom = Inches(0.05)
                col_box.text_frame.margin_top = Inches(0.05)
                col_box.text_frame.word_wrap = True
                p = col_box.text_frame.paragraphs[0]
                p.font.size = Pt(9)
                p.font.bold = True
                p.font.color.rgb = RGBColor(255, 255, 255)
                p.text = col_name
                
                col_shape = slide.shapes.add_shape(1, x_pos, y_pos, col_widths[col_idx], row_height)
                col_shape.fill.solid()
                col_shape.fill.fore_color.rgb = COLORS['walmart_blue']
                col_shape.line.width = Pt(0.5)
                slide.shapes._spTree.remove(col_shape._element)
                slide.shapes._spTree.insert(2, col_shape._element)
                
                x_pos += col_widths[col_idx]
            
            y_pos += row_height
            
            # Data rows
            for row_idx, proj in enumerate(phase_projects[:rows_per_slide]):
                x_pos = Inches(0.3)
                bg_color = RGBColor(245, 245, 245) if row_idx % 2 == 0 else RGBColor(255, 255, 255)
                
                # Initiative
                cell = slide.shapes.add_textbox(x_pos, y_pos, col_widths[0], row_height)
                cell.text_frame.margin_bottom = Inches(0.02)
                cell.text_frame.margin_top = Inches(0.02)
                cell.text_frame.word_wrap = True
                p = cell.text_frame.paragraphs[0]
                p.font.size = Pt(8)
                p.text = proj.get('Initiative - Project Title', '')
                
                cell_shape = slide.shapes.add_shape(1, x_pos, y_pos, col_widths[0], row_height)
                cell_shape.fill.solid()
                cell_shape.fill.fore_color.rgb = bg_color
                cell_shape.line.width = Pt(0.5)
                cell_shape.line.color.rgb = RGBColor(200, 200, 200)
                slide.shapes._spTree.remove(cell_shape._element)
                slide.shapes._spTree.insert(2, cell_shape._element)
                x_pos += col_widths[0]
                
                # Health Status
                health = proj.get('Health Status', '')
                cell = slide.shapes.add_textbox(x_pos, y_pos, col_widths[1], row_height)
                cell.text_frame.margin_bottom = Inches(0.02)
                cell.text_frame.margin_top = Inches(0.02)
                p = cell.text_frame.paragraphs[0]
                p.font.size = Pt(8)
                p.font.bold = True
                p.text = health
                if 'On Track' in health:
                    p.font.color.rgb = COLORS['on_track']
                elif 'At Risk' in health:
                    p.font.color.rgb = COLORS['at_risk']
                elif 'Off Track' in health:
                    p.font.color.rgb = COLORS['off_track']
                
                cell_shape = slide.shapes.add_shape(1, x_pos, y_pos, col_widths[1], row_height)
                cell_shape.fill.solid()
                cell_shape.fill.fore_color.rgb = bg_color
                cell_shape.line.width = Pt(0.5)
                cell_shape.line.color.rgb = RGBColor(200, 200, 200)
                slide.shapes._spTree.remove(cell_shape._element)
                slide.shapes._spTree.insert(2, cell_shape._element)
                x_pos += col_widths[1]
                
                # Phase
                cell = slide.shapes.add_textbox(x_pos, y_pos, col_widths[2], row_height)
                cell.text_frame.margin_bottom = Inches(0.02)
                cell.text_frame.margin_top = Inches(0.02)
                p = cell.text_frame.paragraphs[0]
                p.font.size = Pt(8)
                p.text = proj.get('Phase', '')
                
                cell_shape = slide.shapes.add_shape(1, x_pos, y_pos, col_widths[2], row_height)
                cell_shape.fill.solid()
                cell_shape.fill.fore_color.rgb = bg_color
                cell_shape.line.width = Pt(0.5)
                cell_shape.line.color.rgb = RGBColor(200, 200, 200)
                slide.shapes._spTree.remove(cell_shape._element)
                slide.shapes._spTree.insert(2, cell_shape._element)
                x_pos += col_widths[2]
                
                # WM Week
                cell = slide.shapes.add_textbox(x_pos, y_pos, col_widths[3], row_height)
                cell.text_frame.margin_bottom = Inches(0.02)
                cell.text_frame.margin_top = Inches(0.02)
                p = cell.text_frame.paragraphs[0]
                p.font.size = Pt(8)
                p.text = str(proj.get('WM Week', ''))
                
                cell_shape = slide.shapes.add_shape(1, x_pos, y_pos, col_widths[3], row_height)
                cell_shape.fill.solid()
                cell_shape.fill.fore_color.rgb = bg_color
                cell_shape.line.width = Pt(0.5)
                cell_shape.line.color.rgb = RGBColor(200, 200, 200)
                slide.shapes._spTree.remove(cell_shape._element)
                slide.shapes._spTree.insert(2, cell_shape._element)
                x_pos += col_widths[3]
                
                # Stores
                cell = slide.shapes.add_textbox(x_pos, y_pos, col_widths[4], row_height)
                cell.text_frame.margin_bottom = Inches(0.02)
                cell.text_frame.margin_top = Inches(0.02)
                p = cell.text_frame.paragraphs[0]
                p.font.size = Pt(8)
                stores = proj.get('# of Stores', 0)
                p.text = f"{stores:,}" if isinstance(stores, int) else str(stores)
                p.font.bold = True
                
                cell_shape = slide.shapes.add_shape(1, x_pos, y_pos, col_widths[4], row_height)
                cell_shape.fill.solid()
                cell_shape.fill.fore_color.rgb = bg_color
                cell_shape.line.width = Pt(0.5)
                cell_shape.line.color.rgb = RGBColor(200, 200, 200)
                slide.shapes._spTree.remove(cell_shape._element)
                slide.shapes._spTree.insert(2, cell_shape._element)
                x_pos += col_widths[4]
                
                # Notes
                cell = slide.shapes.add_textbox(x_pos, y_pos, col_widths[5], row_height)
                cell.text_frame.margin_bottom = Inches(0.02)
                cell.text_frame.margin_top = Inches(0.02)
                cell.text_frame.word_wrap = True
                p = cell.text_frame.paragraphs[0]
                p.font.size = Pt(7)
                notes = proj.get('Executive Notes', '')
                p.text = notes[:50] if len(notes) > 50 else notes
                
                cell_shape = slide.shapes.add_shape(1, x_pos, y_pos, col_widths[5], row_height)
                cell_shape.fill.solid()
                cell_shape.fill.fore_color.rgb = bg_color
                cell_shape.line.width = Pt(0.5)
                cell_shape.line.color.rgb = RGBColor(200, 200, 200)
                slide.shapes._spTree.remove(cell_shape._element)
                slide.shapes._spTree.insert(2, cell_shape._element)
                
                y_pos += row_height
        
        # Save
        wk = get_walmart_week()
        filename = f"VET_Executive_Report_{wk}.pptx"
        output_path = OUTPUT_DIR / filename
        
        prs.save(str(output_path))
        
        logger.info(f"PPT generated successfully: {filename}")
        return str(output_path)
    
    except Exception as e:
        logger.error(f"Error generating PPT: {e}", exc_info=True)
        return None
