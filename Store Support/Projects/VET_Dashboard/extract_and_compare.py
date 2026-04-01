"""Extract content from Dallas VET Weekly Report PDF using pptx reference + BQ comparison"""
from pptx import Presentation
from PIL import Image
import io
import os

# The PDF was generated from the dashboard PPT - let's extract the images and analyze them
# First check if we can read the reference PPTX which has the same content
ref_pptx = 'reports/VET_Executive_Report (3).pptx'

prs = Presentation(ref_pptx)
print(f"Extracting images from reference PPT: {ref_pptx}")
print(f"Slides: {len(prs.slides)}")
print()

# Extract each slide image to disk for viewing
output_dir = 'reports/extracted_slides'
os.makedirs(output_dir, exist_ok=True)

for i, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if shape.shape_type == 13:  # Picture
            img_data = shape.image.blob
            img = Image.open(io.BytesIO(img_data))
            output_path = os.path.join(output_dir, f'reference_slide_{i+1}.png')
            img.save(output_path)
            print(f"Slide {i+1}: Saved {img.size[0]}x{img.size[1]}px -> {output_path}")

print()
print("Images extracted! Now querying BigQuery for comparison data...")
print()

# Now query BQ to compare
from google.cloud import bigquery
client = bigquery.Client(project='wmt-assetprotection-prod')

# Get all projects with their details
q = """
SELECT DISTINCT 
    Topic as project_title,
    TDA_Ownership,
    Phase,
    Health_Update as health_status,
    Dallas_POC as dallas_notes
FROM `wmt-assetprotection-prod.Store_Support_Dev.Output- TDA Report`
WHERE Phase != 'Complete'
ORDER BY TDA_Ownership, Phase, Topic
"""

rows = list(client.query(q).result())

# Group by ownership
from collections import defaultdict
by_owner = defaultdict(list)
for r in rows:
    owner = r.TDA_Ownership or 'No Selection Provided'
    by_owner[owner].append(r)

print("=" * 100)
print("ALL 49 ACTIVE PROJECTS BY TDA OWNERSHIP")
print("=" * 100)

for owner in ['Dallas POC', 'Dallas VET', 'Intake & Test', 'Deployment', 'No Selection Provided', '*Select Owner']:
    if owner not in by_owner:
        continue
    projects = by_owner[owner]
    print(f"\n--- TDA Ownership: \"{owner}\" ({len(projects)} projects) ---")
    for p in projects:
        notes = (p.dallas_notes or '')[:60]
        print(f"  [{p.Phase:<12}] [{p.health_status:<10}] {p.project_title}")
        if notes and notes != 'No Note Provided':
            print(f"                                      Notes: {notes}")

# Summary
print()
print("=" * 100)
print("OWNERSHIP SUMMARY")
print("=" * 100)
for owner in ['Dallas POC', 'Dallas VET', 'Intake & Test', 'Deployment', 'No Selection Provided', '*Select Owner']:
    if owner in by_owner:
        count = len(by_owner[owner])
        phases = defaultdict(int)
        for p in by_owner[owner]:
            phases[p.Phase] += 1
        phase_str = ", ".join(f"{ph}: {c}" for ph, c in sorted(phases.items()))
        print(f"  {owner:<25} {count:>3} projects  ({phase_str})")

print()
print("=" * 100)
print("WHAT EACH DASHBOARD SHOWS")
print("=" * 100)

# TDA Dashboard shows ALL projects (all ownership, excludes Complete)
tda_count = len(rows)
print(f"\n  TDA Dashboard: {tda_count} projects (ALL ownership, excludes Complete phase)")

# VET Dashboard shows only Dallas POC
vet_dallas_poc = len(by_owner.get('Dallas POC', []))
vet_dallas_vet = len(by_owner.get('Dallas VET', []))
vet_count = vet_dallas_poc + vet_dallas_vet
print(f"  VET Dashboard: {vet_count} projects (Dallas POC: {vet_dallas_poc} + Dallas VET: {vet_dallas_vet})")

# Reference PDF (Dallas VET Weekly Report) - appears to show ALL 49
print(f"  Reference PDF: ~49 projects (appears to show ALL projects, not filtered)")

print()
print("=" * 100)
print("PROJECTS LABELED 'Dallas POC' IN BQ vs REFERENCE PDF")
print("=" * 100)
print(f"\nBigQuery has {vet_dallas_poc} projects with TDA_Ownership = 'Dallas POC':")
for p in by_owner.get('Dallas POC', []):
    print(f"  {p.project_title}")

print(f"\nBigQuery has {vet_dallas_vet} project(s) with TDA_Ownership = 'Dallas VET':")
for p in by_owner.get('Dallas VET', []):
    print(f"  {p.project_title}")

print(f"\nThe reference PDF 'Dallas VET Weekly Report 3.30.pdf' (7,904 KB)")
print(f"appears to include ALL {tda_count} projects regardless of TDA_Ownership.")
print(f"This suggests the person treats ALL projects as Dallas VET responsibility,")
print(f"not just the ones tagged 'Dallas POC' in the TDA_Ownership column.")
