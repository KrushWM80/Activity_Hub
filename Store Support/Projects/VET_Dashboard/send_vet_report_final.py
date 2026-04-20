#!/usr/bin/env python3
"""
V.E.T. Executive Report - Final Version (Complete Dashboard Integration)
Captures V.E.T. Dashboard content exactly as displayed including:
- Needs Attention section with at-risk initiatives
- Complete data tables with all columns (Initiative, Health, Phase, WM Week, # Stores, Executive Notes)
- Dashboard-rendered screenshots for PPT and PDF
"""

import sys
import os
import json
import subprocess
import tempfile
import io
import base64
import re
from pathlib import Path
from datetime import datetime, timedelta
from xml.sax.saxutils import escape
import urllib.request
import urllib.error
from PIL import Image, ImageChops

sys.path.insert(0, str(Path(__file__).parent))
from email_service import VETEmailService
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Configuration
SCRIPT_DIR = Path(__file__).parent
DASHBOARD_URL = "http://127.0.0.1:5001"
EDGE_PATH = r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe"
OUTPUT_DIR = SCRIPT_DIR / 'reports'
OUTPUT_DIR.mkdir(exist_ok=True)

COLORS = {
    'walmart_blue_dark': '#1E3A8A',
    'walmart_blue': '#0071CE',
    'walmart_yellow': '#FFCC00',
    'on_track': '#107C10',
    'at_risk': '#F7630C',
    'off_track': '#DC3545',
}


def get_current_walmart_week():
    """Get current Walmart Week from BQ Cal_Dim_Data (source of truth), with local fallback"""
    # Try BQ first
    try:
        import os as _os
        _creds = _os.environ.get('GOOGLE_APPLICATION_CREDENTIALS', '')
        if not _creds or not _os.path.exists(_creds):
            _os.environ['GOOGLE_APPLICATION_CREDENTIALS'] = _os.path.join(
                _os.environ.get('APPDATA', ''), 'gcloud', 'application_default_credentials.json')
        from google.cloud import bigquery as _bq
        _client = _bq.Client(project='wmt-assetprotection-prod')
        _q = """SELECT WM_WEEK_NBR FROM `wmt-assetprotection-prod.Store_Support_Dev.Cal_Dim_Data`
                WHERE CALENDAR_DATE = CURRENT_DATE() LIMIT 1"""
        for r in _client.query(_q).result():
            wk = int(r.WM_WEEK_NBR)
            print(f"     [OK] WM Week from BQ Cal_Dim: WK{wk:02d}")
            return f"WK{wk:02d}"
    except Exception as e:
        print(f"     [!] BQ Cal_Dim lookup failed: {e}")
    
    # Fallback: Walmart fiscal year starts Saturday nearest Feb 1
    today = datetime.now().date()
    if today.month >= 2:
        year_start = datetime(today.year, 2, 1).date()
    else:
        year_start = datetime(today.year - 1, 2, 1).date()
    
    days_until_saturday = (5 - year_start.weekday()) % 7
    first_saturday = year_start + timedelta(days=days_until_saturday)
    days_diff = (today - first_saturday).days
    weeks = max(1, (days_diff // 7) + 1)
    print(f"     [!] Using local fallback: WK{weeks:02d}")
    return f"WK{weeks:02d}"


def _fetch_from_bigquery_direct() -> list:
    """Fallback: query BigQuery directly when API server is unavailable"""
    _PHASE_MAP = {'POC/POT': 'Vet', 'Mkt Scale': 'Test Markets'}
    try:
        from google.cloud import bigquery as bq
        creds_path = os.environ.get('GOOGLE_APPLICATION_CREDENTIALS',
            os.path.join(os.environ.get('APPDATA', ''), 'gcloud', 'application_default_credentials.json'))
        if creds_path and os.path.exists(creds_path):
            os.environ['GOOGLE_APPLICATION_CREDENTIALS'] = creds_path
        bq_client = bq.Client(project='wmt-assetprotection-prod')
        query = """
        SELECT
            tda.Topic AS `Initiative - Project Title`,
            tda.Health_Update AS `Health Status`,
            tda.Phase,
            SUM(CASE WHEN tda.Phase = tda.Facility_Phase THEN tda.Facility ELSE 0 END) AS `# of Stores`,
            tda.Dallas_POC AS `Executive Notes`,
            tda.TDA_Ownership,
            tda.Intake_Card_Nbr AS `Project ID`,
            tda.Intake_n_Testing AS `Intake & Testing`,
            tda.Deployment,
            COALESCE(CAST(MIN(intake.WM_Week) AS STRING), 'TBD') AS `WM Week`
        FROM `wmt-assetprotection-prod.Store_Support_Dev.Output- TDA Report` tda
        LEFT JOIN `wmt-assetprotection-prod.Store_Support_Dev.IH_Intake_Data` intake
            ON CAST(tda.Intake_Card_Nbr AS STRING) = CAST(intake.Intake_Card AS STRING)
        WHERE tda.TDA_Ownership IN ('Dallas POC', 'Dallas VET')
        GROUP BY tda.Topic, tda.Health_Update, tda.Phase, tda.Dallas_POC,
                 tda.TDA_Ownership, tda.Intake_Card_Nbr, tda.Intake_n_Testing, tda.Deployment
        ORDER BY tda.Phase ASC, tda.Topic ASC
        """
        rows = list(bq_client.query(query).result())
        data = [dict(r) for r in rows]
        for row in data:
            row['Phase'] = _PHASE_MAP.get(row.get('Phase', ''), row.get('Phase', 'Unknown'))
        print(f"     [OK] Direct BigQuery returned {len(data)} projects")
        return data
    except Exception as bq_err:
        print(f"     [!] Direct BigQuery also failed: {bq_err}")
        return None


def _build_stats(projects: list) -> dict:
    """Build summary stats dict from a project list"""
    current_wm_week = get_current_walmart_week()
    return {
        'total_projects': len(projects),
        'total_stores': sum(int(p.get('# of Stores', 0) or 0) for p in projects),
        'on_track': len([p for p in projects if 'On Track' in str(p.get('Health Status', ''))]),
        'at_risk': len([p for p in projects if 'At Risk' in str(p.get('Health Status', ''))]),
        'off_track': len([p for p in projects if 'Off Track' in str(p.get('Health Status', ''))]),
        'wm_week': current_wm_week,
    }


def fetch_dashboard_data(api_url: str = DASHBOARD_URL) -> tuple:
    """Fetch dashboard data: API first, then direct BigQuery, then sample data"""
    # --- Attempt 1: Live API ---
    try:
        summary_url = f"{api_url}/api/summary"
        with urllib.request.urlopen(summary_url, timeout=5) as response:
            summary_data = json.loads(response.read().decode())
        summary = summary_data.get('summary', {})

        data_url = f"{api_url}/api/data?phase=All"
        with urllib.request.urlopen(data_url, timeout=5) as response:
            data_response = json.loads(response.read().decode())
        projects = data_response.get('data', [])

        if projects:
            current_wm_week = get_current_walmart_week()
            stats = {
                'total_projects': summary.get('total_projects', 0),
                'total_stores': summary.get('total_stores', 0),
                'on_track': summary.get('by_health_status', {}).get('On Track', 0),
                'at_risk': summary.get('by_health_status', {}).get('At Risk', 0),
                'off_track': summary.get('by_health_status', {}).get('Off Track', 0),
                'wm_week': current_wm_week,
            }
            print(f"     [OK] API returned {len(projects)} projects")
            return stats, projects
    except Exception as e:
        print(f"     [!] API fetch failed: {e}")

    # --- Attempt 2: Direct BigQuery ---
    print("     [*] Falling back to direct BigQuery query...")
    projects = _fetch_from_bigquery_direct()
    if projects:
        return _build_stats(projects), projects

    # --- Attempt 3: Sample data (last resort) ---
    print("     [*] Falling back to sample data...")
    try:
        from sample_data import SAMPLE_DATA_49_PROJECTS
        return _build_stats(SAMPLE_DATA_49_PROJECTS), SAMPLE_DATA_49_PROJECTS
    except Exception:
        return None, None


def build_needs_attention_html(projects: list) -> list:
    """Build HTML pages for Needs Attention section (at-risk initiatives).
    Returns a list of HTML strings — one per page (max 10 cards each).
    Returns empty list if no At Risk items.
    """
    
    at_risk_items = [p for p in projects if str(p.get('Health Status', '')).lower() == 'at risk']
    
    if not at_risk_items:
        return []
    
    CARDS_PER_PAGE = 10
    pages = []
    total_pages = -(-len(at_risk_items) // CARDS_PER_PAGE)  # ceil division
    
    for page_idx in range(total_pages):
        start = page_idx * CARDS_PER_PAGE
        end = start + CARDS_PER_PAGE
        page_items = at_risk_items[start:end]
        
        cards_html = ''
        for item in page_items:
            title = escape(item.get('Initiative - Project Title', 'Unknown'))
            stores = item.get('# of Stores', 0)
            notes = escape(item.get('Executive Notes', 'No notes provided'))
            phase = escape(item.get('Phase', 'Unknown'))
            wm_week = escape(str(item.get('WM Week', 'N/A')))
            status_label = escape(str(item.get('Health Status', 'At Risk')))
            deployment = escape(str(item.get('Deployment', 'No Note Provided')))
            
            cards_html += f'''<div style="background:white;border:1px solid #dc3545;border-radius:6px;padding:16px;box-shadow:0 1px 3px rgba(220,53,69,0.1);">
<div style="margin-bottom:8px;">
  <div style="display:inline-block;background:#dc3545;color:white;padding:4px 8px;border-radius:4px;font-size:12px;font-weight:600;margin-bottom:4px;">{status_label}</div>
  <div style="font-size:12px;color:#666;margin-top:4px;">Phase: {phase}</div>
</div>
<div style="font-weight:700;color:#212121;margin-bottom:8px;font-size:14px;">{title}</div>
<div style="font-size:12px;color:#666;margin:4px 0;padding-left:8px;border-left:2px solid #dc3545;"><strong>Stores:</strong> {stores}</div>
<div style="font-size:12px;color:#333;margin:8px 0;line-height:1.4;padding-left:8px;border-left:2px solid #dc3545;"><strong>Notes:</strong> {notes}</div>
<div style="font-size:12px;color:#666;margin:4px 0;padding-left:8px;border-left:2px solid #dc3545;"><strong>WM Week:</strong> {wm_week}</div>
<div style="font-size:12px;color:#666;margin:4px 0;padding-left:8px;border-left:2px solid #dc3545;"><strong>Status:</strong> {deployment}</div>
</div>'''
        
        page_label = f' (Page {page_idx + 1}/{total_pages})' if total_pages > 1 else ''
        
        html = f'''<!DOCTYPE html><html><head><meta charset="utf-8">
<style>*{{margin:0;padding:0;box-sizing:border-box;}} body{{background:white;font-family:-apple-system,BlinkMacSystemFont,'Segoe UI','Helvetica Neue',Arial,sans-serif;}}</style>
</head><body>
<div style="background:#fff3cd;border-left:6px solid #dc3545;border-radius:8px;padding:20px;box-shadow:0 2px 4px rgba(220,53,69,0.15);">
  <div style="color:#dc3545;font-size:18px;font-weight:700;margin-bottom:16px;">Needs Attention{page_label}</div>
  <div style="display:grid;grid-template-columns:repeat(auto-fit,minmax(300px,1fr));gap:12px;">
    {cards_html}
  </div>
</div>
</body></html>'''
        
        pages.append(html)
    
    return pages


def build_executive_summary_html(stats: dict, projects: list) -> str:
    """Build HTML for the Executive Summary slide — matches dashboard layout exactly.
    Includes: header bar + stats cards + needs attention section.
    """
    # Load Spark logo as base64 for embedding
    import base64 as _b64
    spark_logo_b64 = ''
    spark_path = os.path.join(os.path.dirname(__file__), 'Spark_Blank.png')
    if os.path.exists(spark_path):
        with open(spark_path, 'rb') as _f:
            spark_logo_b64 = _b64.b64encode(_f.read()).decode('ascii')

    total = stats.get('total_projects', 0)
    stores = stats.get('total_stores', 0)
    on_track = stats.get('on_track', 0)
    at_risk = stats.get('at_risk', 0)
    off_track = stats.get('off_track', 0)
    continuous = stats.get('continuous', 0) if 'continuous' in stats else 0
    
    # At-risk items for Needs Attention (first 10 only — overflow goes to separate slides)
    at_risk_items = [p for p in projects if str(p.get('Health Status', '')).lower() == 'at risk']
    display_items = at_risk_items[:10]
    
    cards_html = ''
    for item in display_items:
        title = escape(item.get('Initiative - Project Title', 'Unknown'))
        stores_val = item.get('# of Stores', 0)
        notes = escape(item.get('Executive Notes', 'No notes provided'))
        phase = escape(item.get('Phase', 'Unknown'))
        wm_week = escape(str(item.get('WM Week', 'N/A')))
        status_label = escape(str(item.get('Health Status', 'At Risk')))
        deployment = escape(str(item.get('Deployment', 'No Note Provided')))
        
        cards_html += f'''<div style="background:white;border:1px solid #dc3545;border-radius:6px;padding:16px;box-shadow:0 1px 3px rgba(220,53,69,0.1);">
<div style="margin-bottom:8px;">
  <div style="display:inline-block;background:#dc3545;color:white;padding:4px 8px;border-radius:4px;font-size:12px;font-weight:600;">{status_label}</div>
  <div style="font-size:12px;color:#666;margin-top:4px;">Phase: {phase}</div>
</div>
<div style="font-weight:700;color:#212121;margin-bottom:8px;font-size:14px;">{title}</div>
<div style="font-size:12px;color:#666;margin:4px 0;padding-left:8px;border-left:2px solid #dc3545;"><strong>Stores:</strong> {stores_val}</div>
<div style="font-size:12px;color:#333;margin:8px 0;line-height:1.4;padding-left:8px;border-left:2px solid #dc3545;"><strong>Notes:</strong> {notes}</div>
<div style="font-size:12px;color:#666;margin:4px 0;padding-left:8px;border-left:2px solid #dc3545;"><strong>WM Week:</strong> {wm_week}</div>
<div style="font-size:12px;color:#666;margin:4px 0;padding-left:8px;border-left:2px solid #dc3545;"><strong>Status:</strong> {deployment}</div>
</div>'''
    
    needs_attention_section = ''
    if cards_html:
        needs_attention_section = f'''
<div style="background:#fff3cd;border-left:6px solid #dc3545;border-radius:8px;padding:20px;margin-top:20px;box-shadow:0 2px 4px rgba(220,53,69,0.15);">
  <div style="color:#dc3545;font-size:18px;font-weight:700;margin-bottom:16px;">🚨 Needs Attention</div>
  <div style="display:grid;grid-template-columns:repeat(auto-fit,minmax(220px,1fr));gap:12px;">
    {cards_html}
  </div>
</div>'''
    
    html = f'''<!DOCTYPE html><html><head><meta charset="utf-8">
<style>
*{{margin:0;padding:0;box-sizing:border-box;}}
body{{background:white;font-family:-apple-system,BlinkMacSystemFont,'Segoe UI','Helvetica Neue',Arial,sans-serif;}}
</style>
</head><body>

<!-- Header Bar -->
<div style="background:linear-gradient(135deg,#1E3A8A 0%,#2563EB 100%);padding:20px 30px;display:flex;align-items:center;gap:16px;">
  <div style="width:40px;height:40px;background:#FFC220;border-radius:50%;display:flex;align-items:center;justify-content:center;">
    <img src="data:image/png;base64,{spark_logo_b64}" style="width:24px;height:24px;" alt="Spark">
  </div>
  <div>
    <div style="color:white;font-size:22px;font-weight:700;">Dallas Team Report</div>
    <div style="color:rgba(255,255,255,0.8);font-size:13px;">TDA Initiatives - Dallas VET Focus</div>
  </div>
</div>

<!-- Stats Cards Row -->
<div style="display:flex;gap:0;border:1px solid #e5e5e5;margin:20px 0;">
  <div style="flex:1;text-align:center;padding:16px;border-right:1px solid #e5e5e5;">
    <div style="font-size:11px;font-weight:600;color:#666;text-transform:uppercase;letter-spacing:0.5px;">Total Initiatives</div>
    <div style="font-size:28px;font-weight:700;color:#212121;">{total}</div>
  </div>
  <div style="flex:1;text-align:center;padding:16px;border-right:1px solid #e5e5e5;">
    <div style="font-size:11px;font-weight:600;color:#666;text-transform:uppercase;letter-spacing:0.5px;">On Track</div>
    <div style="font-size:28px;font-weight:700;color:#107C10;">{on_track}</div>
  </div>
  <div style="flex:1;text-align:center;padding:16px;border-right:1px solid #e5e5e5;">
    <div style="font-size:11px;font-weight:600;color:#666;text-transform:uppercase;letter-spacing:0.5px;">At Risk</div>
    <div style="font-size:28px;font-weight:700;color:#F7630C;">{at_risk}</div>
  </div>
  <div style="flex:1;text-align:center;padding:16px;border-right:1px solid #e5e5e5;">
    <div style="font-size:11px;font-weight:600;color:#666;text-transform:uppercase;letter-spacing:0.5px;">Off Track</div>
    <div style="font-size:28px;font-weight:700;color:#DC3545;">{off_track}</div>
  </div>
  <div style="flex:1;text-align:center;padding:16px;">
    <div style="font-size:11px;font-weight:600;color:#666;text-transform:uppercase;letter-spacing:0.5px;">Continuous</div>
    <div style="font-size:28px;font-weight:700;color:#212121;">{continuous}</div>
  </div>
</div>

<!-- Needs Attention -->
{needs_attention_section}

</body></html>'''
    
    return html


def build_phase_html(phase: str, rows: list) -> str:
    """Build HTML for a phase section — matches dashboard Generate PPT styling exactly"""
    
    columns = ['Initiative - Project Title', 'Health Status', 'Phase', 'WM Week', '# of Stores', 'Executive Notes']
    
    # Phase banner — matches dashboard blue (#3B82F6)
    banner = f'''<div style="background:#3B82F6;color:white;padding:10px 20px;font-size:16px;font-weight:600;text-align:center;border-radius:4px 4px 0 0;">Phase: {escape(phase)}</div>'''
    
    # Build table matching dashboard styling (14px font, 16px 24px padding)
    table = '<table style="width:100%;border-collapse:collapse;font-size:14px;font-family:-apple-system,BlinkMacSystemFont,\'Segoe UI\',\'Helvetica Neue\',Arial,sans-serif;">'
    table += '<thead style="background-color:#F5F5F5;border-bottom:2px solid #E5E5E5;"><tr>'
    
    for col in columns:
        table += f'<th style="padding:16px 24px;text-align:left;font-weight:600;color:#212121;white-space:nowrap;">{escape(col)}</th>'
    
    table += '</tr></thead><tbody>'
    
    for row in rows:
        table += '<tr>'
        
        for col in columns:
            value = str(row.get(col, '') or '')
            cell_style = 'padding:16px 24px;border-bottom:1px solid #E5E5E5;text-align:left;'
            
            if col == 'Health Status':
                status = value.lower()
                bg_color, color = '#f0f0f0', '#333'
                if 'on track' in status:
                    bg_color, color = 'rgba(16,124,16,0.2)', '#107C10'
                elif 'at risk' in status:
                    bg_color, color = 'rgba(247,99,12,0.2)', '#F7630C'
                elif 'off track' in status:
                    bg_color, color = 'rgba(220,53,69,0.2)', '#DC3545'
                
                table += f'<td style="{cell_style}"><span style="display:inline-block;padding:4px 10px;background-color:{bg_color};color:{color};border-radius:12px;font-weight:600;font-size:12px;white-space:nowrap;">{escape(value)}</span></td>'
            
            elif col == 'Phase':
                table += f'<td style="{cell_style}"><span style="display:inline-block;padding:4px 10px;background-color:#DBEAFE;color:#1E3A8A;border-radius:4px;font-weight:500;font-size:12px;">{escape(value)}</span></td>'
            
            elif col == '# of Stores':
                try:
                    stores = int(row.get(col, 0) or 0)
                    value = f"{stores:,}"
                except:
                    pass
                table += f'<td style="{cell_style}font-weight:600;color:#0071CE;">{escape(value)}</td>'
            
            elif col == 'Initiative - Project Title':
                project_id = row.get('Project ID', '')
                if project_id:
                    table += f'<td style="{cell_style}"><a href="https://hoops.wal-mart.com/intake-hub/projects/{escape(str(project_id))}" style="color:#0071CE;text-decoration:none;font-weight:500;">{escape(value)}</a></td>'
                else:
                    table += f'<td style="{cell_style}">{escape(value)}</td>'
            
            else:
                table += f'<td style="{cell_style}">{escape(value)}</td>'
        
        table += '</tr>'
    
    table += '</tbody></table>'
    
    # Wrap in document
    full_html = f'''<!DOCTYPE html><html><head><meta charset="utf-8">
<style>*{{margin:0;padding:0;box-sizing:border-box;}} body{{background:white;}}</style>
</head><body><div style="margin-bottom:10px;">{banner}{table}</div></body></html>'''
    
    return full_html


def build_combined_phases_html(phase_sections: list) -> str:
    """Build HTML combining multiple small phases onto one page.
    phase_sections: list of (phase_name, rows_list) tuples
    """
    columns = ['Initiative - Project Title', 'Health Status', 'Phase', 'WM Week', '# of Stores', 'Executive Notes']
    
    body = ''
    for phase, rows in phase_sections:
        banner = f'''<div style="background:#3B82F6;color:white;padding:10px 20px;font-size:16px;font-weight:600;text-align:center;border-radius:4px 4px 0 0;">Phase: {escape(phase)}</div>'''
        
        table = '<table style="width:100%;border-collapse:collapse;font-size:14px;font-family:-apple-system,BlinkMacSystemFont,\'Segoe UI\',\'Helvetica Neue\',Arial,sans-serif;">'
        table += '<thead style="background-color:#F5F5F5;border-bottom:2px solid #E5E5E5;"><tr>'
        for col in columns:
            table += f'<th style="padding:16px 24px;text-align:left;font-weight:600;color:#212121;white-space:nowrap;">{escape(col)}</th>'
        table += '</tr></thead><tbody>'
        
        for row in rows:
            table += '<tr>'
            for col in columns:
                value = str(row.get(col, '') or '')
                cell_style = 'padding:16px 24px;border-bottom:1px solid #E5E5E5;text-align:left;'
                
                if col == 'Health Status':
                    status = value.lower()
                    bg_color, color = '#f0f0f0', '#333'
                    if 'on track' in status:
                        bg_color, color = 'rgba(16,124,16,0.2)', '#107C10'
                    elif 'at risk' in status:
                        bg_color, color = 'rgba(247,99,12,0.2)', '#F7630C'
                    elif 'off track' in status:
                        bg_color, color = 'rgba(220,53,69,0.2)', '#DC3545'
                    table += f'<td style="{cell_style}"><span style="display:inline-block;padding:4px 10px;background-color:{bg_color};color:{color};border-radius:12px;font-weight:600;font-size:12px;white-space:nowrap;">{escape(value)}</span></td>'
                elif col == 'Phase':
                    table += f'<td style="{cell_style}"><span style="display:inline-block;padding:4px 10px;background-color:#DBEAFE;color:#1E3A8A;border-radius:4px;font-weight:500;font-size:12px;">{escape(value)}</span></td>'
                elif col == '# of Stores':
                    try:
                        stores = int(row.get(col, 0) or 0)
                        value = f"{stores:,}"
                    except:
                        pass
                    table += f'<td style="{cell_style}font-weight:600;color:#0071CE;">{escape(value)}</td>'
                elif col == 'Initiative - Project Title':
                    project_id = row.get('Project ID', '')
                    if project_id:
                        table += f'<td style="{cell_style}"><a href="https://hoops.wal-mart.com/intake-hub/projects/{escape(str(project_id))}" style="color:#0071CE;text-decoration:none;font-weight:500;">{escape(value)}</a></td>'
                    else:
                        table += f'<td style="{cell_style}">{escape(value)}</td>'
                else:
                    table += f'<td style="{cell_style}">{escape(value)}</td>'
            table += '</tr>'
        
        table += '</tbody></table>'
        body += f'<div style="margin-bottom:20px;">{banner}{table}</div>'
    
    return f'''<!DOCTYPE html><html><head><meta charset="utf-8">
<style>*{{margin:0;padding:0;box-sizing:border-box;}} body{{background:white;}}</style>
</head><body>{body}</body></html>'''


def capture_html_screenshot(html_content: str, output_png: str, width: int = 1280) -> bool:
    """Use Edge headless to capture HTML as PNG"""
    try:
        with tempfile.NamedTemporaryFile(suffix='.html', delete=False, mode='w', encoding='utf-8') as f:
            f.write(html_content)
            html_path = f.name
        
        cmd = [
            EDGE_PATH,
            '--headless',
            '--disable-gpu',
            '--no-sandbox',
            '--hide-scrollbars',
            '--force-device-scale-factor=1.5',
            f'--screenshot={output_png}',
            f'--window-size={width},2000',
            f'file:///{html_path.replace(os.sep, "/")}'
        ]
        
        result = subprocess.run(cmd, capture_output=True, timeout=30)
        
        if not Path(output_png).exists():
            raise RuntimeError("Edge screenshot failed")
        
        # Auto-crop
        img = Image.open(output_png)
        bg = Image.new(img.mode, img.size, (255, 255, 255))
        diff = ImageChops.difference(img, bg)
        bbox = diff.getbbox()
        
        if bbox:
            cropped = img.crop((0, 0, bbox[2], bbox[3] + 4))
            jpg_path = output_png.replace('.png', '.jpg')
            cropped.convert('RGB').save(jpg_path, 'JPEG', quality=78, optimize=True)
            import shutil
            shutil.move(jpg_path, output_png)
        
        os.unlink(html_path)
        return True
    
    except Exception as e:
        print(f"     [!] Screenshot failed: {e}")
        return False


def generate_report_pptx(executive_summary_html: str, sections: list) -> tuple:
    """Generate PPTX: Title + Executive Summary + Phase slides"""
    try:
        print("     [*] Generating PPT from dashboard screenshots...")
        
        prs = Presentation()
        prs.slide_width = Inches(9.6)
        prs.slide_height = Inches(7.2)
        SLIDE_WIDTH = prs.slide_width
        SLIDE_HEIGHT = prs.slide_height
        
        # ── Title Slide ──
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0x1E, 0x3A, 0x8A)
        
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(8.6), Inches(1.8))
        p = txBox.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = "V.E.T. Executive Report"
        run.font.size = Pt(44)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        
        txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(8.6), Inches(0.8))
        p2 = txBox2.text_frame.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = "Walmart Enterprise Transformation"
        run2.font.size = Pt(20)
        run2.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        
        accent = slide.shapes.add_shape(1, Emu(0), Inches(6.7), SLIDE_WIDTH, Inches(0.5))
        accent.fill.solid()
        accent.fill.fore_color.rgb = RGBColor(0xFF, 0xC2, 0x20)
        accent.line.fill.background()
        
        screenshots = []
        temp_files = []
        
        # ── Executive Summary Slide ──
        if executive_summary_html:
            png_path = OUTPUT_DIR / "temp_executive_summary.png"
            temp_files.append(png_path)
            
            if capture_html_screenshot(executive_summary_html, str(png_path)):
                if Path(png_path).exists():
                    png_bytes = Path(png_path).read_bytes()
                    img = Image.open(io.BytesIO(png_bytes))
                    img_w, img_h = img.size
                    aspect = img_h / img_w if img_w > 0 else 0.5625
                    final_w = SLIDE_WIDTH
                    final_h = Emu(int(SLIDE_WIDTH * aspect))
                    if final_h > SLIDE_HEIGHT:
                        final_h = SLIDE_HEIGHT
                    
                    slide = prs.slides.add_slide(prs.slide_layouts[6])
                    slide.shapes.add_picture(io.BytesIO(png_bytes), Emu(0), Emu(0), final_w, final_h)
                    screenshots.append(("Executive Summary", png_bytes))
        
        # ── Phase Slides ──
        for section_name, html_content in sections:
            png_path = OUTPUT_DIR / f"temp_{section_name.replace('/', '')}.png"
            temp_files.append(png_path)
            
            if capture_html_screenshot(html_content, str(png_path)):
                if Path(png_path).exists():
                    png_bytes = Path(png_path).read_bytes()
                    img = Image.open(io.BytesIO(png_bytes))
                    img_w, img_h = img.size
                    aspect = img_h / img_w if img_w > 0 else 0.5625
                    final_w = SLIDE_WIDTH
                    final_h = Emu(int(SLIDE_WIDTH * aspect))
                    if final_h > SLIDE_HEIGHT:
                        final_h = SLIDE_HEIGHT
                    
                    slide = prs.slides.add_slide(prs.slide_layouts[6])
                    slide.shapes.add_picture(io.BytesIO(png_bytes), Emu(0), Emu(0), final_w, final_h)
                    screenshots.append((section_name, png_bytes))
        
        # Save PPT
        pptx_buffer = io.BytesIO()
        prs.save(pptx_buffer)
        pptx_buffer.seek(0)
        
        # Cleanup
        for temp_file in temp_files:
            if temp_file.exists():
                try:
                    os.unlink(temp_file)
                except:
                    pass
        
        print(f"     [OK] PPT generated with {len(screenshots)} sections")
        return pptx_buffer.getvalue(), screenshots
    
    except Exception as e:
        print(f"     [!] Error generating PPT: {e}")
        import traceback
        traceback.print_exc()
        return None, []


def generate_report_pdf(screenshots: list) -> bytes:
    """Convert screenshots to PDF (same as PPT)"""
    try:
        if not screenshots:
            return None
        
        print("     [*] Generating PDF from screenshots...")
        
        pages = []
        for section_name, png_bytes in screenshots:
            img = Image.open(io.BytesIO(png_bytes)).convert('RGB')
            pages.append(img)
        
        pdf_buffer = io.BytesIO()
        if pages:
            pages[0].save(pdf_buffer, 'PDF', save_all=True, append_images=pages[1:] if len(pages) > 1 else [], resolution=150)
            pdf_buffer.seek(0)
            print(f"     [OK] PDF generated with {len(pages)} pages")
            return pdf_buffer.getvalue()
    
    except Exception as e:
        print(f"     [!] Error generating PDF: {e}")
        return None


def send_vet_report_email(recipient_emails: list = None, test_mode: bool = False):
    """Main email generation and sending"""
    
    if recipient_emails is None:
        recipient_emails = [
            "kendall.rush@walmart.com",
            "Matthew.Farnworth@walmart.com",
            "Tina.Budnaitis@walmart.com"
        ]
    
    print("=" * 90)
    print("DALLAS TEAM REPORT - FINAL VERSION")
    print("=" * 90)
    print()
    
    try:
        # Step 1: Fetch data
        print("[1/6] Fetching dashboard data...")
        stats, projects = fetch_dashboard_data()
        if not stats or not projects:
            print("     [!] Failed to fetch data")
            return False
        
        print("     [OK] Data fetched successfully")
        print(f"       * Total Projects: {stats['total_projects']}")
        print(f"       * On Track: {stats['on_track']}")
        print(f"       * At Risk: {stats['at_risk']}")
        print(f"       * Off Track: {stats['off_track']}")
        print(f"       * WM Week: {stats['wm_week']}")
        print()
        
        # Step 2: Build Executive Summary (header + stats + needs attention)
        print("[2/6] Building Executive Summary slide...")
        executive_summary_html = build_executive_summary_html(stats, projects)
        print("     [OK] Executive Summary built (header + stats + needs attention)")
        
        # Build overflow Needs Attention pages (items beyond the first 10 shown in exec summary)
        at_risk_all = [p for p in projects if str(p.get('Health Status', '')).lower() == 'at risk']
        needs_attention_overflow = []
        if len(at_risk_all) > 10:
            overflow_pages = build_needs_attention_html(at_risk_all[10:])
            for i, page_html in enumerate(overflow_pages):
                label = f"Needs Attention (cont. {i + 2})"
                needs_attention_overflow.append((label, page_html))
            print(f"     [OK] {len(needs_attention_overflow)} overflow Needs Attention page(s)")
        print()
        
        # Step 3: Build Phase sections — combine small phases onto one slide
        print("[3/6] Capturing dashboard phase tables...")
        phases_dict = {}
        for proj in projects:
            phase = proj.get('Phase', 'Unknown')
            if phase not in phases_dict:
                phases_dict[phase] = []
            phases_dict[phase].append(proj)
        
        # Build ordered list of (phase, rows) tuples
        phase_order = ['Pending', 'Vet', 'Test', 'Test Markets', 'Roll/Deploy']
        ordered_phases = []
        for phase in phase_order:
            if phase in phases_dict and phases_dict[phase]:
                ordered_phases.append((phase, phases_dict[phase]))
        
        # Combine small phases (≤ 5 rows each, ≤ 10 rows combined) onto one slide
        # Split large phases into pages of MAX_ROWS_PER_PAGE rows each
        MAX_ROWS_PER_PAGE = 8
        MAX_ROWS_TO_COMBINE = 5
        MAX_ROWS_COMBINED_TOTAL = 10
        sections = []
        pending_combine = []
        pending_rows = 0
        
        for phase, rows in ordered_phases:
            # Split large phases into pages first
            if len(rows) > MAX_ROWS_PER_PAGE:
                # Flush pending combines first
                if pending_combine:
                    if len(pending_combine) == 1:
                        html = build_phase_html(pending_combine[0][0], pending_combine[0][1])
                        sections.append((pending_combine[0][0], html))
                    else:
                        html = build_combined_phases_html(pending_combine)
                        names = ' + '.join(p[0] for p in pending_combine)
                        sections.append((names, html))
                    pending_combine = []
                    pending_rows = 0
                
                # Paginate large phase
                for i in range(0, len(rows), MAX_ROWS_PER_PAGE):
                    page_rows = rows[i:i + MAX_ROWS_PER_PAGE]
                    page_num = (i // MAX_ROWS_PER_PAGE) + 1
                    total_pages = -(-len(rows) // MAX_ROWS_PER_PAGE)  # ceil division
                    label = f"{phase} ({page_num}/{total_pages})" if total_pages > 1 else phase
                    html = build_phase_html(phase, page_rows)
                    sections.append((label, html))
            elif len(rows) <= MAX_ROWS_TO_COMBINE and pending_rows + len(rows) <= MAX_ROWS_COMBINED_TOTAL:
                # Small enough to combine
                pending_combine.append((phase, rows))
                pending_rows += len(rows)
            else:
                # Flush any pending combined sections first
                if pending_combine:
                    if len(pending_combine) == 1:
                        html = build_phase_html(pending_combine[0][0], pending_combine[0][1])
                        sections.append((pending_combine[0][0], html))
                    else:
                        html = build_combined_phases_html(pending_combine)
                        names = ' + '.join(p[0] for p in pending_combine)
                        sections.append((names, html))
                    pending_combine = []
                    pending_rows = 0
                
                # This phase is too big to combine — make it its own slide
                if len(rows) <= MAX_ROWS_TO_COMBINE:
                    pending_combine.append((phase, rows))
                    pending_rows = len(rows)
                else:
                    html = build_phase_html(phase, rows)
                    sections.append((phase, html))
        
        # Flush remaining
        if pending_combine:
            if len(pending_combine) == 1:
                html = build_phase_html(pending_combine[0][0], pending_combine[0][1])
                sections.append((pending_combine[0][0], html))
            else:
                html = build_combined_phases_html(pending_combine)
                names = ' + '.join(p[0] for p in pending_combine)
                sections.append((names, html))
        
        print(f"     [OK] Built {len(sections)} slide sections (combined small phases)")
        
        # Insert overflow Needs Attention pages before phase slides
        if needs_attention_overflow:
            sections = needs_attention_overflow + sections
        
        for name, _ in sections:
            print(f"       * {name}")
        print()
        
        # Step 4: Generate PPT
        print("[4/6] Generating PowerPoint...")
        wk = stats['wm_week']
        pptx_filename = f"Dallas_Team_Report_{wk}.pptx"
        pdf_filename = f"Dallas_Team_Report_{wk}.pdf"
        
        pptx_path = OUTPUT_DIR / pptx_filename
        pdf_path = OUTPUT_DIR / pdf_filename
        
        pptx_data, screenshots = generate_report_pptx(executive_summary_html, sections)
        
        if pptx_data:
            pptx_path.write_bytes(pptx_data)
            print(f"     [OK] PPT saved: {pptx_filename}")
        else:
            print("     [!] Failed to generate PPT")
            return False
        
        print()
        
        # Step 5: Generate PDF
        print("[5/6] Generating PDF...")
        pdf_data = generate_report_pdf(screenshots)
        
        if pdf_data:
            pdf_path.write_bytes(pdf_data)
            print(f"     [OK] PDF saved: {pdf_filename}")
        else:
            print("     [!] PDF generation failed")
        
        print()
        
        # Step 6: Send email
        print("[6/6] Sending email...")
        email_service = VETEmailService()
        
        mode_text = "DRAFT" if test_mode else "SEND"
        print(f"     Email Configuration:")
        print(f"       * Mode: {mode_text}")
        print(f"       * Recipients: {', '.join(recipient_emails)}")
        print(f"       * Subject: Dallas Team Report - {wk}")
        
        attachments = [str(pptx_path)]
        if pdf_data:
            attachments.append(str(pdf_path))
        
        print(f"       * Attachments: {', '.join([Path(a).name for a in attachments])}")
        print()
        
        # Prepare report data with dashboard HTML
        report_data = stats.copy()
        report_data['dashboard_html'] = executive_summary_html  # Pass executive summary HTML
        
        # Send
        success = email_service.send_report_email(
            to_recipients=recipient_emails,
            report_data=report_data,
            attachment_paths=attachments,
            test_mode=test_mode
        )
        
        if success:
            print()
            print("=" * 90)
            print("[OK] EMAIL SENT SUCCESSFULLY")
            print("=" * 90)
            print()
            print(f"   Recipients: {', '.join(recipient_emails)}")
            print(f"   Subject: Dallas Team Report - {wk}")
            print(f"   Data: {stats['total_projects']} projects")
            print()
            return True
        else:
            print()
            print("=" * 90)
            print("[ERROR] FAILED TO SEND EMAIL")
            print("=" * 90)
            return False
    
    except Exception as e:
        print(f"[ERROR] {e}")
        import traceback
        traceback.print_exc()
        return False


def main():
    """Main entry point"""
    recipients = None  # Uses default list: Kendall, Matthew, Tina
    test_mode = False
    
    for arg in sys.argv[1:]:
        if arg == '--draft':
            test_mode = True
        elif arg == '--email':
            idx = sys.argv.index(arg)
            if idx + 1 < len(sys.argv):
                recipients = [sys.argv[idx + 1]]
    
    success = send_vet_report_email(recipient_emails=recipients, test_mode=test_mode)
    sys.exit(0 if success else 1)


if __name__ == '__main__':
    main()
