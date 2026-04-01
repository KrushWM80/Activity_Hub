"""Compare dashboard-generated PPT vs our generated PPT"""
from pptx import Presentation
from pptx.util import Inches, Emu
from PIL import Image
import io

# Dashboard-generated PPT (reference)
ref = Presentation('reports/VET_Executive_Report (3).pptx')
print('=== DASHBOARD PPT (Reference - 2,078 KB) ===')
print(f'Slide size: {ref.slide_width / 914400:.1f}" x {ref.slide_height / 914400:.1f}"')
print(f'Slides: {len(ref.slides)}')
for i, slide in enumerate(ref.slides):
    shapes = list(slide.shapes)
    print(f'  Slide {i+1}: {len(shapes)} shapes')
    for s in shapes:
        if s.shape_type == 13:  # Picture
            print(f'    Picture: {s.width/914400:.1f}"x{s.height/914400:.1f}" at ({s.left/914400:.1f}",{s.top/914400:.1f}")')
            img = Image.open(io.BytesIO(s.image.blob))
            print(f'      Image: {img.size[0]}x{img.size[1]}px, {len(s.image.blob):,} bytes')
        elif hasattr(s, 'text') and s.text:
            txt = s.text[:80]
            print(f'    Text: "{txt}"')

print()

# Our generated PPT
ours = Presentation('reports/VET_Executive_Report_WK09.pptx')
print('=== OUR PPT (Generated - Current) ===')
print(f'Slide size: {ours.slide_width / 914400:.1f}" x {ours.slide_height / 914400:.1f}"')
print(f'Slides: {len(ours.slides)}')
for i, slide in enumerate(ours.slides):
    shapes = list(slide.shapes)
    print(f'  Slide {i+1}: {len(shapes)} shapes')
    for s in shapes:
        if s.shape_type == 13:  # Picture
            print(f'    Picture: {s.width/914400:.1f}"x{s.height/914400:.1f}" at ({s.left/914400:.1f}",{s.top/914400:.1f}")')
            img = Image.open(io.BytesIO(s.image.blob))
            print(f'      Image: {img.size[0]}x{img.size[1]}px, {len(s.image.blob):,} bytes')
        elif hasattr(s, 'text') and s.text:
            txt = s.text[:80]
            print(f'    Text: "{txt}"')

print()
print('=== KEY DIFFERENCES ===')
print(f'Dashboard: {len(ref.slides)} slides, {ref.slide_width/914400:.1f}"x{ref.slide_height/914400:.1f}"')
print(f'Ours:      {len(ours.slides)} slides, {ours.slide_width/914400:.1f}"x{ours.slide_height/914400:.1f}"')
