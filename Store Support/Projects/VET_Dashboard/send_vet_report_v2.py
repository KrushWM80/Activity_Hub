#!/usr/bin/env python3
"""
V.E.T. Executive Report - Email Reporter (Dashboard Screenshot Version)
Captures V.E.T. Dashboard tables using Edge headless, generates PPT+PDF, sends via Outlook
Follows TDA Insights approach for consistency and quality

Usage:
  python send_vet_report_v2.py              # Send email
  python send_vet_report_v2.py --draft      # Save as draft only
  python send_vet_report_v2.py --email user@example.com  # Send to custom recipient
"""

import sys
import os
import json
import subprocess
import tempfile
import io
import base64
import re
from pathlib import Path
from datetime import datetime, timedelta
from xml.sax.saxutils import escape
import urllib.request
import urllib.error
from PIL import Image, ImageChops

# Imports
sys.path.insert(0, str(Path(__file__).parent))
from email_service import VETEmailService
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Configuration
SCRIPT_DIR = Path(__file__).parent
DASHBOARD_URL = "http://127.0.0.1:5001"
EDGE_PATH = r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe"
OUTPUT_DIR = SCRIPT_DIR / 'reports'
OUTPUT_DIR.mkdir(exist_ok=True)

COLORS = {
    'walmart_blue_dark': '#1E3A8A',
    'walmart_blue': '#0071CE',
    'walmart_yellow': '#FFCC00',
    'on_track': '#107C10',
    'at_risk': '#F7630C',
    'off_track': '#DC3545',
    'continuous': '#6F42C1',
}


def get_current_walmart_week():
    """Calculate current Walmart Week (WK format)"""
    today = datetime.now().date()
    
    # Walmart fiscal year starts in February
    if today.month >= 2:
        year_start = datetime(today.year, 2, 1).date()
    else:
        year_start = datetime(today.year - 1, 2, 1).date()
    
    # Find first Saturday
    days_until_saturday = (5 - year_start.weekday()) % 7
    first_saturday = year_start + timedelta(days=days_until_saturday)
    
    # Calculate week number
    days_diff = (today - first_saturday).days
    weeks = max(1, (days_diff // 7) + 1)
    
    return f"WK{weeks:02d}"


def fetch_dashboard_stats(api_url: str = DASHBOARD_URL) -> dict:
    """Fetch current dashboard statistics from the running backend API"""
    try:
        summary_url = f"{api_url}/api/summary"
        with urllib.request.urlopen(summary_url, timeout=5) as response:
            response_data = json.loads(response.read().decode())
        
        summary = response_data.get('summary', {})
        health_status = summary.get('by_health_status', {})
        
        current_wm_week = get_current_walmart_week()
        
        return {
            'total_projects': summary.get('total_projects', 0),
            'total_stores': summary.get('total_stores', 0),
            'on_track': health_status.get('On Track', 0),
            'at_risk': health_status.get('At Risk', 0),
            'off_track': health_status.get('Off Track', 0),
            'continuous': 0,
            'wm_week': current_wm_week
        }
    
    except Exception as e:
        print(f"     [!] Error fetching dashboard data: {e}")
        return None


def capture_dashboard_html(api_url: str = DASHBOARD_URL) -> list:
    """
    Capture dashboard table data as HTML for screenshot conversion.
    Returns list of (section_name, html_content) tuples, organized by phase.
    """
    try:
        # Fetch all project data with All phases
        data_url = f"{api_url}/api/data?phase=All"
        with urllib.request.urlopen(data_url, timeout=5) as response:
            response_data = json.loads(response.read().decode())
        
        projects = response_data.get('data', [])
        
        # Group by phase
        phases_dict = {}
        for proj in projects:
            phase = proj.get('Phase', 'Unknown')
            if phase not in phases_dict:
                phases_dict[phase] = []
            phases_dict[phase].append(proj)
        
        sections = []
        
        # For each phase in order, create HTML table
        phase_order = ['Pending', 'POC/POT', 'Test', 'Mkt Scale', 'Roll/Deploy']
        for phase in phase_order:
            if phase not in phases_dict or not phases_dict[phase]:
                continue
            
            rows = phases_dict[phase]
            html = _build_phase_html(phase, rows)
            sections.append((phase, html))
        
        return sections
    
    except Exception as e:
        print(f"     [!] Error capturing dashboard HTML: {e}")
        import traceback
        traceback.print_exc()
        return []


def _build_phase_html(phase: str, rows: list) -> str:
    """Build HTML for a phase section matching V.E.T. Dashboard styling"""
    
    columns = ['Initiative - Project Title', 'Health Status', 'Phase', '# of Stores']
    
    # Phase banner (dark blue)
    banner = f'''<div style="background:#1E3A8A;color:white;padding:15px 20px;font-size:18px;font-weight:700;text-align:center;border-radius:4px 4px 0 0;">
    {escape(phase)} ({len(rows)} Projects)
</div>'''
    
    # Build table
    table = '<table style="width:100%;border-collapse:collapse;font-size:14px;font-family:-apple-system,BlinkMacSystemFont,\'Segoe UI\',\'Helvetica Neue\',Arial,sans-serif;">'
    table += '<thead style="background-color:#F5F5F5;border-bottom:2px solid #E5E5E5;"><tr>'
    
    for col in columns:
        table += f'<th style="padding:16px 24px;text-align:left;font-weight:600;color:#212121;white-space:nowrap;">{escape(col)}</th>'
    
    table += '</tr></thead><tbody>'
    
    for i, row in enumerate(rows):
        bg = '#ffffff' if i % 2 == 0 else '#fafafa'
        table += f'<tr style="background:{bg};">'
        
        for col in columns:
            # Get value from row using the exact key
            value = str(row.get(col, '') or '')
            cell_style = 'padding:16px 24px;border-bottom:1px solid #E5E5E5;text-align:left;'
            
            if col == 'Health Status':
                status = value.lower()
                bg_color, color = '#f0f0f0', '#333'
                if 'on track' in status:
                    bg_color, color = 'rgba(16,124,16,0.2)', '#107C10'
                elif 'at risk' in status:
                    bg_color, color = 'rgba(247,99,12,0.2)', '#F7630C'
                elif 'off track' in status:
                    bg_color, color = 'rgba(220,53,69,0.2)', '#DC3545'
                
                table += f'<td style="{cell_style}"><span style="display:inline-block;padding:4px 10px;background-color:{bg_color};color:{color};border-radius:12px;font-weight:600;font-size:12px;white-space:nowrap;">{escape(value)}</span></td>'
            
            elif col == '# of Stores':
                # Format as number with commas
                try:
                    stores = int(row.get(col, 0) or 0)
                    value = f"{stores:,}"
                except:
                    pass
                table += f'<td style="{cell_style}font-weight:600;color:#0071CE;">{escape(value)}</td>'
            
            else:
                table += f'<td style="{cell_style}">{escape(value)}</td>'
        
        table += '</tr>'
    
    table += '</tbody></table>'
    
    # Wrap in HTML document
    full_html = f'''<!DOCTYPE html><html><head><meta charset="utf-8">
<style>*{{margin:0;padding:0;box-sizing:border-box;}} body{{background:white;}}</style>
</head><body>{banner}{table}</body></html>'''
    
    return full_html


def capture_html_screenshot(html_content: str, output_png: str, width: int = 1280) -> bool:
    """Use Edge headless to capture HTML as PNG screenshot"""
    try:
        with tempfile.NamedTemporaryFile(suffix='.html', delete=False, mode='w', encoding='utf-8') as f:
            f.write(html_content)
            html_path = f.name
        
        cmd = [
            EDGE_PATH,
            '--headless',
            '--disable-gpu',
            '--no-sandbox',
            '--hide-scrollbars',
            '--force-device-scale-factor=1.5',
            f'--screenshot={output_png}',
            f'--window-size={width},2000',
            f'file:///{html_path.replace(os.sep, "/")}'
        ]
        
        result = subprocess.run(cmd, capture_output=True, timeout=30)
        
        if not Path(output_png).exists():
            raise RuntimeError(f"Edge screenshot failed")
        
        # Auto-crop whitespace
        img = Image.open(output_png)
        bg = Image.new(img.mode, img.size, (255, 255, 255))
        diff = ImageChops.difference(img, bg)
        bbox = diff.getbbox()
        
        if bbox:
            cropped = img.crop((0, 0, bbox[2], bbox[3] + 4))
            jpg_path = output_png.replace('.png', '.jpg')
            cropped.convert('RGB').save(jpg_path, 'JPEG', quality=78, optimize=True)
            import shutil
            shutil.move(jpg_path, output_png)
        
        os.unlink(html_path)
        return True
    
    except Exception as e:
        print(f"     [!] Screenshot capture failed: {e}")
        return False


def generate_report_pptx(sections: list) -> tuple:
    """Generate PPTX with title slide + screenshot slides"""
    try:
        print("     [*] Generating PPT from dashboard screenshots...")
        
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        SLIDE_WIDTH = prs.slide_width
        SLIDE_HEIGHT = prs.slide_height
        
        # Title slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0x1E, 0x3A, 0x8A)
        
        # Main title
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.8))
        p = txBox.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = "V.E.T. Executive Report"
        run.font.size = Pt(44)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        
        # Subtitle
        txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(0.8))
        p2 = txBox2.text_frame.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = "Walmart Enterprise Transformation"
        run2.font.size = Pt(20)
        run2.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        run2.font.bold = False
        
        # Yellow accent bar
        accent = slide.shapes.add_shape(1, Emu(0), Inches(7.0), SLIDE_WIDTH, Inches(0.5))
        accent.fill.solid()
        accent.fill.fore_color.rgb = RGBColor(0xFF, 0xC2, 0x20)
        accent.line.fill.background()
        
        screenshots = []
        temp_files = []
        
        # Add content slides from screenshots
        for section_name, html_content in sections:
            png_path = OUTPUT_DIR / f"temp_{section_name.replace('/', '')}.png"
            temp_files.append(png_path)
            
            if capture_html_screenshot(html_content, str(png_path)):
                if Path(png_path).exists():
                    # Read file into memory before adding to PPT
                    png_bytes = Path(png_path).read_bytes()
                    
                    img = Image.open(io.BytesIO(png_bytes))
                    img_w, img_h = img.size
                    aspect = img_h / img_w if img_w > 0 else 0.5625
                    final_w = SLIDE_WIDTH
                    final_h = Emu(int(SLIDE_WIDTH * aspect))
                    if final_h > SLIDE_HEIGHT:
                        final_h = SLIDE_HEIGHT
                    
                    slide = prs.slides.add_slide(prs.slide_layouts[6])
                    slide.shapes.add_picture(io.BytesIO(png_bytes), Emu(0), Emu(0), final_w, final_h)
                    
                    screenshots.append((section_name, png_bytes))
        
        # Save PPT
        pptx_buffer = io.BytesIO()
        prs.save(pptx_buffer)
        pptx_buffer.seek(0)
        
        # Cleanup temp files after PPT is generated
        for temp_file in temp_files:
            if temp_file.exists():
                try:
                    os.unlink(temp_file)
                except:
                    pass
        
        print(f"     [OK] PPT generated with {len(screenshots)} sections")
        return pptx_buffer.getvalue(), screenshots
    
    except Exception as e:
        print(f"     [!] Error generating PPT: {e}")
        import traceback
        traceback.print_exc()
        return None, []


def generate_report_pdf(screenshots: list) -> bytes:
    """Convert screenshots to multi-page PDF"""
    try:
        if not screenshots:
            return None
        
        print("     [*] Generating PDF from screenshots...")
        
        pages = []
        for section_name, png_bytes in screenshots:
            img = Image.open(io.BytesIO(png_bytes)).convert('RGB')
            pages.append(img)
        
        pdf_buffer = io.BytesIO()
        if pages:
            pages[0].save(pdf_buffer, 'PDF', save_all=True, append_images=pages[1:] if len(pages) > 1 else [], resolution=150)
            pdf_buffer.seek(0)
            print(f"     [OK] PDF generated with {len(pages)} pages")
            return pdf_buffer.getvalue()
    
    except Exception as e:
        print(f"     [!] Error generating PDF: {e}")
        return None


def send_vet_report_email(recipient_email: str = "kendall.rush@walmart.com", test_mode: bool = False):
    """Fetch dashboard data, generate PPT+PDF, send email"""
    
    print("=" * 90)
    print("V.E.T. EXECUTIVE REPORT - EMAIL GENERATOR (Dashboard Screenshot Version)")
    print("=" * 90)
    print()
    
    try:
        # Step 1: Fetch dashboard data
        print("[1/5] Fetching live dashboard data...")
        print()
        
        stats = fetch_dashboard_stats()
        if not stats:
            print("     [!] Failed to fetch dashboard data")
            return False
        
        print("     [OK] Data fetched successfully")
        print(f"       * Total Projects: {stats['total_projects']}")
        print(f"       * Stores Impacted: {stats['total_stores']:,}")
        print(f"       * On Track: {stats['on_track']}")
        print(f"       * At Risk: {stats['at_risk']}")
        print(f"       * Off Track: {stats['off_track']}")
        print(f"       * WM Week: {stats['wm_week']}")
        print()
        
        # Step 2: Capture dashboard HTML
        print("[2/5] Capturing dashboard tables...")
        sections = capture_dashboard_html()
        if not sections:
            print("     [!] No dashboard sections captured")
            return False
        print(f"     [OK] Captured {len(sections)} phase sections")
        print()
        
        # Step 3: Generate PPT
        print("[3/5] Generating PowerPoint report...")
        wk = stats['wm_week']
        pptx_filename = f"VET_Executive_Report_{wk}.pptx"
        pdf_filename = f"VET_Executive_Report_{wk}.pdf"
        
        pptx_path = OUTPUT_DIR / pptx_filename
        pdf_path = OUTPUT_DIR / pdf_filename
        
        pptx_data, screenshots = generate_report_pptx(sections)
        
        if pptx_data:
            pptx_path.write_bytes(pptx_data)
            print(f"     [OK] PPT saved: {pptx_filename}")
        else:
            print("     [!] Failed to generate PPT")
            return False
        
        print()
        
        # Step 4: Generate PDF
        print("[4/5] Generating PDF report...")
        pdf_data = generate_report_pdf(screenshots)
        
        if pdf_data:
            pdf_path.write_bytes(pdf_data)
            print(f"     [OK] PDF saved: {pdf_filename}")
        else:
            print("     [!] PDF generation not available")
        
        print()
        
        # Step 5: Send email
        print("[5/5] Sending email...")
        email_service = VETEmailService()
        
        mode_text = "DRAFT" if test_mode else "SEND"
        print(f"     Email Configuration:")
        print(f"       * Mode: {mode_text}")
        print(f"       * Recipient: {recipient_email}")
        print(f"       * Subject: V.E.T. Executive Report - {wk}")
        attachments = [str(pptx_path)]
        if pdf_data:
            attachments.append(str(pdf_path))
        print(f"       * Attachments: {', '.join([Path(a).name for a in attachments])}")
        print()
        
        # Send email
        success = email_service.send_report_email(
            to_recipients=[recipient_email],
            report_data=stats,
            attachment_paths=attachments,
            test_mode=test_mode
        )
        
        if success:
            print()
            print("=" * 90)
            print("[OK] EMAIL SENT SUCCESSFULLY")
            print("=" * 90)
            print()
            print(f"   Recipient: {recipient_email}")
            print(f"   Subject: V.E.T. Executive Report - {wk}")
            print(f"   Data: {stats['total_projects']} projects, {stats['total_stores']:,} stores")
            print()
            return True
        else:
            print()
            print("=" * 90)
            print("[ERROR] FAILED TO SEND EMAIL")
            print("=" * 90)
            return False
    
    except Exception as e:
        print(f"[ERROR] {e}")
        import traceback
        traceback.print_exc()
        return False


def main():
    """Main entry point"""
    recipient = "kendall.rush@walmart.com"
    test_mode = False
    
    # Parse command line arguments
    for arg in sys.argv[1:]:
        if arg == '--draft':
            test_mode = True
        elif arg == '--email':
            idx = sys.argv.index(arg)
            if idx + 1 < len(sys.argv):
                recipient = sys.argv[idx + 1]
    
    success = send_vet_report_email(recipient_email=recipient, test_mode=test_mode)
    sys.exit(0 if success else 1)


if __name__ == '__main__':
    main()
