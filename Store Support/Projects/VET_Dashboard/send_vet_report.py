#!/usr/bin/env python3
"""
V.E.T. Executive Report - Weekly Email
Mirrors the TDA Insights process exactly:
  1. Fetch data from dashboard API
  2. Build HTML tables matching dashboard styling
  3. Screenshot with Edge headless (same as TDA)
  4. Build PPT from screenshots (same as TDA)
  5. Build PDF from screenshots (same as TDA)
  6. Build email HTML with summary + phase tables (same as TDA)
  7. Send via Outlook with PPT + PDF attachments
"""

import sys
import os
import io
import json
import re
import tempfile
import subprocess
from pathlib import Path
from datetime import datetime, timedelta
from xml.sax.saxutils import escape
import urllib.request
import urllib.error

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image, ImageChops

# ── Configuration ──────────────────────────────────────────────────────────────
SCRIPT_DIR = Path(__file__).parent
DASHBOARD_URL = "http://WEUS42608431466:5001"
EDGE_PATH = r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe"
OUTPUT_DIR = SCRIPT_DIR / 'reports'
OUTPUT_DIR.mkdir(exist_ok=True)

RECIPIENTS = ["Kendall.rush@walmart.com", "Matthew.Farnworth@walmart.com", "Tina.Budnaitis@walmart.com"]
PHASE_ORDER = ['Pending', 'Vet', 'Test', 'Test Markets', 'Roll/Deploy']

# Phase normalization (until BQ data reflects new names)
_PHASE_MAP = {'POC/POT': 'Vet', 'Mkt Scale': 'Test Markets'}
def _normalize_phase(phase):
    return _PHASE_MAP.get(phase, phase)

COLORS = {
    'walmart_blue': '#0071CE',
    'walmart_blue_dark': '#004C91',
    'walmart_yellow': '#FFC220',
    'on_track': '#107C10',
    'at_risk': '#F7630C',
    'off_track': '#DC3545',
}


# ── Walmart Week (from BQ Cal_Dim_Data, like TDA) ──────────────────────────────
def _walmart_week():
    try:
        from google.cloud import bigquery
        client = bigquery.Client(project='wmt-assetprotection-prod')
        query = """
        SELECT WM_WEEK_NBR
        FROM `wmt-assetprotection-prod.Store_Support_Dev.Cal_Dim_Data`
        WHERE CALENDAR_DATE = CURRENT_DATE()
        LIMIT 1
        """
        rows = list(client.query(query).result())
        if rows:
            return rows[0]['WM_WEEK_NBR']
    except Exception:
        pass
    # Fallback: manual calc
    today = datetime.now().date()
    if today.month >= 2:
        year_start = datetime(today.year, 2, 1).date()
    else:
        year_start = datetime(today.year - 1, 2, 1).date()
    days_until_sat = (5 - year_start.weekday()) % 7
    first_sat = year_start + timedelta(days=days_until_sat)
    return max(1, (today - first_sat).days // 7 + 1)

WEEK_NUM = _walmart_week()
PPT_OUTPUT = OUTPUT_DIR / f"VET_Executive_Report_WK{WEEK_NUM:02d}.pptx"
PDF_OUTPUT = OUTPUT_DIR / f"VET_Executive_Report_WK{WEEK_NUM:02d}.pdf"


# ── Data fetching ──────────────────────────────────────────────────────────────
def fetch_data():
    """Fetch project data from VET dashboard API (fallback to sample data)."""
    try:
        url = f"{DASHBOARD_URL}/api/data?phase=All"
        with urllib.request.urlopen(url, timeout=10) as resp:
            body = json.loads(resp.read().decode())
        data = body.get('data', [])
        if data:
            # Normalize phases
            for row in data:
                row['Phase'] = _normalize_phase(row.get('Phase', 'Unknown'))
            print(f"  Loaded {len(data)} projects from dashboard API")
            return data
    except Exception as e:
        print(f"  API fetch failed: {e}")

    # Fallback to sample data
    print("  Falling back to sample data...")
    sys.path.insert(0, str(SCRIPT_DIR))
    from sample_data import SAMPLE_DATA_49_PROJECTS
    # Normalize phases in sample data too
    for row in SAMPLE_DATA_49_PROJECTS:
        row['Phase'] = _normalize_phase(row.get('Phase', 'Unknown'))
    print(f"  Loaded {len(SAMPLE_DATA_49_PROJECTS)} projects from sample data")
    return SAMPLE_DATA_49_PROJECTS


# ── Email HTML (mirrors TDA build_email_html) ─────────────────────────────────
def build_email_html(data):
    """Build HTML email body with summary cards + phase tables. No dashboard cards."""
    today = datetime.now().strftime('%B %d, %Y')

    total = len(data)
    total_stores = sum(int(r.get('# of Stores', 0) or 0) for r in data)
    health_counts = {}
    for r in data:
        h = r.get('Health Status', 'Unknown')
        health_counts[h] = health_counts.get(h, 0) + 1

    on_track = health_counts.get('On Track', 0)
    at_risk = health_counts.get('At Risk', 0)
    off_track = health_counts.get('Off Track', 0)

    def badge_style(status):
        s = status.lower()
        if 'on track' in s:
            return f"background-color:rgba(16,124,16,0.15);color:{COLORS['on_track']};padding:3px 10px;border-radius:12px;font-weight:600;font-size:12px;display:inline-block;"
        elif 'at risk' in s:
            return f"background-color:rgba(247,99,12,0.15);color:{COLORS['at_risk']};padding:3px 10px;border-radius:12px;font-weight:600;font-size:12px;display:inline-block;"
        elif 'off track' in s:
            return f"background-color:rgba(220,53,69,0.15);color:{COLORS['off_track']};padding:3px 10px;border-radius:12px;font-weight:600;font-size:12px;display:inline-block;"
        return "padding:3px 10px;border-radius:12px;font-weight:600;font-size:12px;display:inline-block;"

    html = f"""<!DOCTYPE html>
<html><head><meta charset="utf-8"></head>
<body style="font-family: 'Segoe UI', Arial, sans-serif; color: #333; max-width: 1000px; margin: 0 auto; padding: 20px;">

<!-- Header -->
<table width="100%" cellpadding="0" cellspacing="0" style="margin-bottom: 24px;">
<tr>
<td bgcolor="#004C91" style="background-color:#004C91; padding: 24px 30px;">
    <div style="color: white; font-size: 26px; font-weight: 700;">V.E.T. Executive Report</div>
    <div style="color: #cccccc; font-size: 13px; margin-top: 2px;">Weekly Report &mdash; {today}</div>
</td>
</tr>
</table>

<!-- Dashboard Button -->
<table width="100%" cellpadding="0" cellspacing="0" style="margin-bottom: 24px;">
<tr>
<td align="center" style="padding: 16px 0;">
    <table cellpadding="0" cellspacing="0" style="border-collapse: separate;">
    <tr>
    <td align="center" bgcolor="#0052CC" style="background-color: #0052CC; border-radius: 30px; mso-padding-alt: 14px 40px;">
        <a href="http://WEUS42608431466:5001/VET_Executive_Report" target="_blank" style="background-color: #0052CC; color: #FFFFFF; font-family: 'Segoe UI', Arial, sans-serif; font-size: 16px; font-weight: 700; text-decoration: none; padding: 14px 40px; border-radius: 30px; display: inline-block; letter-spacing: 0.5px; mso-padding-alt: 0; border: 2px solid #0052CC;">Go to Dashboard</a>
    </td>
    </tr>
    </table>
</td>
</tr>
</table>

<!-- Context -->
<table width="100%" cellpadding="0" cellspacing="0" style="margin-bottom: 24px;"><tr>
<td bgcolor="#f8f9fa" style="background-color:#f8f9fa; border: 1px solid #e5e5e5; padding: 16px 20px; font-size: 13px; line-height: 1.6; color: #444;">
    <strong style="color: {COLORS['walmart_blue_dark']};">About This Report</strong><br>
    This is your weekly V.E.T. (Validation, Engineering &amp; Testing) Executive Report. It provides a snapshot of all active Dallas VET projects tracked through the Intake Hub. Projects are grouped by phase. A PowerPoint version is attached for offline review and sharing.
</td></tr></table>

<!-- Summary Cards -->
<table width="100%" cellpadding="0" cellspacing="0" style="margin-bottom: 24px;">
<tr>
    <td width="20%" style="padding: 0 4px;">
        <div style="background:white; border:1px solid #e5e5e5; border-left:4px solid {COLORS['walmart_yellow']}; border-radius:6px; padding:12px; text-align:center;">
            <div style="font-size:11px; color:#666; text-transform:uppercase; letter-spacing:0.5px;">Total Initiatives</div>
            <div style="font-size:24px; font-weight:700; color:{COLORS['walmart_blue_dark']};">{total}</div>
        </div>
    </td>
    <td width="20%" style="padding: 0 4px;">
        <div style="background:white; border:1px solid #e5e5e5; border-left:4px solid {COLORS['walmart_yellow']}; border-radius:6px; padding:12px; text-align:center;">
            <div style="font-size:11px; color:#666; text-transform:uppercase; letter-spacing:0.5px;">Total Stores</div>
            <div style="font-size:24px; font-weight:700; color:{COLORS['walmart_blue_dark']};">{total_stores:,}</div>
        </div>
    </td>
    <td width="20%" style="padding: 0 4px;">
        <div style="background:white; border:1px solid #e5e5e5; border-left:4px solid {COLORS['on_track']}; border-radius:6px; padding:12px; text-align:center;">
            <div style="font-size:11px; color:#666; text-transform:uppercase; letter-spacing:0.5px;">On Track</div>
            <div style="font-size:24px; font-weight:700; color:{COLORS['on_track']};">{on_track}</div>
        </div>
    </td>
    <td width="20%" style="padding: 0 4px;">
        <div style="background:white; border:1px solid #e5e5e5; border-left:4px solid {COLORS['at_risk']}; border-radius:6px; padding:12px; text-align:center;">
            <div style="font-size:11px; color:#666; text-transform:uppercase; letter-spacing:0.5px;">At Risk</div>
            <div style="font-size:24px; font-weight:700; color:{COLORS['at_risk']};">{at_risk}</div>
        </div>
    </td>
    <td width="20%" style="padding: 0 4px;">
        <div style="background:white; border:1px solid #e5e5e5; border-left:4px solid {COLORS['off_track']}; border-radius:6px; padding:12px; text-align:center;">
            <div style="font-size:11px; color:#666; text-transform:uppercase; letter-spacing:0.5px;">Off Track</div>
            <div style="font-size:24px; font-weight:700; color:{COLORS['off_track']};">{off_track}</div>
        </div>
    </td>
</tr>
</table>
"""

    # Phase tables — grouped by phase
    for phase in PHASE_ORDER:
        phase_data = [r for r in data if r.get('Phase') == phase]
        if not phase_data:
            continue
        count = len(phase_data)

        html += f"""
<div style="background:#1E3A8A; color:white; padding:10px 16px; font-size:15px; font-weight:700; border-radius:6px 6px 0 0; margin-top:16px;">
    Phase: {escape(phase)} &mdash; {count} Project{'s' if count != 1 else ''}
</div>
<table width="100%" cellpadding="0" cellspacing="0" style="border:1px solid #e5e5e5; border-top:none; border-collapse:collapse; margin-bottom:8px; font-size:14px;">
<thead>
<tr style="background:#f5f5f5;">
    <th style="padding:10px 12px; text-align:left; border-bottom:2px solid #ddd; font-weight:700; color:#333; font-size:14px;">Initiative - Project Title</th>
    <th style="padding:10px 12px; text-align:left; border-bottom:2px solid #ddd; font-weight:700; color:#333; font-size:14px; width:100px;">Health</th>
    <th style="padding:10px 12px; text-align:center; border-bottom:2px solid #ddd; font-weight:700; color:#333; font-size:14px; width:80px;"># of Stores</th>
</tr>
</thead>
<tbody>
"""
        for i, row in enumerate(phase_data):
            bg = '#ffffff' if i % 2 == 0 else '#fafafa'
            title = escape(row.get('Initiative - Project Title', ''))
            health = escape(row.get('Health Status', ''))
            stores = int(row.get('# of Stores', 0) or 0)
            bs = badge_style(row.get('Health Status', ''))

            html += f"""<tr style="background:{bg};">
    <td style="padding:10px 12px; border-bottom:1px solid #eee; color:{COLORS['walmart_blue']}; font-weight:500; font-size:14px;">{title}</td>
    <td style="padding:10px 12px; border-bottom:1px solid #eee;"><span style="{bs}">{health}</span></td>
    <td style="padding:10px 12px; border-bottom:1px solid #eee; text-align:center; font-weight:600; color:{COLORS['walmart_blue']}; font-size:14px;">{stores:,}</td>
</tr>
"""
        html += "</tbody></table>\n"

    # Footer
    html += f"""
<div style="margin-top:24px; padding:16px; border-top:2px solid #e5e5e5; text-align:center; color:#999; font-size:11px;">
    <p>This report was automatically generated from the V.E.T. Executive Dashboard.<br>
    PPT report is attached for offline review.</p>
    <p style="margin-top:8px;">Generated on {today} at {datetime.now().strftime('%I:%M %p')} CT</p>
</div>
</body></html>"""
    return html


# ── PPT HTML builder (mirrors TDA _build_phase_html) ──────────────────────────
def _build_summary_html(data):
    """Build the summary+needs-attention slide HTML matching dashboard layout."""
    total = len(data)
    total_stores = sum(int(r.get('# of Stores', 0) or 0) for r in data)
    health_counts = {}
    for r in data:
        h = r.get('Health Status', 'Unknown')
        health_counts[h] = health_counts.get(h, 0) + 1
    on_track = health_counts.get('On Track', 0)
    at_risk_count = health_counts.get('At Risk', 0)
    off_track = health_counts.get('Off Track', 0)
    at_risk_items = [r for r in data if str(r.get('Health Status', '')).lower() == 'at risk']

    html = f'''<!DOCTYPE html><html><head><meta charset="utf-8">
<style>*{{margin:0;padding:0;box-sizing:border-box;}} body{{background:white;font-family:-apple-system,BlinkMacSystemFont,'Segoe UI',Arial,sans-serif;padding:30px;}}</style>
</head><body>
<div style="background:#004C91;color:white;padding:24px 30px;border-radius:8px;margin-bottom:24px;">
<div style="font-size:28px;font-weight:700;">V.E.T. Executive Report</div>
<div style="font-size:14px;color:#ccc;margin-top:4px;">Walmart Enterprise Transformation Dashboard</div>
</div>
<div style="display:flex;gap:16px;margin-bottom:24px;">
<div style="flex:1;background:white;border:1px solid #e5e5e5;border-left:4px solid #FFC220;border-radius:6px;padding:16px;text-align:center;">
<div style="font-size:12px;color:#666;text-transform:uppercase;">Total Projects</div>
<div style="font-size:32px;font-weight:700;color:#004C91;">{total}</div>
</div>
<div style="flex:1;background:white;border:1px solid #e5e5e5;border-left:4px solid #FFC220;border-radius:6px;padding:16px;text-align:center;">
<div style="font-size:12px;color:#666;text-transform:uppercase;">Total Stores</div>
<div style="font-size:32px;font-weight:700;color:#004C91;">{total_stores:,}</div>
</div>
<div style="flex:1;background:white;border:1px solid #e5e5e5;border-left:4px solid #107C10;border-radius:6px;padding:16px;text-align:center;">
<div style="font-size:12px;color:#666;text-transform:uppercase;">On Track</div>
<div style="font-size:32px;font-weight:700;color:#107C10;">{on_track}</div>
</div>
<div style="flex:1;background:white;border:1px solid #e5e5e5;border-left:4px solid #F7630C;border-radius:6px;padding:16px;text-align:center;">
<div style="font-size:12px;color:#666;text-transform:uppercase;">At Risk</div>
<div style="font-size:32px;font-weight:700;color:#F7630C;">{at_risk_count}</div>
</div>
<div style="flex:1;background:white;border:1px solid #e5e5e5;border-left:4px solid #DC3545;border-radius:6px;padding:16px;text-align:center;">
<div style="font-size:12px;color:#666;text-transform:uppercase;">Off Track</div>
<div style="font-size:32px;font-weight:700;color:#DC3545;">{off_track}</div>
</div>
</div>'''

    # Needs Attention section
    if at_risk_items:
        html += '<div style="margin-bottom:20px;">'
        html += '<div style="font-size:20px;font-weight:700;color:#DC3545;margin-bottom:12px;">&#128680; Needs Attention</div>'
        html += '<div style="display:grid;grid-template-columns:repeat(auto-fill,minmax(280px,1fr));gap:12px;">'
        for item in at_risk_items:
            title = escape(item.get('Initiative - Project Title', ''))
            notes = escape(str(item.get('Executive Notes', '') or ''))
            phase = escape(str(item.get('Phase', '')))
            stores = int(item.get('# of Stores', 0) or 0)
            html += f'''<div style="background:#FFF3CD;border:1px solid #F7630C;border-radius:8px;padding:14px;">
<div style="display:inline-block;background:#DC3545;color:white;padding:2px 8px;border-radius:4px;font-size:11px;font-weight:700;margin-bottom:6px;">At Risk</div>
<div style="font-size:14px;font-weight:600;color:#212121;margin-bottom:4px;">{title}</div>
<div style="font-size:12px;color:#666;">Phase: {phase} | Stores: {stores:,}</div>
<div style="font-size:12px;color:#333;margin-top:4px;">{notes}</div>
</div>'''
        html += '</div></div>'

    html += '</body></html>'
    return html


def _build_phase_html(phase, rows):
    """Build HTML table for a phase — matches dashboard's PPT table columns exactly."""
    columns = ['Initiative - Project Title', 'Health Status', 'Phase',
                'WM Week', '# of Stores', 'Executive Notes']

    banner = f'<div style="background:#3B82F6;color:white;padding:10px 20px;font-size:16px;font-weight:600;text-align:center;border-radius:4px 4px 0 0;">Phase: {escape(phase)}</div>'

    table = '<table style="width:100%;border-collapse:collapse;font-size:14px;font-family:-apple-system,BlinkMacSystemFont,\'Segoe UI\',\'Helvetica Neue\',Arial,sans-serif;">'
    table += '<thead style="background-color:#F5F5F5;border-bottom:2px solid #E5E5E5;"><tr>'
    for col in columns:
        table += f'<th style="padding:16px 24px;text-align:left;font-weight:600;color:#212121;white-space:nowrap;">{escape(col)}</th>'
    table += '</tr></thead><tbody>'

    for row in rows:
        table += '<tr>'
        for col in columns:
            value = str(row.get(col, '') or '')
            if col == '# of Stores':
                try:
                    value = f"{int(row.get(col, 0) or 0):,}"
                except (ValueError, TypeError):
                    pass
            cell_style = 'padding:16px 24px;border-bottom:1px solid #E5E5E5;text-align:left;'

            if col == 'Health Status':
                status = value.lower()
                bg, clr = '#f0f0f0', '#333'
                if 'on track' in status:
                    bg, clr = 'rgba(16,124,16,0.2)', '#107C10'
                elif 'at risk' in status:
                    bg, clr = 'rgba(247,99,12,0.2)', '#F7630C'
                elif 'off track' in status:
                    bg, clr = 'rgba(220,53,69,0.2)', '#DC3545'
                table += f'<td style="{cell_style}"><span style="display:inline-block;padding:4px 10px;background-color:{bg};color:{clr};border-radius:12px;font-weight:600;font-size:12px;white-space:nowrap;">{escape(value)}</span></td>'
            elif col == 'Phase':
                table += f'<td style="{cell_style}"><span style="display:inline-block;padding:4px 10px;background-color:#DBEAFE;color:#1E3A8A;border-radius:4px;font-weight:500;font-size:12px;">{escape(value)}</span></td>'
            elif col == '# of Stores':
                table += f'<td style="{cell_style}font-weight:600;color:#0071CE;">{escape(value)}</td>'
            else:
                table += f'<td style="{cell_style}">{escape(value)}</td>'
        table += '</tr>'
    table += '</tbody></table>'

    return f'''<!DOCTYPE html><html><head><meta charset="utf-8">
<style>*{{margin:0;padding:0;box-sizing:border-box;}} body{{background:white;}}</style>
</head><body>{banner}{table}</body></html>'''


# ── Edge headless screenshot (same as TDA) ─────────────────────────────────────
def _capture_html_screenshot(html_content, output_png, width=1280):
    """Use Edge headless to capture HTML as a PNG screenshot."""
    with tempfile.NamedTemporaryFile(suffix='.html', delete=False, mode='w', encoding='utf-8') as f:
        f.write(html_content)
        html_path = f.name
    try:
        cmd = [
            EDGE_PATH,
            '--headless',
            '--disable-gpu',
            '--no-sandbox',
            '--hide-scrollbars',
            '--force-device-scale-factor=1.5',
            f'--screenshot={output_png}',
            f'--window-size={width},2000',
            f'file:///{html_path.replace(os.sep, "/")}'
        ]
        subprocess.run(cmd, capture_output=True, timeout=30)
        if not Path(output_png).exists():
            raise RuntimeError("Edge screenshot failed")
        # Auto-crop whitespace
        img = Image.open(output_png)
        bg = Image.new(img.mode, img.size, (255, 255, 255))
        diff = ImageChops.difference(img, bg)
        bbox = diff.getbbox()
        if bbox:
            cropped = img.crop((0, 0, bbox[2], bbox[3] + 4))
            jpg_path = output_png.replace('.png', '.jpg')
            cropped.convert('RGB').save(jpg_path, 'JPEG', quality=78, optimize=True)
            import shutil
            shutil.move(jpg_path, output_png)
    finally:
        os.unlink(html_path)


# ── Row-height measurement & pagination (same as TDA) ─────────────────────────
def _measure_row_heights(rows):
    """Use Edge --dump-dom to measure rendered row heights."""
    columns = ['Initiative - Project Title', 'Health Status', 'Phase',
                'WM Week', '# of Stores', 'Executive Notes']
    table = '<table style="width:100%;border-collapse:collapse;font-size:14px;font-family:Segoe UI,Arial,sans-serif;line-height:1.5;">'
    table += '<thead><tr style="background-color:#f0f0f0;border-bottom:2px solid #999;height:50px;">'
    for col in columns:
        table += f'<th style="padding:8px;text-align:left;border:1px solid #ddd;">{escape(col)}</th>'
    table += '</tr></thead><tbody>'
    for row in rows:
        table += '<tr style="border:1px solid #ddd;">'
        for col in columns:
            value = str(row.get(col, '') or '')
            if col == '# of Stores':
                try:
                    value = str(int(row.get(col, 0) or 0))
                except (ValueError, TypeError):
                    pass
            table += f'<td style="padding:8px;border:1px solid #ddd;word-wrap:break-word;white-space:normal;">{escape(value)}</td>'
        table += '</tr>'
    table += '</tbody></table>'

    meas_html = f'''<!DOCTYPE html><html><head><meta charset="utf-8">
<style>*{{margin:0;padding:0;box-sizing:border-box;}}</style>
</head><body style="width:1280px;">
{table}
<script>
window.onload=function(){{
  var rs=document.querySelectorAll("tbody tr");
  var h=[];rs.forEach(function(r){{h.push(r.offsetHeight);}});
  document.getElementById("_rh").textContent="ROW_HEIGHTS:"+JSON.stringify(h);
}};
</script><pre id="_rh"></pre>
</body></html>'''

    with tempfile.NamedTemporaryFile(suffix='.html', delete=False, mode='w', encoding='utf-8') as f:
        f.write(meas_html)
        html_path = f.name
    try:
        cmd = [
            EDGE_PATH, '--headless', '--disable-gpu', '--no-sandbox',
            '--dump-dom', '--virtual-time-budget=5000',
            '--window-size=1280,4000',
            f'file:///{html_path.replace(os.sep, "/")}'
        ]
        result = subprocess.run(cmd, capture_output=True, timeout=30)
        stdout = result.stdout.decode('utf-8', errors='replace')
        match = re.search(r'ROW_HEIGHTS:\[([\d,]+)\]', stdout)
        if match:
            return [int(x) for x in match.group(1).split(',')]
    finally:
        os.unlink(html_path)
    return [55] * len(rows)


def _paginate_by_height(rows, heights):
    """Split rows into pages when cumulative height exceeds slide space."""
    AVAILABLE_HEIGHT = 580  # matches dashboard: 720 - 60 title - 50 header - 30 padding
    pages = []
    current_page = []
    current_h = 0
    for row, h in zip(rows, heights):
        if current_h + h > AVAILABLE_HEIGHT and current_page:
            pages.append(current_page)
            current_page = [row]
            current_h = h
        else:
            current_page.append(row)
            current_h += h
    if current_page:
        pages.append(current_page)
    return pages


# ── PPT generation (mirrors TDA generate_report_pptx) ─────────────────────────
def generate_report_pptx(data):
    """Generate PPTX using Edge headless screenshots — identical to dashboard Generate PPT."""
    print("    Using Edge headless for screenshot-based PPT (matches dashboard)...")

    # Group data by phase
    screenshots = []  # list of (label, png_bytes)

    with tempfile.TemporaryDirectory() as tmp_dir:
        # Step 1: Summary slide (header + stats + needs attention)
        print("    Capturing executive summary slide...")
        summary_html = _build_summary_html(data)
        summary_png = os.path.join(tmp_dir, 'slide_summary.png')
        _capture_html_screenshot(summary_html, summary_png)
        if Path(summary_png).exists():
            png_bytes = Path(summary_png).read_bytes()
            screenshots.append(('Executive Summary', png_bytes))
            print(f"    Captured: Executive Summary ({len(png_bytes):,} bytes)")

        # Step 2: Phase table slides
        for phase in PHASE_ORDER:
            phase_rows = [r for r in data if r.get('Phase') == phase]
            if not phase_rows:
                continue

            heights = _measure_row_heights(phase_rows)
            pages = _paginate_by_height(phase_rows, heights)
            total_pages = len(pages)

            for page_idx, page_rows in enumerate(pages):
                label = f"Phase: {phase}"
                if total_pages > 1:
                    label += f" ({page_idx + 1}/{total_pages})"

                html = _build_phase_html(phase, page_rows)
                png_path = os.path.join(tmp_dir, f'slide_{len(screenshots) + 1}.png')
                _capture_html_screenshot(html, png_path)
                png_bytes = Path(png_path).read_bytes()
                screenshots.append((label, png_bytes))
                print(f"    Captured: {label} ({len(png_bytes):,} bytes)")

    # Build PPTX (same dimensions as dashboard: 9.6 x 7.2)
    SLIDE_WIDTH = Inches(9.6)
    SLIDE_HEIGHT = Inches(7.2)

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Screenshot slides (no separate title slide — summary IS the first slide, like dashboard)
    for label, png_bytes in screenshots:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        img = Image.open(io.BytesIO(png_bytes))
        img_w, img_h = img.size
        aspect = img_h / img_w if img_w > 0 else 0.75

        max_w = 9.2  # inches, with margins
        max_h = 6.8
        if aspect > (max_h / max_w):
            slide.shapes.add_picture(io.BytesIO(png_bytes), Inches(0.2), Inches(0.2), height=Inches(max_h))
        else:
            slide.shapes.add_picture(io.BytesIO(png_bytes), Inches(0.2), Inches(0.2), width=Inches(max_w))

    pptx_buffer = io.BytesIO()
    prs.save(pptx_buffer)
    pptx_buffer.seek(0)
    return pptx_buffer.getvalue(), screenshots


# ── PDF generation (mirrors TDA generate_report_pdf) ──────────────────────────
def generate_report_pdf(pptx_screenshots):
    """Convert screenshot images to a multi-page PDF using Pillow."""
    if not pptx_screenshots:
        return None
    pages = []
    for label, png_bytes in pptx_screenshots:
        img = Image.open(io.BytesIO(png_bytes)).convert('RGB')
        pages.append(img)
    pdf_buffer = io.BytesIO()
    pages[0].save(pdf_buffer, 'PDF', save_all=True, append_images=pages[1:], resolution=150)
    pdf_buffer.seek(0)
    return pdf_buffer.getvalue()


# ── Outlook send (mirrors TDA send_outlook_email) ─────────────────────────────
def send_outlook_email(html_body, attachments, recipients, test_mode=False):
    """Send email via Outlook with attachments."""
    import win32com.client
    outlook = win32com.client.Dispatch('Outlook.Application')
    mail = outlook.CreateItem(0)
    mail.Subject = f"V.E.T. Executive Report - WK{WEEK_NUM:02d}"
    mail.HTMLBody = html_body
    mail.To = '; '.join(recipients)
    for att in attachments:
        mail.Attachments.Add(str(att))
    if test_mode:
        mail.Save()
        print(f"[OK] Email saved as DRAFT")
    else:
        mail.Send()
        print(f"[OK] Email sent to: {', '.join(recipients)}")


# ── Main ───────────────────────────────────────────────────────────────────────
def main():
    preview_only = '--preview' in sys.argv
    test_mode = '--draft' in sys.argv

    # Parse recipient
    recipients = list(RECIPIENTS)
    if '--email' in sys.argv:
        idx = sys.argv.index('--email')
        if idx + 1 < len(sys.argv):
            recipients = [sys.argv[idx + 1]]

    print("=" * 60)
    print("V.E.T. Executive Report - Weekly Email")
    print("=" * 60)

    # 1. Fetch data
    print("\n[1/5] Fetching data...")
    data = fetch_data()

    # 2. Generate PPT + PDF
    print("\n[2/5] Generating PPT report...")
    pptx_data, screenshots = generate_report_pptx(data)
    PPT_OUTPUT.write_bytes(pptx_data)
    print(f"  Saved: {PPT_OUTPUT.name}")
    print(f"  Size: {len(pptx_data):,} bytes")

    print("\n[3/5] Generating PDF report...")
    pdf_data = generate_report_pdf(screenshots)
    if pdf_data:
        PDF_OUTPUT.write_bytes(pdf_data)
        print(f"  Saved: {PDF_OUTPUT.name}")
        print(f"  Size: {len(pdf_data):,} bytes")

    # 3. Build email HTML
    print("\n[4/5] Building email HTML...")
    html = build_email_html(data)

    if preview_only:
        preview_path = SCRIPT_DIR / "email_preview.html"
        preview_path.write_text(html, encoding='utf-8')
        print(f"\n  PREVIEW MODE - files saved:")
        print(f"  PPT:   {PPT_OUTPUT}")
        print(f"  PDF:   {PDF_OUTPUT}")
        print(f"  HTML:  {preview_path}")
        return

    # 4. Send email with PPT + PDF attachments
    print(f"\n[5/5] Sending email to {', '.join(recipients)}...")
    attachments = [PPT_OUTPUT]
    if PDF_OUTPUT.exists():
        attachments.append(PDF_OUTPUT)
    send_outlook_email(html, attachments, recipients, test_mode=test_mode)

    print("\n" + "=" * 60)
    print("Done!")
    print("=" * 60)


if __name__ == '__main__':
    main()

