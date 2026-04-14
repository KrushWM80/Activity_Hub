#!/usr/bin/env python3
"""
Impact Platform Report Generator
Generates PPT + PDF reports from project data
Uses python-pptx and PIL (Pillow) following TDA Insights pattern

Generates:
- Title slide (Impact Platform Report)
- Metrics slide (Project Count, % Updated, Unique Owners)
- Project list slides (sorted by Business Area)
"""

import io
import os
import uuid
from pathlib import Path
from datetime import datetime
from typing import List, Tuple

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image, ImageDraw, ImageFont
import tempfile

# Configuration
TEMP_DIR = Path(tempfile.gettempdir()) / "impact_reports"
TEMP_DIR.mkdir(exist_ok=True)

# Walmart brand colors
WALMART_BLUE = RGBColor(0, 113, 206)  # #0071CE
WALMART_YELLOW = RGBColor(255, 194, 32)  # #FFC220
GREEN = RGBColor(16, 124, 16)
YELLOW = RGBColor(247, 99, 12)
RED = RGBColor(220, 53, 69)

# Slide dimensions (in EMUs)
SLIDE_WIDTH = Inches(10)
SLIDE_HEIGHT = Inches(7.5)

def get_health_color(health_status: str) -> RGBColor:
    """Map health status to RGB color"""
    status_lower = (health_status or "").lower()
    if "green" in status_lower:
        return GREEN
    elif "yellow" in status_lower or "risk" in status_lower:
        return YELLOW
    elif "red" in status_lower:
        return RED
    return RGBColor(128, 128, 128)  # Gray default

def create_title_slide(prs: Presentation) -> None:
    """Create title slide with Impact Platform branding"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Set background to Walmart blue
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WALMART_BLUE
    
    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5), Inches(9), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = "Impact Platform"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(72)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(4.2), Inches(9), Inches(1)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Projects Impact Dashboard Report"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(32)
    subtitle_para.font.color.rgb = WALMART_YELLOW
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Date
    date_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(6), Inches(9), Inches(0.8)
    )
    date_frame = date_box.text_frame
    date_frame.text = datetime.now().strftime("%B %d, %Y")
    date_para = date_frame.paragraphs[0]
    date_para.font.size = Pt(18)
    date_para.font.color.rgb = RGBColor(255, 255, 255)
    date_para.alignment = PP_ALIGN.CENTER

def create_metrics_slide(prs: Presentation, metrics: dict) -> None:
    """Create metrics overview slide with 2x2 grid layout"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # White background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_frame.text = "Dashboard Metrics"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = WALMART_BLUE
    
    # Metrics boxes - 2x2 grid centered
    metrics_data = [
        ("Active Projects", str(metrics.get('active_projects', 0)), WALMART_BLUE),
        ("Percent Updated", f"{metrics.get('percent_updated', 0)}%", WALMART_YELLOW),
        ("Updated This Week", str(metrics.get('projects_updated_this_week', 0)), GREEN),
        ("Unique Owners", str(metrics.get('unique_owners', 0)), RGBColor(0, 102, 204)),
    ]
    
    # Adjust positioning for better centering
    box_width = 3.8
    box_height = 2
    gap = 0.4
    total_width = (box_width * 2) + gap
    start_x = (10 - total_width) / 2  # Center horizontally
    start_y = 1.8
    
    for idx, (label, value, color) in enumerate(metrics_data):
        row = idx // 2
        col = idx % 2
        left = start_x + col * (box_width + gap)
        top = start_y + row * (box_height + gap)
        
        # Box background
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left), Inches(top), Inches(box_width), Inches(box_height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = color
        shape.line.width = 0
        
        # Value text (large)
        text_frame = shape.text_frame
        text_frame.clear()
        text_frame.margin_bottom = Inches(0)
        text_frame.margin_top = Inches(0.2)
        p = text_frame.paragraphs[0]
        p.text = value
        p.font.size = Pt(52)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        # Label text below box
        label_box = slide.shapes.add_textbox(
            Inches(left), Inches(top + box_height + 0.15), Inches(box_width), Inches(0.5)
        )
        label_frame = label_box.text_frame
        label_frame.text = label
        label_para = label_frame.paragraphs[0]
        label_para.font.size = Pt(16)
        label_para.font.bold = True
        label_para.font.color.rgb = RGBColor(50, 50, 50)
        label_para.alignment = PP_ALIGN.CENTER

def create_projects_slide(prs: Presentation, projects: List[dict], slide_number: int) -> int:
    """
    Create project list slides with table layout (max 8 projects per slide)
    Returns number of slides created
    """
    projects_per_slide = 8
    slides_created = 0
    
    for i in range(0, len(projects), projects_per_slide):
        slide_batch = projects[i:i + projects_per_slide]
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)
        
        # White background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)
        
        # Title with page number
        title_box = slide.shapes.add_textbox(
            Inches(0.4), Inches(0.2), Inches(9.2), Inches(0.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = f"Active Projects (Page {slide_number + 1})"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = WALMART_BLUE
        
        # Add horizontal line below title
        line_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.4), Inches(0.75), Inches(9.2), Inches(0.02)
        )
        line_shape.fill.solid()
        line_shape.fill.fore_color.rgb = WALMART_BLUE
        line_shape.line.width = 0
        
        # Table with projects
        rows = len(slide_batch) + 1  # +1 for header
        cols = 4
        left = Inches(0.3)
        top = Inches(1)
        width = Inches(9.4)
        height = Inches(6)
        
        table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
        table = table_shape.table
        
        # Set column widths
        table.columns[0].width = Inches(3.5)  # Title
        table.columns[1].width = Inches(2)    # Health
        table.columns[2].width = Inches(2)    # Owner
        table.columns[3].width = Inches(1.9)  # Latest Update
        
        # Header row styling
        headers = ["Project Title", "Health", "Owner", "Latest Update"]
        for col_idx, header in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WALMART_BLUE
            cell.text_frame.clear()
            text_frame = cell.text_frame
            text_frame.text = header
            text_frame.margin_bottom = Inches(0.1)
            text_frame.margin_top = Inches(0.1)
            text_para = text_frame.paragraphs[0]
            text_para.font.size = Pt(12)
            text_para.font.bold = True
            text_para.font.color.rgb = RGBColor(255, 255, 255)
        
        # Data rows with alternating background
        for row_idx, project in enumerate(slide_batch, start=1):
            # Alternating row colors for readability
            row_color = RGBColor(245, 247, 250) if row_idx % 2 == 0 else RGBColor(255, 255, 255)
            
            # Title
            cell = table.cell(row_idx, 0)
            cell.fill.solid()
            cell.fill.fore_color.rgb = row_color
            cell.text_frame.clear()
            cell.text_frame.text = project.get('title', '')[:50]
            cell.text_frame.paragraphs[0].font.size = Pt(11)
            cell.text_frame.margin_bottom = Inches(0.05)
            cell.text_frame.margin_top = Inches(0.05)
            
            # Health (colored background)
            cell = table.cell(row_idx, 1)
            health = project.get('health_status', 'Gray')
            cell.fill.solid()
            cell.fill.fore_color.rgb = get_health_color(health)
            cell.text_frame.clear()
            cell.text_frame.text = health
            cell.text_frame.paragraphs[0].font.size = Pt(10)
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            cell.text_frame.margin_bottom = Inches(0.05)
            cell.text_frame.margin_top = Inches(0.05)
            
            # Owner
            cell = table.cell(row_idx, 2)
            cell.fill.solid()
            cell.fill.fore_color.rgb = row_color
            cell.text_frame.clear()
            cell.text_frame.text = project.get('owner_name', '')[:20]
            cell.text_frame.paragraphs[0].font.size = Pt(10)
            cell.text_frame.margin_bottom = Inches(0.05)
            cell.text_frame.margin_top = Inches(0.05)
            
            # Latest update
            cell = table.cell(row_idx, 3)
            cell.fill.solid()
            cell.fill.fore_color.rgb = row_color
            cell.text_frame.clear()
            update_text = "Updated" if project.get('latest_update') else "Pending"
            cell.text_frame.text = update_text
            cell.text_frame.paragraphs[0].font.size = Pt(9)
            cell.text_frame.margin_bottom = Inches(0.05)
            cell.text_frame.margin_top = Inches(0.05)
        
        slides_created += 1
        slide_number += 1
    
    return slides_created

def generate_pptx_report(projects: List[dict]) -> Tuple[bytes, str]:
    """
    Generate complete PPT report from projects list
    Returns: (pptx_bytes, report_id)
    """
    # Calculate metrics
    active = [p for p in projects if p.get('project_status') == 'Active']
    updated_this_week = [p for p in active if p.get('current_wm_week_update')]
    unique_owners = set(p.get('owner_id') for p in active)
    
    metrics = {
        'active_projects': len(active),
        'projects_updated_this_week': len(updated_this_week),
        'percent_updated': round(len(updated_this_week) / max(len(active), 1) * 100, 1),
        'unique_owners': len(unique_owners)
    }
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    
    # Add slides
    create_title_slide(prs)
    create_metrics_slide(prs, metrics)
    create_projects_slide(prs, sorted(active, key=lambda p: (p.get('business_area', ''), p.get('title', ''))), 0)
    
    # Save to bytes
    pptx_buffer = io.BytesIO()
    prs.save(pptx_buffer)
    pptx_buffer.seek(0)
    pptx_data = pptx_buffer.getvalue()
    
    # Save to temp file
    report_id = str(uuid.uuid4())
    pptx_path = TEMP_DIR / f"{report_id}.pptx"
    pptx_path.write_bytes(pptx_data)
    
    return pptx_data, report_id

def generate_pdf_report(pptx_bytes: bytes, report_id: str) -> bytes:
    """
    Convert PPT to PDF (simplified - uses PIL to create multi-page image-based PDF)
    Note: For production, consider using LibreOffice or python-pptx with proper conversion
    """
    try:
        # For now, create a placeholder PDF
        # In production, use proper PPT->PDF conversion
        from PyPDF2 import PdfWriter
    except ImportError:
        print("Note: PyPDF2 not installed. PDF generation skipped.")
        return None
    
    return pptx_bytes  # Placeholder

if __name__ == "__main__":
    # Test
    test_projects = [
        {
            'impact_id': '1',
            'title': 'Test Project 1',
            'owner_name': 'John Doe',
            'owner_id': 'jdoe',
            'business_area': 'Operations',
            'health_status': 'Green',
            'project_status': 'Active',
            'current_wm_week_update': 'On track',
            'current_wm_week_update_timestamp': datetime.now().isoformat()
        }
    ]
    
    pptx_data, report_id = generate_pptx_report(test_projects)
    print(f"Generated PPT: {report_id}")
    print(f"File size: {len(pptx_data)} bytes")
