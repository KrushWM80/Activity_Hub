"""
Activity Hub Consolidated Presentation Generator

This script generates a 4-slide PowerPoint presentation focused on:
1. Title Slide
2. The Challenges - Pain Points & Fragmentation
3. The Activity Hub - Production Path & Future Solutions
4. Path to Activity Hub Playground - Solution, Investment, ROI

Requirements:
    pip install python-pptx

Usage:
    python activity_hub_consolidated.py

Output:
    Activity_Hub_Consolidated_Proposal.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Walmart brand colors
WALMART_BLUE = RGBColor(0, 72, 144)
WALMART_YELLOW = RGBColor(252, 200, 16)
WHITE = RGBColor(255, 255, 255)
DARK_GRAY = RGBColor(50, 50, 50)
RED = RGBColor(200, 0, 0)
LIGHT_GRAY = RGBColor(240, 240, 240)


def add_title_slide(prs, title, subtitle, bottom_text):
    """Create title slide with Walmart branding."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WALMART_BLUE
    shapes = slide.shapes
    
    # Title
    title_box = shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.2))
    title_tf = title_box.text_frame
    title_tf.text = title
    title_tf.paragraphs[0].font.size = Pt(66)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = WALMART_YELLOW
    title_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(1))
    subtitle_tf = subtitle_box.text_frame
    subtitle_tf.text = subtitle
    subtitle_tf.paragraphs[0].font.size = Pt(40)
    subtitle_tf.paragraphs[0].font.color.rgb = WHITE
    subtitle_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Bottom text
    bottom_box = shapes.add_textbox(Inches(1), Inches(6.8), Inches(8), Inches(0.5))
    bottom_tf = bottom_box.text_frame
    bottom_tf.text = bottom_text
    bottom_tf.paragraphs[0].font.size = Pt(18)
    bottom_tf.paragraphs[0].font.color.rgb = WHITE
    bottom_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Walmart spark logo placeholder (top right)
    logo_box = shapes.add_textbox(Inches(8), Inches(0.3), Inches(1.5), Inches(1))
    logo_tf = logo_box.text_frame
    logo_tf.text = "[Walmart Spark Logo]"
    logo_tf.paragraphs[0].font.size = Pt(14)
    logo_tf.paragraphs[0].font.color.rgb = WALMART_YELLOW
    logo_tf.paragraphs[0].alignment = PP_ALIGN.RIGHT


def add_challenges_slide(prs):
    """Slide 2: The Challenges - Pain Points & Fragmentation"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shapes = slide.shapes
    
    # Title
    title_box = shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
    title_tf = title_box.text_frame
    title_tf.text = "The Challenge: Fragmented Operations"
    title_tf.paragraphs[0].font.size = Pt(48)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = WALMART_BLUE
    
    # Left column - Fragmentation
    left_title_box = shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4.5), Inches(0.5))
    left_title_tf = left_title_box.text_frame
    left_title_tf.text = "No Centralized Location"
    left_title_tf.paragraphs[0].font.size = Pt(24)
    left_title_tf.paragraphs[0].font.bold = True
    left_title_tf.paragraphs[0].font.color.rgb = RED
    
    left_bullets = [
        "16 disconnected platforms",
        "Multiple communication channels",
        "No single source of truth",
        "Information silos across teams",
        "Duplicate work and effort"
    ]
    
    left_box = shapes.add_textbox(Inches(0.5), Inches(1.9), Inches(4.5), Inches(4.5))
    left_tf = left_box.text_frame
    for i, bullet in enumerate(left_bullets):
        if i == 0:
            p = left_tf.paragraphs[0]
        else:
            p = left_tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_GRAY
    
    # Right column - Pain Points
    right_title_box = shapes.add_textbox(Inches(5.2), Inches(1.3), Inches(4.3), Inches(0.5))
    right_title_tf = right_title_box.text_frame
    right_title_tf.text = "Pain Points"
    right_title_tf.paragraphs[0].font.size = Pt(24)
    right_title_tf.paragraphs[0].font.bold = True
    right_title_tf.paragraphs[0].font.color.rgb = WALMART_BLUE
    
    right_bullets = [
        "120+ hours/week manual data work",
        "30-40% communication overlap",
        "Store confusion & initiative fatigue",
        "No way to measure cumulative impact",
        "Slow decision-making from lack of data"
    ]
    
    right_box = shapes.add_textbox(Inches(5.2), Inches(1.9), Inches(4.3), Inches(4.5))
    right_tf = right_box.text_frame
    for i, bullet in enumerate(right_bullets):
        if i == 0:
            p = right_tf.paragraphs[0]
        else:
            p = right_tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_GRAY


def add_activity_hub_slide(prs):
    """Slide 3: The Activity Hub - Production Path & Future Solutions"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shapes = slide.shapes
    
    # Title
    title_box = shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
    title_tf = title_box.text_frame
    title_tf.text = "The Activity Hub: Production Path"
    title_tf.paragraphs[0].font.size = Pt(48)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = WALMART_BLUE
    
    # Mission statement
    mission_box = shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(0.8))
    mission_tf = mission_box.text_frame
    mission_tf.text = "Mission: Bring together all store-impacting initiatives into a unified, data-driven platform"
    mission_tf.paragraphs[0].font.size = Pt(24)
    mission_tf.paragraphs[0].font.italic = True
    mission_tf.paragraphs[0].font.color.rgb = DARK_GRAY
    
    # Current state
    current_box = shapes.add_textbox(Inches(0.5), Inches(2.3), Inches(9), Inches(1))
    current_tf = current_box.text_frame
    current_tf.text = "Current: 16 integrated platforms creating the Activity Hub ecosystem"
    p = current_tf.paragraphs[0]
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WALMART_BLUE
    
    # Future solutions
    future_title = shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(0.5))
    future_title_tf = future_title.text_frame
    future_title_tf.text = "A Production Path Enables Future Solutions:"
    future_title_tf.paragraphs[0].font.size = Pt(22)
    future_title_tf.paragraphs[0].font.bold = True
    future_title_tf.paragraphs[0].font.color.rgb = WALMART_YELLOW
    
    # Future solutions bullets
    future_solutions = [
        "Tour Guides - Coordinated store visit management",
        "Calendar Events Resource - Unified event scheduling across Walmart",
        "Store Meeting Planner - Integrated meeting and event coordination",
        "And many others - Built on the same unified platform foundation"
    ]
    
    future_box = shapes.add_textbox(Inches(0.8), Inches(4.1), Inches(8.7), Inches(2.5))
    future_tf = future_box.text_frame
    for i, solution in enumerate(future_solutions):
        if i == 0:
            p = future_tf.paragraphs[0]
        else:
            p = future_tf.add_paragraph()
        p.text = f"✓ {solution}"
        p.font.size = Pt(19)
        p.font.color.rgb = DARK_GRAY
    
    # Bottom note
    note_box = shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.5))
    note_tf = note_box.text_frame
    note_tf.text = "With a robust production foundation, we can rapidly scale to build the solutions Walmart needs tomorrow"
    note_tf.paragraphs[0].font.size = Pt(16)
    note_tf.paragraphs[0].font.italic = True
    note_tf.paragraphs[0].font.color.rgb = RED


def add_path_to_playground_slide(prs):
    """Slide 4: Path to Activity Hub Playground - Solution, Investment, ROI"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shapes = slide.shapes
    
    # Title
    title_box = shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(0.65))
    title_tf = title_box.text_frame
    title_tf.text = "Path to Activity Hub Playground"
    title_tf.paragraphs[0].font.size = Pt(46)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = WALMART_BLUE
    
    # Section 1: Solution
    sol_title = shapes.add_textbox(Inches(0.5), Inches(1.15), Inches(4.3), Inches(0.4))
    sol_title_tf = sol_title.text_frame
    sol_title_tf.text = "Solution"
    sol_title_tf.paragraphs[0].font.size = Pt(22)
    sol_title_tf.paragraphs[0].font.bold = True
    sol_title_tf.paragraphs[0].font.color.rgb = WALMART_YELLOW
    
    solution_items = [
        "Unified Activity Hub Platform",
        "16 Integrated Tools",
        "Production Infrastructure",
        "Data Integration & APIs"
    ]
    
    sol_box = shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(4.3), Inches(2.3))
    sol_tf = sol_box.text_frame
    for i, item in enumerate(solution_items):
        if i == 0:
            p = sol_tf.paragraphs[0]
        else:
            p = sol_tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
    
    # Section 2: Investment
    inv_title = shapes.add_textbox(Inches(3.35), Inches(1.15), Inches(3.3), Inches(0.4))
    inv_title_tf = inv_title.text_frame
    inv_title_tf.text = "Investment"
    inv_title_tf.paragraphs[0].font.size = Pt(22)
    inv_title_tf.paragraphs[0].font.bold = True
    inv_title_tf.paragraphs[0].font.color.rgb = WALMART_YELLOW
    
    inv_details = [
        "Year 1: $184,000 (one-time)",
        "Year 1-3 Ops: $165K-$180K/yr",
        "",
        "Total 3-Year: $699,000"
    ]
    
    inv_box = shapes.add_textbox(Inches(3.35), Inches(1.6), Inches(3.3), Inches(2.3))
    inv_tf = inv_box.text_frame
    for i, detail in enumerate(inv_details):
        if i == 0:
            p = inv_tf.paragraphs[0]
        else:
            p = inv_tf.add_paragraph()
        if detail:
            p.text = detail
            if "Total" in detail:
                p.font.bold = True
                p.font.color.rgb = WALMART_BLUE
                p.font.size = Pt(17)
            else:
                p.font.size = Pt(16)
                p.font.color.rgb = DARK_GRAY
        else:
            p.text = ""
    
    # Section 3: Requirements/Next Steps
    req_title = shapes.add_textbox(Inches(6.7), Inches(1.15), Inches(2.8), Inches(0.4))
    req_title_tf = req_title.text_frame
    req_title_tf.text = "What We Need"
    req_title_tf.paragraphs[0].font.size = Pt(22)
    req_title_tf.paragraphs[0].font.bold = True
    req_title_tf.paragraphs[0].font.color.rgb = WALMART_YELLOW
    
    req_items = [
        "Budget Approval",
        "Cross-Team Alignment",
        "Executive Sponsorship",
        "Go/No-Go Decision"
    ]
    
    req_box = shapes.add_textbox(Inches(6.7), Inches(1.6), Inches(2.8), Inches(2.3))
    req_tf = req_box.text_frame
    for i, item in enumerate(req_items):
        if i == 0:
            p = req_tf.paragraphs[0]
        else:
            p = req_tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
    
    # ROI Section
    roi_divider = shapes.add_shape(1, Inches(0.5), Inches(4.1), Inches(9), Inches(0.02))
    roi_divider.fill.solid()
    roi_divider.fill.fore_color.rgb = WALMART_BLUE
    roi_divider.line.color.rgb = WALMART_BLUE
    
    roi_title = shapes.add_textbox(Inches(0.5), Inches(4.3), Inches(9), Inches(0.4))
    roi_title_tf = roi_title.text_frame
    roi_title_tf.text = "Return on Investment (Real Numbers)"
    roi_title_tf.paragraphs[0].font.size = Pt(28)
    roi_title_tf.paragraphs[0].font.bold = True
    roi_title_tf.paragraphs[0].font.color.rgb = WALMART_BLUE
    
    # ROI Details - Broken into real numbers
    roi_col1_title = shapes.add_textbox(Inches(0.5), Inches(4.85), Inches(3), Inches(0.3))
    roi_col1_title_tf = roi_col1_title.text_frame
    roi_col1_title_tf.text = "Annual Savings by Year"
    roi_col1_title_tf.paragraphs[0].font.size = Pt(18)
    roi_col1_title_tf.paragraphs[0].font.bold = True
    roi_col1_title_tf.paragraphs[0].font.color.rgb = WALMART_YELLOW
    
    roi_col1_items = [
        "Time Savings: $1.2M/year",
        "Error Reduction: $450K/year",
        "Training Efficiency: $200K/year",
        "Total Annual Benefit: $1.85M"
    ]
    
    roi_col1 = shapes.add_textbox(Inches(0.5), Inches(5.25), Inches(3), Inches(1.8))
    roi_col1_tf = roi_col1.text_frame
    for i, item in enumerate(roi_col1_items):
        if i == 0:
            p = roi_col1_tf.paragraphs[0]
        else:
            p = roi_col1_tf.add_paragraph()
        p.text = item
        if "Total" in item:
            p.font.bold = True
            p.font.size = Pt(17)
            p.font.color.rgb = WALMART_BLUE
        else:
            p.font.size = Pt(15)
            p.font.color.rgb = DARK_GRAY
    
    # ROI Timeline
    roi_col2_title = shapes.add_textbox(Inches(3.7), Inches(4.85), Inches(2.8), Inches(0.3))
    roi_col2_title_tf = roi_col2_title.text_frame
    roi_col2_title_tf.text = "Break-Even Analysis"
    roi_col2_title_tf.paragraphs[0].font.size = Pt(18)
    roi_col2_title_tf.paragraphs[0].font.bold = True
    roi_col2_title_tf.paragraphs[0].font.color.rgb = WALMART_YELLOW
    
    roi_col2_items = [
        "Year 1: -$184K (investment)",
        "Year 2: +$820K (net)",
        "Year 3: +$1.67M (cumulative)",
        "Break-even: 12-14 months"
    ]
    
    roi_col2 = shapes.add_textbox(Inches(3.7), Inches(5.25), Inches(2.8), Inches(1.8))
    roi_col2_tf = roi_col2.text_frame
    for i, item in enumerate(roi_col2_items):
        if i == 0:
            p = roi_col2_tf.paragraphs[0]
        else:
            p = roi_col2_tf.add_paragraph()
        p.text = item
        if "Break-even" in item:
            p.font.bold = True
            p.font.color.rgb = WALMART_BLUE
            p.font.size = Pt(16)
        else:
            p.font.size = Pt(14)
            p.font.color.rgb = DARK_GRAY
    
    # 3-Year ROI
    roi_col3_title = shapes.add_textbox(Inches(6.7), Inches(4.85), Inches(2.8), Inches(0.3))
    roi_col3_title_tf = roi_col3_title.text_frame
    roi_col3_title_tf.text = "3-Year Summary"
    roi_col3_title_tf.paragraphs[0].font.size = Pt(18)
    roi_col3_title_tf.paragraphs[0].font.bold = True
    roi_col3_title_tf.paragraphs[0].font.color.rgb = WALMART_YELLOW
    
    roi_col3_items = [
        "Total Investment: $699K",
        "Total Benefit: $5.55M",
        "Net Benefit: $4.85M",
        "ROI: 693%"
    ]
    
    roi_col3 = shapes.add_textbox(Inches(6.7), Inches(5.25), Inches(2.8), Inches(1.8))
    roi_col3_tf = roi_col3.text_frame
    for i, item in enumerate(roi_col3_items):
        if i == 0:
            p = roi_col3_tf.paragraphs[0]
        else:
            p = roi_col3_tf.add_paragraph()
        p.text = item
        if "ROI" in item:
            p.font.bold = True
            p.font.color.rgb = WALMART_BLUE
            p.font.size = Pt(18)
        else:
            p.font.size = Pt(15)
            p.font.color.rgb = DARK_GRAY


def create_presentation():
    """Generate the complete 4-slide presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # SLIDE 1: Title Slide
    add_title_slide(prs, 
                    "Activity Hub", 
                    "Production Path Initiative",
                    "Executive Proposal | December 2025")
    
    # SLIDE 2: The Challenges
    add_challenges_slide(prs)
    
    # SLIDE 3: The Activity Hub - Production Path
    add_activity_hub_slide(prs)
    
    # SLIDE 4: Path to Activity Hub Playground
    add_path_to_playground_slide(prs)
    
    # Save presentation
    prs.save('Activity_Hub_Consolidated_Proposal.pptx')
    print("✓ Generated: Activity_Hub_Consolidated_Proposal.pptx")
    print(f"✓ Total slides: {len(prs.slides)}")
    print("✓ Slide breakdown:")
    print("  - Slide 1: Title")
    print("  - Slide 2: The Challenges (Fragmentation & Pain Points)")
    print("  - Slide 3: The Activity Hub (Production Path & Future Solutions)")
    print("  - Slide 4: Path to Activity Hub Playground (Solution, Investment, ROI)")


if __name__ == "__main__":
    print("Generating Activity Hub Consolidated Proposal PowerPoint...")
    create_presentation()
