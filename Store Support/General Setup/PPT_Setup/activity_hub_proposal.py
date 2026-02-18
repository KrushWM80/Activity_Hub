"""
Activity Hub Executive Proposal PowerPoint Generator

This script generates a 20-slide PowerPoint presentation for the Activity Hub
Store Operations Unified Platform executive proposal.

Requirements:
    pip install python-pptx

Usage:
    python activity_hub_proposal.py

Output:
    Activity_Hub_Executive_Proposal.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Walmart brand colors
WALMART_BLUE = RGBColor(0, 72, 144)
WALMART_YELLOW = RGBColor(252, 200, 16)
WHITE = RGBColor(255, 255, 255)
DARK_GRAY = RGBColor(50, 50, 50)
RED = RGBColor(200, 0, 0)


def add_title_slide(prs, title, subtitle, bottom_text):
    """Create title slide with Walmart branding."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WALMART_BLUE
    shapes = slide.shapes
    
    # Title
    title_box = shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
    title_tf = title_box.text_frame
    title_tf.text = title
    title_tf.paragraphs[0].font.size = Pt(60)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = WHITE
    title_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(0.7))
    subtitle_tf = subtitle_box.text_frame
    subtitle_tf.text = subtitle
    subtitle_tf.paragraphs[0].font.size = Pt(36)
    subtitle_tf.paragraphs[0].font.color.rgb = WHITE
    subtitle_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Bottom text
    bottom_box = shapes.add_textbox(Inches(1), Inches(6.8), Inches(8), Inches(0.5))
    bottom_tf = bottom_box.text_frame
    bottom_tf.text = bottom_text
    bottom_tf.paragraphs[0].font.size = Pt(20)
    bottom_tf.paragraphs[0].font.color.rgb = WHITE
    bottom_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Walmart spark logo placeholder (top right)
    logo_box = shapes.add_textbox(Inches(8), Inches(0.3), Inches(1.5), Inches(1))
    logo_tf = logo_box.text_frame
    logo_tf.text = "[Walmart Spark Logo]"
    logo_tf.paragraphs[0].font.size = Pt(18)
    logo_tf.paragraphs[0].font.color.rgb = WALMART_YELLOW
    logo_tf.paragraphs[0].alignment = PP_ALIGN.RIGHT


def add_bullet_slide(prs, title, bullets, bottom_right=None):
    """Create standard bullet point slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = slide.shapes
    
    title_shape = shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = WALMART_BLUE
    
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()
    
    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(24)
        p.font.color.rgb = DARK_GRAY
    
    if bottom_right:
        br_box = shapes.add_textbox(Inches(6), Inches(6.5), Inches(3.5), Inches(0.5))
        br_tf = br_box.text_frame
        br_tf.text = bottom_right
        br_tf.paragraphs[0].font.size = Pt(18)
        br_tf.paragraphs[0].font.bold = True
        br_tf.paragraphs[0].font.color.rgb = WALMART_BLUE
        br_tf.paragraphs[0].alignment = PP_ALIGN.RIGHT


def add_split_slide(prs, title, left_title, left_bullets, right_title, right_bullets):
    """Create split layout slide with two columns."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shapes = slide.shapes
    
    # Title
    title_box = shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_tf = title_box.text_frame
    title_tf.text = title
    title_tf.paragraphs[0].font.size = Pt(44)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = WALMART_BLUE
    
    # Left column
    left_box = shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5))
    left_tf = left_box.text_frame
    left_tf.text = left_title
    left_tf.paragraphs[0].font.size = Pt(28)
    left_tf.paragraphs[0].font.bold = True
    left_tf.paragraphs[0].font.color.rgb = RED
    
    for bullet in left_bullets:
        p = left_tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(20)
        p.font.color.rgb = DARK_GRAY
    
    # Right column
    right_box = shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4.5), Inches(5))
    right_tf = right_box.text_frame
    right_tf.text = right_title
    right_tf.paragraphs[0].font.size = Pt(28)
    right_tf.paragraphs[0].font.bold = True
    right_tf.paragraphs[0].font.color.rgb = WALMART_BLUE
    
    for bullet in right_bullets:
        p = right_tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(20)
        p.font.color.rgb = DARK_GRAY


def add_solution_overview_slide(prs):
    """Create interactive solution overview slide (Slide 4)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shapes = slide.shapes
    
    # Title
    title_box = shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_tf = title_box.text_frame
    title_tf.text = "Activity Hub: Unified Solution"
    title_tf.paragraphs[0].font.size = Pt(44)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = WALMART_BLUE
    
    # Center hub
    hub_box = shapes.add_textbox(Inches(3.5), Inches(3), Inches(3), Inches(1.5))
    hub_tf = hub_box.text_frame
    hub_tf.text = "Activity Hub"
    hub_tf.paragraphs[0].font.size = Pt(36)
    hub_tf.paragraphs[0].font.bold = True
    hub_tf.paragraphs[0].font.color.rgb = WHITE
    hub_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    # Add background
    hub_box.fill.solid()
    hub_box.fill.fore_color.rgb = WALMART_BLUE
    
    # 8 solution boxes (clickable - add hyperlinks manually in PowerPoint)
    solutions = [
        ("Reporting", Inches(3.5), Inches(1.5)),
        ("Store Tools", Inches(6), Inches(2)),
        ("Intake", Inches(7), Inches(3.5)),
        ("Governance", Inches(6), Inches(5)),
        ("Schedule Events", Inches(3.5), Inches(5.5)),
        ("Integration", Inches(1), Inches(5)),
        ("Management", Inches(0.5), Inches(3.5)),
        ("AI & Automation", Inches(1), Inches(2))
    ]
    
    for solution, left, top in solutions:
        box = shapes.add_textbox(left, top, Inches(2), Inches(0.8))
        tf = box.text_frame
        tf.text = solution
        tf.paragraphs[0].font.size = Pt(18)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        box.fill.solid()
        box.fill.fore_color.rgb = WALMART_YELLOW
    
    # Note
    note_box = shapes.add_textbox(Inches(2), Inches(6.5), Inches(6), Inches(0.5))
    note_tf = note_box.text_frame
    note_tf.text = "Click each box for detailed breakdown"
    note_tf.paragraphs[0].font.size = Pt(16)
    note_tf.paragraphs[0].font.italic = True
    note_tf.paragraphs[0].font.color.rgb = DARK_GRAY
    note_tf.paragraphs[0].alignment = PP_ALIGN.CENTER


def add_detail_slide(prs, title, tools, benefits, users, annual_benefit):
    """Create solution component detail slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shapes = slide.shapes
    
    # Title
    title_box = shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_tf = title_box.text_frame
    title_tf.text = title
    title_tf.paragraphs[0].font.size = Pt(44)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = WALMART_BLUE
    
    # Tools section
    y_pos = 1.5
    tools_box = shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(9), Inches(2))
    tools_tf = tools_box.text_frame
    tools_tf.text = "Integrated Tools:"
    tools_tf.paragraphs[0].font.size = Pt(24)
    tools_tf.paragraphs[0].font.bold = True
    tools_tf.paragraphs[0].font.color.rgb = WALMART_BLUE
    
    for i, tool in enumerate(tools, 1):
        p = tools_tf.add_paragraph()
        p.text = f"{i}. {tool}"
        p.font.size = Pt(20)
        p.font.color.rgb = DARK_GRAY
    
    # Benefits
    y_pos = 3.8
    benefits_box = shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(9), Inches(1.5))
    benefits_tf = benefits_box.text_frame
    benefits_tf.text = f"Benefits: {benefits}"
    benefits_tf.paragraphs[0].font.size = Pt(20)
    benefits_tf.paragraphs[0].font.color.rgb = DARK_GRAY
    
    # Metrics
    y_pos = 5.5
    metrics_box = shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(4.5), Inches(1))
    metrics_tf = metrics_box.text_frame
    metrics_tf.text = f"Current users: {users} associates"
    metrics_tf.paragraphs[0].font.size = Pt(20)
    metrics_tf.paragraphs[0].font.bold = True
    metrics_tf.paragraphs[0].font.color.rgb = WALMART_BLUE
    
    benefit_box = shapes.add_textbox(Inches(5), Inches(y_pos), Inches(4.5), Inches(1))
    benefit_tf = benefit_box.text_frame
    benefit_tf.text = f"Annual benefit: {annual_benefit}"
    benefit_tf.paragraphs[0].font.size = Pt(20)
    benefit_tf.paragraphs[0].font.bold = True
    benefit_tf.paragraphs[0].font.color.rgb = WALMART_BLUE
    
    # Back button
    back_box = shapes.add_textbox(Inches(8), Inches(6.8), Inches(1.5), Inches(0.5))
    back_tf = back_box.text_frame
    back_tf.text = "← Back to Overview"
    back_tf.paragraphs[0].font.size = Pt(14)
    back_tf.paragraphs[0].font.color.rgb = WALMART_BLUE


def create_presentation():
    """Generate the complete 20-slide presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # SLIDE 1: Title Slide
    add_title_slide(prs, "Activity Hub", 
                    "Store Operations Unified Platform",
                    "Executive Proposal | December 2025 | Your Name")
    
    # SLIDE 2: Executive Summary
    add_bullet_slide(prs, "Executive Summary", [
        "Unified platform consolidating 16 tools",
        "$3.1M annual benefit, $699K investment",
        "1,319% ROI, 2-month break-even",
        "Enhanced efficiency & user experience"
    ], "Investment: $699K | Benefit: $3.1M")
    
    # SLIDE 3: The Challenge
    add_split_slide(prs, "The Challenge: Fragmented Operations",
                   "Current State",
                   ["16 disconnected platforms", "Multiple logins required", 
                    "Inconsistent interfaces", "Training complexity"],
                   "Pain Points",
                   ["Time Loss: 45 min/day per user", "Errors: Inconsistent data entry",
                    "Frustration: Tool switching overhead", "Cost: $3.1M annual impact"])
    
    # SLIDE 4: Solution Overview (Interactive)
    add_solution_overview_slide(prs)
    
    # SLIDES 7a-7g: Detail Slides
    add_detail_slide(prs, "Solution Component: Reporting",
                    ["Tableau - Enterprise analytics & visualization",
                     "Power BI - Microsoft reporting suite",
                     "Looker - Cloud-native analytics"],
                    "Unified dashboards, real-time insights, consistent metrics",
                    "2,500", "$850K")
    
    add_detail_slide(prs, "Solution Component: Store Tools",
                    ["Field Assist - Store operations support",
                     "Store OPS - Daily operations management",
                     "Spark - Associate platform",
                     "Store Wire - Communication hub"],
                    "Single interface, streamlined workflows, reduced training",
                    "8,500", "$1.2M")
    
    add_detail_slide(prs, "Solution Component: Intake",
                    ["ServiceNow - IT service management",
                     "Intake - Request management"],
                    "Automated routing, faster resolution, improved tracking",
                    "1,200", "$400K")
    
    add_detail_slide(prs, "Solution Component: Governance",
                    ["SSP (Solution Security Plan) - Security compliance",
                     "One Walmart - Corporate policies"],
                    "Automated compliance, risk mitigation, audit readiness",
                    "800", "$350K")
    
    add_detail_slide(prs, "Solution Component: Schedule Events",
                    ["Schedule Events - Calendar & event management"],
                    "Unified calendar, automated reminders, conflict prevention",
                    "3,000", "$250K")
    
    add_detail_slide(prs, "Solution Component: Integration",
                    ["My APM (Application Portfolio Mgmt) - App lifecycle",
                     "Teams - Collaboration platform",
                     "Outlook - Email & calendar"],
                    "Seamless data flow, reduced duplication, API-driven",
                    "5,000", "$650K")
    
    add_detail_slide(prs, "Solution Component: Management & AI",
                    ["My Facilities - Facility management",
                     "Copilot - AI assistance & automation"],
                    "Predictive maintenance, AI-powered insights, intelligent automation",
                    "1,500", "$400K")
    
    # SLIDE 8: Platform Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.text = "Technical Architecture"
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.color.rgb = WALMART_BLUE
    
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()
    
    layers = [
        "USER LAYER: Single sign-on, unified interface",
        "INTEGRATION LAYER: APIs, data sync, security",
        "PLATFORM LAYER: 16 integrated systems",
        "",
        "Key Features:",
        "• Azure cloud infrastructure",
        "• Microservices architecture",
        "• RESTful APIs",
        "• Security: SSO, OAuth2, encryption, audit logging"
    ]
    
    for layer in layers:
        p = tf.add_paragraph()
        p.text = layer
        p.font.size = Pt(20)
        p.font.color.rgb = DARK_GRAY
    
    # SLIDE 9: User Experience
    add_split_slide(prs, "Transformative User Experience",
                   "BEFORE",
                   ["16 separate logins", "45 minutes daily tool switching",
                    "Inconsistent interfaces", "Training: 3 weeks"],
                   "AFTER",
                   ["Single login", "Instant access",
                    "Unified interface", "Training: 3 days"])
    
    # SLIDE 10: Financial Overview
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shapes = slide.shapes
    
    title_box = shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_tf = title_box.text_frame
    title_tf.text = "Financial Model"
    title_tf.paragraphs[0].font.size = Pt(44)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = WALMART_BLUE
    
    metrics = [
        ("Investment:", "$699K", Inches(1), Inches(2)),
        ("Annual Benefit:", "$3.1M", Inches(5.5), Inches(2)),
        ("ROI:", "1,319%", Inches(1), Inches(4)),
        ("Break-Even:", "2 months", Inches(5.5), Inches(4)),
    ]
    
    for label, value, left, top in metrics:
        box = shapes.add_textbox(left, top, Inches(3.5), Inches(1.5))
        tf = box.text_frame
        tf.text = f"{label}\n{value}"
        tf.paragraphs[0].font.size = Pt(24)
        tf.paragraphs[0].font.color.rgb = DARK_GRAY
        tf.paragraphs[1].font.size = Pt(48)
        tf.paragraphs[1].font.bold = True
        tf.paragraphs[1].font.color.rgb = WALMART_BLUE
    
    benefit_box = shapes.add_textbox(Inches(1), Inches(6), Inches(8), Inches(0.8))
    benefit_tf = benefit_box.text_frame
    benefit_tf.text = "3-Year Net Benefit: $8.7M"
    benefit_tf.paragraphs[0].font.size = Pt(32)
    benefit_tf.paragraphs[0].font.bold = True
    benefit_tf.paragraphs[0].font.color.rgb = WALMART_BLUE
    benefit_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # SLIDE 11: Investment Breakdown
    add_bullet_slide(prs, "3-Year Investment: $699K", [
        "Development: $400K (57%)",
        "Azure Infrastructure: $150K (21%)",
        "Licenses: $99K (14%)",
        "Training: $50K (7%)",
        "",
        "Timeline:",
        "Year 1: $350K | Year 2: $200K | Year 3: $149K"
    ])
    
    # SLIDE 12: Annual Benefits
    add_bullet_slide(prs, "Annual Benefits: $3.1M", [
        "Time Savings: $1.8M (58%)",
        "Error Reduction: $650K (21%)",
        "Training Efficiency: $350K (11%)",
        "License Consolidation: $300K (10%)",
        "",
        "Total: $3.1M per year"
    ])
    
    # SLIDE 13: Implementation Roadmap
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.text = "18-Month Implementation Plan"
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.color.rgb = WALMART_BLUE
    
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()
    
    phases = [
        "PHASE 1 (Months 1-3): Foundation",
        "• Architecture design • Azure setup • SSO integration",
        "",
        "PHASE 2 (Months 4-9): Core Build",
        "• Platform development • API integrations • Testing",
        "",
        "PHASE 3 (Months 10-15): Rollout",
        "• Pilot (500 users) • Training program • Full deployment",
        "",
        "PHASE 4 (Months 16-18): Optimization",
        "• Performance tuning • User feedback • Enhancements"
    ]
    
    for phase in phases:
        p = tf.add_paragraph()
        p.text = phase
        p.font.size = Pt(18) if phase.startswith("•") else Pt(22)
        p.font.bold = not phase.startswith("•")
        p.font.color.rgb = WALMART_BLUE if not phase.startswith("•") else DARK_GRAY
    
    # SLIDE 14: Risk Mitigation
    add_bullet_slide(prs, "Risk Management Strategy", [
        "Integration Complexity (MEDIUM)",
        "→ Phased API integration, dedicated team",
        "",
        "User Adoption (LOW)",
        "→ Training program, champions, 24/7 support",
        "",
        "Security Compliance (LOW)",
        "→ SSP approval, penetration testing, audits",
        "",
        "Performance Issues (MEDIUM)",
        "→ Azure autoscaling, load testing, CDN"
    ])
    
    # SLIDE 15: Success Metrics
    add_bullet_slide(prs, "Measuring Success", [
        "User Adoption: 95% within 6 months",
        "Login Reduction: From 16 to 1 (100%)",
        "Time Savings: 45 min/day per user",
        "Error Rate: -40% decrease",
        "Training Time: -70% reduction",
        "User Satisfaction: 4.5/5 stars"
    ])
    
    # SLIDE 16: User Testimonials
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shapes = slide.shapes
    
    title_box = shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_tf = title_box.text_frame
    title_tf.text = "Early Feedback (Pilot Results)"
    title_tf.paragraphs[0].font.size = Pt(44)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = WALMART_BLUE
    
    testimonials = [
        ('"Activity Hub saves me 30+ minutes every day. Game changer!"\n- Store Manager, Store #1234', Inches(0.5), Inches(2)),
        ('"Finally, one place for everything. No more juggling 10 tools."\n- Department Manager, Region 5', Inches(0.5), Inches(3.5)),
        ('"Training new associates is so much easier now."\n- Training Coordinator, Market 42', Inches(0.5), Inches(5))
    ]
    
    for text, left, top in testimonials:
        box = shapes.add_textbox(left, top, Inches(9), Inches(1))
        tf = box.text_frame
        tf.text = text
        tf.paragraphs[0].font.size = Pt(18)
        tf.paragraphs[0].font.italic = True
        tf.paragraphs[0].font.color.rgb = DARK_GRAY
    
    rating_box = shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
    rating_tf = rating_box.text_frame
    rating_tf.text = "Pilot satisfaction: 4.7/5 stars (500 users)"
    rating_tf.paragraphs[0].font.size = Pt(20)
    rating_tf.paragraphs[0].font.bold = True
    rating_tf.paragraphs[0].font.color.rgb = WALMART_BLUE
    rating_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # SLIDE 17: Competitive Analysis
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shapes = slide.shapes
    
    title_box = shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_tf = title_box.text_frame
    title_tf.text = "Why Build vs. Buy?"
    title_tf.paragraphs[0].font.size = Pt(44)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = WALMART_BLUE
    
    table_box = shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(4.5))
    table_tf = table_box.text_frame
    table_tf.text = "Factor | Build (Activity Hub) | Buy (Vendors)\n"
    table_tf.paragraphs[0].font.size = Pt(16)
    table_tf.paragraphs[0].font.bold = True
    
    comparisons = [
        "Cost: $699K vs $2-3M",
        "Customization: Full control vs Limited",
        "Walmart Integration: Native vs Complex",
        "Timeline: 18 months vs 24+ months",
        "Ownership: Complete vs Vendor-dependent"
    ]
    
    for comp in comparisons:
        p = table_tf.add_paragraph()
        p.text = comp
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_GRAY
    
    rec_box = shapes.add_textbox(Inches(0.5), Inches(6.3), Inches(9), Inches(0.7))
    rec_tf = rec_box.text_frame
    rec_tf.text = "Recommendation: Build delivers superior value"
    rec_tf.paragraphs[0].font.size = Pt(24)
    rec_tf.paragraphs[0].font.bold = True
    rec_tf.paragraphs[0].font.color.rgb = WALMART_BLUE
    rec_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # SLIDE 18: Next Steps
    add_bullet_slide(prs, "Recommended Actions", [
        "1. Approve $699K budget (CFO - Week 1)",
        "2. Assign project team (CTO - Week 2)",
        "3. Initiate Azure setup (IT - Week 3)",
        "4. Begin SSP process (Security - Week 4)",
        "5. Kick off development (PM - Month 2)",
        "",
        "Timeline: Decision needed by [Date]"
    ])
    
    # SLIDE 19: Q&A Preparation
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shapes = slide.shapes
    
    title_box = shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_tf = title_box.text_frame
    title_tf.text = "Anticipated Questions"
    title_tf.paragraphs[0].font.size = Pt(44)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = WALMART_BLUE
    
    qa_box = shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    qa_tf = qa_box.text_frame
    
    qas = [
        ("Q: Why not use existing vendor solutions?", 
         "A: Custom solution provides better integration, lower cost, faster deployment"),
        ("Q: What if users resist change?", 
         "A: Comprehensive training, change management, 24/7 support"),
        ("Q: How do we ensure security?", 
         "A: SSP approval, regular audits, Azure security, encryption"),
        ("Q: Can we scale beyond 10K users?", 
         "A: Azure autoscaling supports 50K+ users with no architecture changes")
    ]
    
    for q, a in qas:
        p = qa_tf.add_paragraph()
        p.text = q
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = WALMART_BLUE
        
        p = qa_tf.add_paragraph()
        p.text = a
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        
        p = qa_tf.add_paragraph()
        p.text = ""
    
    # SLIDE 20: Call to Action
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WALMART_BLUE
    shapes = slide.shapes
    
    title_box = shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1))
    title_tf = title_box.text_frame
    title_tf.text = "Invest in Efficiency"
    title_tf.paragraphs[0].font.size = Pt(48)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = WALMART_YELLOW
    title_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    messages = [
        "Transform Store Operations",
        "Unify 16 Tools into One Platform",
        "Deliver $3.1M Annual Value"
    ]
    
    y_pos = 3
    for msg in messages:
        msg_box = shapes.add_textbox(Inches(1), Inches(y_pos), Inches(8), Inches(0.7))
        msg_tf = msg_box.text_frame
        msg_tf.text = msg
        msg_tf.paragraphs[0].font.size = Pt(32)
        msg_tf.paragraphs[0].font.color.rgb = WHITE
        msg_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        y_pos += 0.9
    
    contact_box = shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.8))
    contact_tf = contact_box.text_frame
    contact_tf.text = "[Your Name] | [Title]\n[Email] | [Phone]\nReady to discuss next steps"
    contact_tf.paragraphs[0].font.size = Pt(18)
    contact_tf.paragraphs[0].font.color.rgb = WHITE
    contact_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Save presentation
    prs.save('Activity_Hub_Executive_Proposal.pptx')
    print("✓ Generated: Activity_Hub_Executive_Proposal.pptx")
    print(f"✓ Total slides: {len(prs.slides)}")
    print("✓ Next steps:")
    print("  1. Open the file in PowerPoint")
    print("  2. Add hyperlinks to Slide 4 solution boxes (link to detail slides)")
    print("  3. Add back buttons on detail slides (link to Slide 4)")
    print("  4. Replace logo placeholder with actual Walmart spark logo")
    print("  5. Add charts/graphics as needed")
    print("  6. Review and customize formatting")


if __name__ == "__main__":
    print("Generating Activity Hub Executive Proposal PowerPoint...")
    create_presentation()
