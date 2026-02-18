"""
Store Support Path To Production - Executive Brief (5 Slides)
Condensed version for quick executive review

This script generates a 5-slide PowerPoint presentation summarizing the
Store Support Path to Production investment proposal.

Requirements:
    pip install python-pptx

Usage:
    python store_support_executive_brief.py

Output:
    Store_Support_Executive_Brief.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Walmart brand colors
WALMART_BLUE = RGBColor(0, 113, 206)  # #0071ce
WALMART_YELLOW = RGBColor(255, 194, 32)  # #ffc220
WHITE = RGBColor(255, 255, 255)
DARK_GRAY = RGBColor(50, 50, 50)
LIGHT_GRAY = RGBColor(200, 200, 200)
GREEN = RGBColor(0, 150, 0)
RED = RGBColor(200, 0, 0)


def add_title_slide(prs, title, subtitle, bottom_text):
    """Create title slide with Walmart branding."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WALMART_BLUE
    shapes = slide.shapes
    
    # Title
    title_box = shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.2))
    title_tf = title_box.text_frame
    title_tf.text = title
    title_tf.paragraphs[0].font.size = Pt(54)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = WHITE
    title_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = shapes.add_textbox(Inches(1), Inches(3.2), Inches(8), Inches(0.8))
    subtitle_tf = subtitle_box.text_frame
    subtitle_tf.text = subtitle
    subtitle_tf.paragraphs[0].font.size = Pt(32)
    subtitle_tf.paragraphs[0].font.color.rgb = WHITE
    subtitle_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Bottom text (investment/ROI)
    bottom_box = shapes.add_textbox(Inches(1), Inches(6), Inches(8), Inches(0.8))
    bottom_tf = bottom_box.text_frame
    bottom_tf.text = bottom_text
    bottom_tf.paragraphs[0].font.size = Pt(24)
    bottom_tf.paragraphs[0].font.bold = True
    bottom_tf.paragraphs[0].font.color.rgb = WALMART_YELLOW
    bottom_tf.paragraphs[0].alignment = PP_ALIGN.CENTER


def add_callout_box(slide, text, left, top, width, height, color):
    """Add a highlighted callout box."""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    tf = box.text_frame
    tf.text = text
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = 1  # Middle


def create_presentation():
    """Create the executive brief presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    print("Creating Store Support Executive Brief Presentation...")
    
    # Slide 1: Title Slide
    print("  Adding Slide 1: Title Slide")
    add_title_slide(
        prs,
        "Store Support Path to Production",
        "Executive Brief",
        "$464K Investment | $25.6M Value | 5,420% ROI"
    )
    
    # Slide 2: The Ask & The Problem
    print("  Adding Slide 2: The Ask & The Problem")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shapes = slide.shapes
    
    title_box = shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_tf = title_box.text_frame
    title_tf.text = "The Two-Part Ask"
    title_tf.paragraphs[0].font.size = Pt(44)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = WALMART_BLUE
    
    # Left: Part 1
    left_box = shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(2.5))
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = RGBColor(240, 248, 255)
    left_tf = left_box.text_frame
    left_tf.text = "PART 1: Path to Production"
    left_tf.paragraphs[0].font.size = Pt(22)
    left_tf.paragraphs[0].font.bold = True
    left_tf.paragraphs[0].font.color.rgb = WALMART_BLUE
    
    part1_items = [
        "🏗️ Production infrastructure",
        "🔒 Security & compliance",
        "⚙️ CI/CD pipeline",
        "📊 Support framework",
        "",
        "$409K (3 years)",
        "Enables ANY product"
    ]
    
    for item in part1_items:
        p = left_tf.add_paragraph()
        p.text = item
        p.font.size = Pt(16) if "$409K" not in item else Pt(18)
        p.font.bold = "$409K" in item or "Enables" in item
        p.font.color.rgb = WALMART_BLUE if "$409K" in item else DARK_GRAY
    
    # Right: Part 2
    right_box = shapes.add_textbox(Inches(5.5), Inches(1.5), Inches(4), Inches(2.5))
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = RGBColor(240, 255, 240)
    right_tf = right_box.text_frame
    right_tf.text = "PART 2: Activity Hub"
    right_tf.paragraphs[0].font.size = Pt(22)
    right_tf.paragraphs[0].font.bold = True
    right_tf.paragraphs[0].font.color.rgb = GREEN
    
    part2_items = [
        "🎯 Central integration platform",
        "🔗 Unifies 16 tools",
        "📈 Real-time visibility",
        "⚡ Automated workflows",
        "",
        "$55K (one-time)",
        "Immediate ROI proof"
    ]
    
    for item in part2_items:
        p = right_tf.add_paragraph()
        p.text = item
        p.font.size = Pt(16) if "$55K" not in item else Pt(18)
        p.font.bold = "$55K" in item or "Immediate" in item
        p.font.color.rgb = GREEN if "$55K" in item else DARK_GRAY
    
    # Bottom: Current Problem
    problem_box = shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2.5))
    problem_tf = problem_box.text_frame
    problem_tf.text = "The Problem We're Solving:"
    problem_tf.paragraphs[0].font.size = Pt(20)
    problem_tf.paragraphs[0].font.bold = True
    problem_tf.paragraphs[0].font.color.rgb = RED
    
    problems = [
        "❌ 15+ disconnected tools creating 275 hrs/week inefficiency",
        "❌ No production deployment capability ($184K per product waste)",
        "❌ Manual processes costing $3.1M annually",
        "❌ Store communication chaos (30-40% message overlap)"
    ]
    
    for problem in problems:
        p = problem_tf.add_paragraph()
        p.text = problem
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(6)
    
    # Slide 3: Platform Assessment Tool & ROI
    print("  Adding Slide 3: Platform Assessment Tool & ROI")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shapes = slide.shapes
    
    title_box = shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_tf = title_box.text_frame
    title_tf.text = "How We Ensure Value"
    title_tf.paragraphs[0].font.size = Pt(44)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = WALMART_BLUE
    
    # Platform Assessment Tool section
    tool_box = shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(2))
    tool_box.fill.solid()
    tool_box.fill.fore_color.rgb = RGBColor(240, 248, 255)
    tool_tf = tool_box.text_frame
    tool_tf.text = "📊 Platform Assessment Tool - Quality Control Gateway"
    tool_tf.paragraphs[0].font.size = Pt(20)
    tool_tf.paragraphs[0].font.bold = True
    tool_tf.paragraphs[0].font.color.rgb = WALMART_BLUE
    
    tool_features = [
        "✅ Business owners self-evaluate product ideas before requesting development",
        "✅ Instant complexity scoring, cost estimates, and timeline recommendations",
        "✅ Only validated, high-value products enter the expensive Path to Production pipeline",
        "✅ Built-in safeguard ensures $464K infrastructure investment has quality control"
    ]
    
    for feature in tool_features:
        p = tool_tf.add_paragraph()
        p.text = feature
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(6)
    
    # ROI Comparison
    roi_title = shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(0.4))
    roi_tf = roi_title.text_frame
    roi_tf.text = "The Numbers: Deploy 5 Products Over 3 Years"
    roi_tf.paragraphs[0].font.size = Pt(20)
    roi_tf.paragraphs[0].font.bold = True
    roi_tf.paragraphs[0].font.color.rgb = DARK_GRAY
    
    # Without Path to Production
    without_box = shapes.add_textbox(Inches(0.5), Inches(4), Inches(4.5), Inches(2.3))
    without_box.fill.solid()
    without_box.fill.fore_color.rgb = RGBColor(255, 240, 240)
    without_tf = without_box.text_frame
    without_tf.text = "❌ Without Path to Production:"
    without_tf.paragraphs[0].font.size = Pt(18)
    without_tf.paragraphs[0].font.bold = True
    without_tf.paragraphs[0].font.color.rgb = RED
    
    without_items = [
        "Product 1: $184K infrastructure",
        "Product 2: $184K infrastructure",
        "Product 3: $184K infrastructure",
        "Product 4: $184K infrastructure",
        "Product 5: $184K infrastructure",
        "",
        "Total: $920K"
    ]
    
    for item in without_items:
        p = without_tf.add_paragraph()
        p.text = item
        p.font.size = Pt(13) if "Total" not in item else Pt(16)
        p.font.bold = "Total" in item
        p.font.color.rgb = RED if "Total" in item else DARK_GRAY
    
    # With Path to Production
    with_box = shapes.add_textbox(Inches(5.5), Inches(4), Inches(4), Inches(2.3))
    with_box.fill.solid()
    with_box.fill.fore_color.rgb = RGBColor(240, 255, 240)
    with_tf = with_box.text_frame
    with_tf.text = "✅ With Path to Production:"
    with_tf.paragraphs[0].font.size = Pt(18)
    with_tf.paragraphs[0].font.bold = True
    with_tf.paragraphs[0].font.color.rgb = GREEN
    
    with_items = [
        "Path to Production: $409K",
        "Product 1 (Activity Hub): $55K",
        "Product 2: $55K",
        "Product 3: $55K",
        "Product 4: $55K",
        "Product 5: $55K",
        "",
        "Total: $684K | Save $236K (26%)"
    ]
    
    for item in with_items:
        p = with_tf.add_paragraph()
        p.text = item
        p.font.size = Pt(13) if "Total" not in item else Pt(16)
        p.font.bold = "Total" in item
        p.font.color.rgb = GREEN if "Total" in item else DARK_GRAY
    
    # Slide 4: Timeline & Benefits
    print("  Adding Slide 4: Timeline & Benefits")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shapes = slide.shapes
    
    title_box = shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_tf = title_box.text_frame
    title_tf.text = "Timeline & Benefits"
    title_tf.paragraphs[0].font.size = Pt(44)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = WALMART_BLUE
    
    # Timeline
    timeline_box = shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4.5), Inches(2.5))
    timeline_tf = timeline_box.text_frame
    timeline_tf.text = "⏱️ Implementation Timeline"
    timeline_tf.paragraphs[0].font.size = Pt(20)
    timeline_tf.paragraphs[0].font.bold = True
    timeline_tf.paragraphs[0].font.color.rgb = WALMART_BLUE
    
    timeline_items = [
        "",
        "Months 1-2: Foundation",
        "  APM registration, Azure setup",
        "",
        "Months 3-4: Security",
        "  SSP approval, compliance",
        "",
        "Month 5: Testing & Deploy",
        "  CI/CD pipeline, Activity Hub",
        "",
        "Month 6+: Value Realization",
        "  Additional products follow"
    ]
    
    for item in timeline_items:
        p = timeline_tf.add_paragraph()
        p.text = item
        p.font.size = Pt(14) if item.startswith("  ") else Pt(16)
        p.font.bold = not item.startswith("  ") and item != ""
        p.font.color.rgb = DARK_GRAY
    
    # Benefits
    benefits_box = shapes.add_textbox(Inches(5.5), Inches(1.3), Inches(4), Inches(5.3))
    benefits_box.fill.solid()
    benefits_box.fill.fore_color.rgb = RGBColor(240, 255, 240)
    benefits_tf = benefits_box.text_frame
    benefits_tf.text = "💰 Annual Benefits Breakdown"
    benefits_tf.paragraphs[0].font.size = Pt(20)
    benefits_tf.paragraphs[0].font.bold = True
    benefits_tf.paragraphs[0].font.color.rgb = GREEN
    
    benefits_items = [
        "",
        "Process Efficiency:",
        "  $3.1M (275 hrs/week saved)",
        "",
        "Infrastructure Savings:",
        "  $2.1M (reusable platform)",
        "",
        "Better Decision Making:",
        "  $1.4M (real-time data)",
        "",
        "Reduced Store Confusion:",
        "  $900K (communication efficiency)",
        "",
        "Total Annual: $7.5M",
        "3-Year Total: $25.6M"
    ]
    
    for item in benefits_items:
        p = benefits_tf.add_paragraph()
        p.text = item
        p.font.size = Pt(14) if item.startswith("  ") else Pt(15)
        p.font.bold = "Total" in item or (":" in item and not item.startswith(" "))
        p.font.color.rgb = GREEN if "Total" in item else DARK_GRAY
    
    # Bottom callout
    add_callout_box(slide, "Break-Even: Month 2 | Every month of delay costs $262K",
                   0.5, 6.5, 9, 0.7, RED)
    
    # Slide 5: The Ask & Call to Action
    print("  Adding Slide 5: The Ask & Call to Action")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WALMART_BLUE
    shapes = slide.shapes
    
    title_box = shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    title_tf = title_box.text_frame
    title_tf.text = "Ready to Transform Store Support?"
    title_tf.paragraphs[0].font.size = Pt(44)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = WHITE
    title_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # The Ask
    ask_box = shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(7), Inches(2))
    ask_tf = ask_box.text_frame
    ask_tf.text = "We Need Two Approvals:"
    ask_tf.paragraphs[0].font.size = Pt(28)
    ask_tf.paragraphs[0].font.bold = True
    ask_tf.paragraphs[0].font.color.rgb = WALMART_YELLOW
    ask_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    approvals = [
        "",
        "☑️ Path to Production Infrastructure",
        "   $409K over 3 years",
        "",
        "☑️ Activity Hub Platform",
        "   $55K one-time"
    ]
    
    for approval in approvals:
        p = ask_tf.add_paragraph()
        p.text = approval
        p.font.size = Pt(24) if approval.startswith("☑️") else Pt(18)
        p.font.bold = approval.startswith("☑️")
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
    
    # Key Stats
    stats_box = shapes.add_textbox(Inches(1.5), Inches(4.5), Inches(7), Inches(1.5))
    stats_tf = stats_box.text_frame
    stats_tf.text = "$464K Investment | $25.6M Value"
    stats_tf.paragraphs[0].font.size = Pt(32)
    stats_tf.paragraphs[0].font.bold = True
    stats_tf.paragraphs[0].font.color.rgb = WALMART_YELLOW
    stats_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    p = stats_tf.add_paragraph()
    p.text = "5,420% ROI | Break-Even: Month 2"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WALMART_YELLOW
    p.alignment = PP_ALIGN.CENTER
    
    p = stats_tf.add_paragraph()
    p.text = ""
    
    p = stats_tf.add_paragraph()
    p.text = "Target Launch: August 2026"
    p.font.size = Pt(22)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # CTA
    cta_box = shapes.add_textbox(Inches(2), Inches(6.3), Inches(6), Inches(0.8))
    cta_tf = cta_box.text_frame
    cta_tf.text = "Let's Build the Future of Store Support"
    cta_tf.paragraphs[0].font.size = Pt(32)
    cta_tf.paragraphs[0].font.bold = True
    cta_tf.paragraphs[0].font.color.rgb = WHITE
    cta_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    return prs


def main():
    """Main execution function."""
    print("\n" + "="*60)
    print("Store Support Executive Brief Generator (5 Slides)")
    print("="*60 + "\n")
    
    prs = create_presentation()
    
    output_file = "Store_Support_Executive_Brief.pptx"
    prs.save(output_file)
    
    print("\n" + "="*60)
    print(f"✓ Executive Brief created successfully!")
    print(f"✓ File saved as: {output_file}")
    print(f"✓ Total slides: {len(prs.slides)}")
    print("="*60 + "\n")
    print("This condensed version includes:")
    print("  • Slide 1: Title")
    print("  • Slide 2: The Ask & The Problem")
    print("  • Slide 3: Platform Assessment Tool & ROI")
    print("  • Slide 4: Timeline & Benefits")
    print("  • Slide 5: The Ask & Call to Action")
    print("\n")


if __name__ == "__main__":
    main()
