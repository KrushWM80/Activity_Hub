"""
Two-Part Approval Request PowerPoint Generator
Path to Production + Activity Hub

This script generates a 24-slide PowerPoint presentation for the Two-Part
Approval Request (Path to Production Infrastructure + Activity Hub Platform).

Requirements:
    pip install python-pptx

Usage:
    python store_support_path_to_production.py

Output:
    Store_Support_Path_To_Production.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Walmart brand colors
WALMART_BLUE = RGBColor(0, 113, 206)  # #0071ce
WALMART_YELLOW = RGBColor(255, 194, 32)  # #ffc220
WHITE = RGBColor(255, 255, 255)
DARK_GRAY = RGBColor(50, 50, 50)
LIGHT_GRAY = RGBColor(200, 200, 200)
GREEN = RGBColor(0, 150, 0)
RED = RGBColor(200, 0, 0)


def add_title_slide(prs, title, subtitle, bottom_text):
    """Create title slide with Walmart branding."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WALMART_BLUE
    shapes = slide.shapes
    
    # Title
    title_box = shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.2))
    title_tf = title_box.text_frame
    title_tf.text = title
    title_tf.paragraphs[0].font.size = Pt(54)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = WHITE
    title_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = shapes.add_textbox(Inches(1), Inches(3.2), Inches(8), Inches(0.8))
    subtitle_tf = subtitle_box.text_frame
    subtitle_tf.text = subtitle
    subtitle_tf.paragraphs[0].font.size = Pt(32)
    subtitle_tf.paragraphs[0].font.color.rgb = WHITE
    subtitle_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Bottom text (investment/ROI)
    bottom_box = shapes.add_textbox(Inches(1), Inches(6), Inches(8), Inches(0.8))
    bottom_tf = bottom_box.text_frame
    bottom_tf.text = bottom_text
    bottom_tf.paragraphs[0].font.size = Pt(24)
    bottom_tf.paragraphs[0].font.bold = True
    bottom_tf.paragraphs[0].font.color.rgb = WALMART_YELLOW
    bottom_tf.paragraphs[0].alignment = PP_ALIGN.CENTER


def add_section_title(prs, title, color=WALMART_BLUE):
    """Create section divider slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color
    shapes = slide.shapes
    
    title_box = shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1.5))
    title_tf = title_box.text_frame
    title_tf.text = title
    title_tf.paragraphs[0].font.size = Pt(60)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = WHITE
    title_tf.paragraphs[0].alignment = PP_ALIGN.CENTER


def add_content_slide(prs, title, content_dict):
    """Create standard content slide with title and text boxes."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shapes = slide.shapes
    
    # Title
    title_box = shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_tf = title_box.text_frame
    title_tf.text = title
    title_tf.paragraphs[0].font.size = Pt(40)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = WALMART_BLUE
    
    # Content (passed as dict with positions)
    y_position = 1.5
    for key, text in content_dict.items():
        box = shapes.add_textbox(Inches(0.8), Inches(y_position), Inches(8.4), Inches(0.6))
        tf = box.text_frame
        tf.text = text
        tf.paragraphs[0].font.size = Pt(18)
        tf.paragraphs[0].font.color.rgb = DARK_GRAY
        tf.word_wrap = True
        y_position += 0.7


def add_two_column_slide(prs, title, left_content, right_content, left_color=WALMART_BLUE, right_color=GREEN):
    """Create two-column comparison slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shapes = slide.shapes
    
    # Title
    title_box = shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_tf = title_box.text_frame
    title_tf.text = title
    title_tf.paragraphs[0].font.size = Pt(40)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = WALMART_BLUE
    
    # Left column header
    left_header = shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(0.6))
    left_header.fill.solid()
    left_header.fill.fore_color.rgb = left_color
    left_tf = left_header.text_frame
    left_tf.text = left_content['title']
    left_tf.paragraphs[0].font.size = Pt(24)
    left_tf.paragraphs[0].font.bold = True
    left_tf.paragraphs[0].font.color.rgb = WHITE
    left_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Left column content
    left_box = shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(4.5), Inches(4.5))
    left_text_tf = left_box.text_frame
    for i, bullet in enumerate(left_content['bullets']):
        if i == 0:
            left_text_tf.text = bullet
        else:
            p = left_text_tf.add_paragraph()
            p.text = bullet
        left_text_tf.paragraphs[i].font.size = Pt(16)
        left_text_tf.paragraphs[i].font.color.rgb = DARK_GRAY
        left_text_tf.paragraphs[i].space_after = Pt(8)
    
    # Right column header
    right_header = shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4.5), Inches(0.6))
    right_header.fill.solid()
    right_header.fill.fore_color.rgb = right_color
    right_tf = right_header.text_frame
    right_tf.text = right_content['title']
    right_tf.paragraphs[0].font.size = Pt(24)
    right_tf.paragraphs[0].font.bold = True
    right_tf.paragraphs[0].font.color.rgb = WHITE
    right_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Right column content
    right_box = shapes.add_textbox(Inches(5.2), Inches(2.2), Inches(4.5), Inches(4.5))
    right_text_tf = right_box.text_frame
    for i, bullet in enumerate(right_content['bullets']):
        if i == 0:
            right_text_tf.text = bullet
        else:
            p = right_text_tf.add_paragraph()
            p.text = bullet
        right_text_tf.paragraphs[i].font.size = Pt(16)
        right_text_tf.paragraphs[i].font.color.rgb = DARK_GRAY
        right_text_tf.paragraphs[i].space_after = Pt(8)


def add_table_slide(prs, title, headers, rows, highlight_row=None):
    """Create slide with data table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shapes = slide.shapes
    
    # Title
    title_box = shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_tf = title_box.text_frame
    title_tf.text = title
    title_tf.paragraphs[0].font.size = Pt(40)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = WALMART_BLUE
    
    # Create table
    num_cols = len(headers)
    num_rows = len(rows) + 1  # +1 for header
    
    table = shapes.add_table(num_rows, num_cols, Inches(0.8), Inches(1.8), Inches(8.4), Inches(4.5)).table
    
    # Set column widths
    col_width = Inches(8.4 / num_cols)
    for col_idx in range(num_cols):
        table.columns[col_idx].width = col_width
    
    # Header row
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = WALMART_BLUE
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(16)
            paragraph.font.bold = True
            paragraph.font.color.rgb = WHITE
            paragraph.alignment = PP_ALIGN.CENTER
    
    # Data rows
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_data in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_data)
            
            # Highlight specific row if needed
            if highlight_row and row_idx + 1 == highlight_row:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WALMART_YELLOW
            
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(14)
                paragraph.font.color.rgb = DARK_GRAY
                if col_idx > 0:  # Right-align numbers
                    paragraph.alignment = PP_ALIGN.RIGHT


def add_callout_box(slide, text, left, top, width, height, color=WALMART_YELLOW):
    """Add a callout box to a slide."""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    tf = box.text_frame
    tf.text = text
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = DARK_GRAY
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER


def create_presentation():
    """Generate the complete Two-Part Approval Request presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    print("Creating Two-Part Approval Request Presentation...")
    
    # Slide 1: Title Slide
    print("  Adding Slide 1: Title Slide")
    add_title_slide(
        prs,
        "Two-Part Investment Proposal",
        "Path to Production Infrastructure + Activity Hub Platform",
        "Investment: $464K over 3 years | ROI: 5,420%"
    )
    
    # Slide 2: The Two-Part Ask
    print("  Adding Slide 2: The Two-Part Ask")
    add_two_column_slide(
        prs,
        "We Need Two Approvals",
        {
            'title': 'PART 1: Path to Production',
            'bullets': [
                '🏗️ Infrastructure Foundation',
                '',
                '$129K one-time + $90K/year',
                '',
                'Enables ANY product deployment',
                '',
                '3-year total: $409K'
            ]
        },
        {
            'title': 'PART 2: Activity Hub',
            'bullets': [
                '🎯 Flagship Integration Platform',
                '',
                '$55K one-time development',
                '',
                'Unifies 16 Store Support tools',
                '',
                '3-year total: $55K'
            ]
        }
    )
    
    # Slide 3: Why Two Separate Approvals?
    print("  Adding Slide 3: Why Two Separate Approvals?")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shapes = slide.shapes
    
    title_box = shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_tf = title_box.text_frame
    title_tf.text = "Why Two Separate Approvals?"
    title_tf.paragraphs[0].font.size = Pt(40)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = WALMART_BLUE
    
    # Three boxes
    boxes_data = [
        ("Path to Production = Foundation", 1.5, [
            "✅ Production Azure environment",
            "✅ Security & compliance framework",
            "✅ Development pipeline",
            "✅ Testing & deployment automation",
            "Enables: Any future product"
        ]),
        ("Activity Hub = Proof of Value", 3.7, [
            "✅ Central integration platform",
            "✅ 16-tool ecosystem",
            "✅ Real-time visibility",
            "✅ Automated workflows",
            "Demonstrates: Immediate ROI"
        ]),
        ("Combined = Maximum Efficiency", 5.9, [
            "✅ No duplicate infrastructure costs",
            "✅ Faster time to value",
            "✅ Proven deployment model",
            "✅ Scalable for 15+ future products",
            "Delivers: $25.6M over 3 years"
        ])
    ]
    
    for title, top, bullets in boxes_data:
        # Box header
        header = shapes.add_textbox(Inches(0.8), Inches(top), Inches(8.4), Inches(0.4))
        header.fill.solid()
        header.fill.fore_color.rgb = WALMART_BLUE
        header_tf = header.text_frame
        header_tf.text = title
        header_tf.paragraphs[0].font.size = Pt(18)
        header_tf.paragraphs[0].font.bold = True
        header_tf.paragraphs[0].font.color.rgb = WHITE
        header_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Box content
        content = shapes.add_textbox(Inches(0.8), Inches(top + 0.45), Inches(8.4), Inches(1.8))
        content_tf = content.text_frame
        for i, bullet in enumerate(bullets):
            if i == 0:
                content_tf.text = bullet
            else:
                p = content_tf.add_paragraph()
                p.text = bullet
            content_tf.paragraphs[i].font.size = Pt(14)
            content_tf.paragraphs[i].font.color.rgb = DARK_GRAY
    
    # Slide 4: Platform Assessment Tool
    print("  Adding Slide 4: Platform Assessment Tool")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shapes = slide.shapes
    
    title_box = shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_tf = title_box.text_frame
    title_tf.text = "How We Evaluate New Product Requests"
    title_tf.paragraphs[0].font.size = Pt(40)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = WALMART_BLUE
    
    # Left side: Tool Overview
    left_box = shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(3))
    left_tf = left_box.text_frame
    left_tf.text = "📊 Store Support Platform Assessment Tool"
    left_tf.paragraphs[0].font.size = Pt(18)
    left_tf.paragraphs[0].font.bold = True
    left_tf.paragraphs[0].font.color.rgb = WALMART_BLUE
    
    tool_features = [
        "• Self-service evaluation for business owners",
        "• Simple questionnaire (no technical knowledge needed)",
        "• Instant complexity & cost estimates",
        "• Guides prioritization decisions"
    ]
    
    for feature in tool_features:
        p = left_tf.add_paragraph()
        p.text = feature
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(8)
    
    # Right side: What It Provides
    right_box = shapes.add_textbox(Inches(5.5), Inches(1.5), Inches(4), Inches(3))
    right_tf = right_box.text_frame
    right_tf.text = "What It Provides:"
    right_tf.paragraphs[0].font.size = Pt(18)
    right_tf.paragraphs[0].font.bold = True
    right_tf.paragraphs[0].font.color.rgb = WALMART_BLUE
    
    outputs = [
        "✅ Complexity score (Low/Medium/High)",
        "✅ Technical requirements summary",
        "✅ Cost estimates (Year 1)",
        "✅ Timeline recommendations",
        "✅ PDF export for proposals"
    ]
    
    for output in outputs:
        p = right_tf.add_paragraph()
        p.text = output
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(8)
    
    # Bottom callout
    callout_box = shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(9), Inches(1.2))
    callout_box.fill.solid()
    callout_box.fill.fore_color.rgb = RGBColor(240, 248, 255)  # Light blue
    callout_tf = callout_box.text_frame
    callout_tf.text = "Use Case:"
    callout_tf.paragraphs[0].font.size = Pt(14)
    callout_tf.paragraphs[0].font.bold = True
    callout_tf.paragraphs[0].font.color.rgb = WALMART_BLUE
    
    p = callout_tf.add_paragraph()
    p.text = "Before requesting Path to Production deployment, business owners use this tool to determine if a product idea should be added to our development pipeline."
    p.font.size = Pt(14)
    p.font.color.rgb = DARK_GRAY
    
    # File location note
    location_box = shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(9), Inches(0.6))
    location_tf = location_box.text_frame
    location_tf.text = "Tool Location: General Setup/Platform-Assessment/assessment_tool.html"
    location_tf.paragraphs[0].font.size = Pt(11)
    location_tf.paragraphs[0].font.italic = True
    location_tf.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)
    location_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Slide 5: Current State - The Problem
    print("  Adding Slide 5: Current State - The Problem")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shapes = slide.shapes
    
    title_box = shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_tf = title_box.text_frame
    title_tf.text = "Why We Need Both Investments NOW"
    title_tf.paragraphs[0].font.size = Pt(40)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = WALMART_BLUE
    
    # Four quadrants
    quadrants = [
        ("Fragmented Tools 🔴", 0.5, 1.5, [
            "15+ disconnected platforms",
            "Manual data entry across systems",
            "No single source of truth",
            "Impact: 275 hrs/week wasted"
        ]),
        ("No Production Capability 🔴", 5.2, 1.5, [
            "No secure deployment process",
            "Each tool requires separate infrastructure",
            "Compliance repeated per product",
            "Impact: $184K per product waste"
        ]),
        ("Limited Visibility 🔴", 0.5, 4.5, [
            "Leadership lacks real-time data",
            "Store communication chaos (30-40% overlap)",
            "No systematic governance",
            "Impact: Delayed decisions, confused stores"
        ]),
        ("Manual Processes 🔴", 5.2, 4.5, [
            "Spreadsheet-based tracking",
            "Email coordination",
            "No automation",
            "Impact: $3.1M/year in inefficiency"
        ])
    ]
    
    for title, left, top, bullets in quadrants:
        box = shapes.add_textbox(Inches(left), Inches(top), Inches(4.5), Inches(2.5))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(255, 240, 240)  # Light red
        tf = box.text_frame
        tf.text = title
        tf.paragraphs[0].font.size = Pt(18)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RED
        
        for bullet in bullets:
            p = tf.add_paragraph()
            p.text = bullet
            p.font.size = Pt(14)
            p.font.color.rgb = DARK_GRAY
            p.space_after = Pt(4)
    
    # PART 1: PATH TO PRODUCTION SECTION
    print("  Adding Slide 6: Section Title - Path to Production")
    add_section_title(prs, "PART 1:\nPath to Production")
    
    # Slide 7: Path to Production Explained
    print("  Adding Slide 7: Path to Production Explained")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shapes = slide.shapes
    
    title_box = shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_tf = title_box.text_frame
    title_tf.text = "What is Path to Production?"
    title_tf.paragraphs[0].font.size = Pt(40)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = WALMART_BLUE
    
    # Building blocks visual
    blocks = [
        ("Production Infrastructure", 1.8, "$45K", [
            "Enterprise Azure environment",
            "Security controls & monitoring",
            "Disaster recovery"
        ]),
        ("Compliance & Governance", 3.3, "$15K", [
            "APM registration process",
            "Security certifications (SSP)",
            "Data classification framework"
        ]),
        ("Development Pipeline", 4.8, "$25K", [
            "CI/CD automation",
            "Testing frameworks",
            "Quality assurance"
        ]),
        ("Operational Capability", 6.3, "$44K", [
            "Help desk & support",
            "Training standards",
            "Change management"
        ])
    ]
    
    x_position = 0.8
    for title, top, cost, bullets in blocks:
        # Block
        block = shapes.add_textbox(Inches(x_position), Inches(top), Inches(2), Inches(1))
        block.fill.solid()
        block.fill.fore_color.rgb = WALMART_BLUE
        block_tf = block.text_frame
        block_tf.text = f"{title}\n{cost}"
        block_tf.paragraphs[0].font.size = Pt(14)
        block_tf.paragraphs[0].font.bold = True
        block_tf.paragraphs[0].font.color.rgb = WHITE
        block_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        x_position += 2.1
    
    # Total
    total_box = shapes.add_textbox(Inches(1), Inches(7), Inches(8), Inches(0.4))
    total_tf = total_box.text_frame
    total_tf.text = "Total Stack: $129K one-time → Enables unlimited products"
    total_tf.paragraphs[0].font.size = Pt(20)
    total_tf.paragraphs[0].font.bold = True
    total_tf.paragraphs[0].font.color.rgb = WALMART_BLUE
    total_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Slide 8: Path to Production ROI
    print("  Adding Slide 8: Path to Production ROI")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shapes = slide.shapes
    
    title_box = shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_tf = title_box.text_frame
    title_tf.text = "The Infrastructure Investment That Keeps Giving"
    title_tf.paragraphs[0].font.size = Pt(36)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = WALMART_BLUE
    
    # Scenario comparison
    scenario_title = shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.4), Inches(0.4))
    scenario_tf = scenario_title.text_frame
    scenario_tf.text = "Scenario: Deploy 5 Products Over 3 Years"
    scenario_tf.paragraphs[0].font.size = Pt(22)
    scenario_tf.paragraphs[0].font.bold = True
    scenario_tf.paragraphs[0].font.color.rgb = DARK_GRAY
    scenario_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Without Path to Production
    without_box = shapes.add_textbox(Inches(0.8), Inches(2), Inches(4.2), Inches(3.5))
    without_box.fill.solid()
    without_box.fill.fore_color.rgb = RGBColor(255, 240, 240)
    without_tf = without_box.text_frame
    without_tf.text = "Without Path to Production (Traditional):"
    without_tf.paragraphs[0].font.size = Pt(18)
    without_tf.paragraphs[0].font.bold = True
    without_tf.paragraphs[0].font.color.rgb = RED
    
    products = ["Product 1: $184K infrastructure",
                "Product 2: $184K infrastructure",
                "Product 3: $184K infrastructure",
                "Product 4: $184K infrastructure",
                "Product 5: $184K infrastructure",
                "",
                "Total: $920K"]
    
    for product in products:
        p = without_tf.add_paragraph()
        p.text = product
        p.font.size = Pt(14) if "Total" not in product else Pt(16)
        p.font.bold = "Total" in product
        p.font.color.rgb = DARK_GRAY
    
    # With Path to Production
    with_box = shapes.add_textbox(Inches(5.2), Inches(2), Inches(4.2), Inches(3.5))
    with_box.fill.solid()
    with_box.fill.fore_color.rgb = RGBColor(240, 255, 240)
    with_tf = with_box.text_frame
    with_tf.text = "With Path to Production:"
    with_tf.paragraphs[0].font.size = Pt(18)
    with_tf.paragraphs[0].font.bold = True
    with_tf.paragraphs[0].font.color.rgb = GREEN
    
    products_with = ["Path to Production: $409K",
                     "Product 1 (Activity Hub): $55K",
                     "Product 2: $55K",
                     "Product 3: $55K",
                     "Product 4: $55K",
                     "Product 5: $55K",
                     "",
                     "Total: $684K"]
    
    for product in products_with:
        p = with_tf.add_paragraph()
        p.text = product
        p.font.size = Pt(14) if "Total" not in product else Pt(16)
        p.font.bold = "Total" in product
        p.font.color.rgb = DARK_GRAY
    
    # Savings callout
    add_callout_box(slide, "Savings: $236,000 (26% reduction)\nBreak-Even: After 3 products",
                   2.5, 5.8, 5, 1, WALMART_YELLOW)
    
    # Slide 9: Path to Production Timeline
    print("  Adding Slide 9: Path to Production Timeline")
    add_table_slide(
        prs,
        "Path to Production: 5 Months to Operational",
        ["Timeline", "Activities", "Deliverable"],
        [
            ["Month 1-2:\nFoundation", "APM registration, Team roster, Azure provisioning", "Infrastructure live"],
            ["Month 3-4:\nSecurity", "Solution Security Plan, Security controls, Compliance docs", "SSP Approved status"],
            ["Month 5:\nTesting", "CI/CD pipeline, Test environment, Training framework", "Production-ready platform"],
            ["Month 6+:\nDeploy", "Activity Hub deployment, Additional products follow", "Value realization begins"]
        ]
    )
    
    # PART 2: ACTIVITY HUB SECTION
    print("  Adding Slide 10: Section Title - Activity Hub")
    add_section_title(prs, "PART 2:\nActivity Hub", GREEN)
    
    # Slide 11: Activity Hub Explained
    print("  Adding Slide 11: Activity Hub Explained")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shapes = slide.shapes
    
    title_box = shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_tf = title_box.text_frame
    title_tf.text = "What is Activity Hub?"
    title_tf.paragraphs[0].font.size = Pt(40)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = WALMART_BLUE
    
    # Center hub
    hub_box = shapes.add_textbox(Inches(3.5), Inches(3), Inches(3), Inches(1.5))
    hub_box.fill.solid()
    hub_box.fill.fore_color.rgb = WALMART_BLUE
    hub_tf = hub_box.text_frame
    hub_tf.text = "Activity Hub\n\nCentral Integration Platform"
    hub_tf.paragraphs[0].font.size = Pt(24)
    hub_tf.paragraphs[0].font.bold = True
    hub_tf.paragraphs[0].font.color.rgb = WHITE
    hub_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Spokes (16 platforms grouped)
    spokes = [
        ("Reporting (6)", 1.5, 1.8),
        ("Intake (4)", 7.5, 1.8),
        ("Store Tools (3)", 1.5, 5.8),
        ("Specialized (3)", 7.5, 5.8)
    ]
    
    for spoke, left, top in spokes:
        spoke_box = shapes.add_textbox(Inches(left), Inches(top), Inches(1.8), Inches(0.8))
        spoke_box.fill.solid()
        spoke_box.fill.fore_color.rgb = WALMART_YELLOW
        spoke_tf = spoke_box.text_frame
        spoke_tf.text = spoke
        spoke_tf.paragraphs[0].font.size = Pt(16)
        spoke_tf.paragraphs[0].font.bold = True
        spoke_tf.paragraphs[0].font.color.rgb = DARK_GRAY
        spoke_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Bottom text
    bottom = shapes.add_textbox(Inches(1), Inches(6.8), Inches(8), Inches(0.5))
    bottom_tf = bottom.text_frame
    bottom_tf.text = "One Way of Working for 4,700 stores and 100+ Home Office stakeholders"
    bottom_tf.paragraphs[0].font.size = Pt(18)
    bottom_tf.paragraphs[0].font.bold = True
    bottom_tf.paragraphs[0].font.color.rgb = WALMART_BLUE
    bottom_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Slide 12: Current vs Future State
    print("  Adding Slide 12: Current vs Future State")
    add_two_column_slide(
        prs,
        "The Transformation",
        {
            'title': 'BEFORE (Current State)',
            'bullets': [
                '🔴 15+ disconnected tools',
                '🔴 275 hours/week manual work',
                '🔴 30-40% message overlap',
                '🔴 No real-time visibility',
                '🔴 Spreadsheet tracking',
                '🔴 Email coordination',
                '🔴 Store confusion'
            ]
        },
        {
            'title': 'AFTER (With Activity Hub)',
            'bullets': [
                '✅ 1 integrated platform',
                '✅ 58 hours/week automated work',
                '✅ 5% message overlap (smart detection)',
                '✅ Real-time executive dashboard',
                '✅ Automated tracking',
                '✅ Workflow automation',
                '✅ Clear store communication'
            ]
        },
        RED,
        GREEN
    )
    
    # Add banner at bottom
    slide = prs.slides[-1]
    banner = slide.shapes.add_textbox(Inches(1), Inches(6.8), Inches(8), Inches(0.5))
    banner.fill.solid()
    banner.fill.fore_color.rgb = WALMART_YELLOW
    banner_tf = banner.text_frame
    banner_tf.text = "Time Saved: 217 hours/week | Efficiency Gained: 79%"
    banner_tf.paragraphs[0].font.size = Pt(20)
    banner_tf.paragraphs[0].font.bold = True
    banner_tf.paragraphs[0].font.color.rgb = DARK_GRAY
    banner_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Slide 13: Activity Hub Benefits Breakdown
    print("  Adding Slide 13: Activity Hub Benefits Breakdown")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shapes = slide.shapes
    
    title_box = shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_tf = title_box.text_frame
    title_tf.text = "$3.1M Annual Benefits Across Four Categories"
    title_tf.paragraphs[0].font.size = Pt(36)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = WALMART_BLUE
    
    # Four benefit boxes
    benefits = [
        ("Time Savings", "$564,200/year", 0.5, 1.8, [
            "Manual reporting: 100 hrs/week saved",
            "Communication: 30 hrs/week saved",
            "Issue tracking: 45 hrs/week saved",
            "Tour data: 25 hrs/week saved"
        ]),
        ("Communication Efficiency", "$481,000/year", 5.2, 1.8, [
            "30% reduction in message overlap",
            "25% increase in targeted messaging",
            "Reduced store confusion/rework"
        ]),
        ("Governance & Decisions", "$470,000/year", 0.5, 4.5, [
            "Real-time visibility → faster decisions",
            "Conflict detection → prevents duplicates",
            "Data-driven resource allocation"
        ]),
        ("Store Productivity", "$1,631,000/year", 5.2, 4.5, [
            "4,700 stores × 2.5 hrs/week saved",
            "Streamlined issue reporting",
            "Efficient tour coordination"
        ])
    ]
    
    for title, amount, left, top, bullets in benefits:
        box = shapes.add_textbox(Inches(left), Inches(top), Inches(4.5), Inches(2.2))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(240, 250, 255)
        tf = box.text_frame
        tf.text = f"{title} 💰"
        tf.paragraphs[0].font.size = Pt(18)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = WALMART_BLUE
        
        p = tf.add_paragraph()
        p.text = amount
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = GREEN
        
        for bullet in bullets:
            p = tf.add_paragraph()
            p.text = f"• {bullet}"
            p.font.size = Pt(12)
            p.font.color.rgb = DARK_GRAY
            p.space_after = Pt(3)
    
    # Total
    total = shapes.add_textbox(Inches(1), Inches(7), Inches(8), Inches(0.4))
    total_tf = total.text_frame
    total_tf.text = "Total Annual Benefit: $3,146,200"
    total_tf.paragraphs[0].font.size = Pt(24)
    total_tf.paragraphs[0].font.bold = True
    total_tf.paragraphs[0].font.color.rgb = WALMART_BLUE
    total_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Slide 14: Activity Hub Investment Breakdown
    print("  Adding Slide 14: Activity Hub Investment Breakdown")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shapes = slide.shapes
    
    title_box = shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_tf = title_box.text_frame
    title_tf.text = "$55K Development Leverages $129K Infrastructure"
    title_tf.paragraphs[0].font.size = Pt(36)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = WALMART_BLUE
    
    # Left: Activity Hub Development
    left_box = shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(4.2), Inches(3))
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = RGBColor(240, 255, 240)
    left_tf = left_box.text_frame
    left_tf.text = "Activity Hub Development Only: $55,000"
    left_tf.paragraphs[0].font.size = Pt(18)
    left_tf.paragraphs[0].font.bold = True
    left_tf.paragraphs[0].font.color.rgb = GREEN
    
    dev_items = [
        "Integration Development: $35,000 (64%)",
        "  • Core API development",
        "  • 16 platform connectors",
        "  • Authentication & data model",
        "",
        "Data Migration: $20,000 (36%)",
        "  • Legacy data extraction",
        "  • Cleanup & standardization",
        "  • Schema transformation"
    ]
    
    for item in dev_items:
        p = left_tf.add_paragraph()
        p.text = item
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(2)
    
    # Right: Leverages Path to Production
    right_box = shapes.add_textbox(Inches(5.2), Inches(1.8), Inches(4.2), Inches(3))
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = RGBColor(240, 240, 255)
    right_tf = right_box.text_frame
    right_tf.text = "Leverages Path to Production (Already Paid):"
    right_tf.paragraphs[0].font.size = Pt(16)
    right_tf.paragraphs[0].font.bold = True
    right_tf.paragraphs[0].font.color.rgb = WALMART_BLUE
    
    infra_items = [
        "Azure infrastructure: $45K ✓",
        "Security & compliance: $15K ✓",
        "Testing & QA: $25K ✓",
        "Training & docs: $20K ✓",
        "Contingency: $24K ✓"
    ]
    
    for item in infra_items:
        p = right_tf.add_paragraph()
        p.text = item
        p.font.size = Pt(15)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(6)
    
    # Bottom summary
    summary_box = shapes.add_textbox(Inches(1.5), Inches(5.2), Inches(7), Inches(1.5))
    summary_box.fill.solid()
    summary_box.fill.fore_color.rgb = WALMART_YELLOW
    summary_tf = summary_box.text_frame
    summary_tf.text = "Total Activity Hub Cost: $55K + $129K = $184K\nRecurring Costs: $0 (covered by Path to Production)"
    summary_tf.paragraphs[0].font.size = Pt(20)
    summary_tf.paragraphs[0].font.bold = True
    summary_tf.paragraphs[0].font.color.rgb = DARK_GRAY
    summary_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Slide 15: Activity Hub ROI
    print("  Adding Slide 15: Activity Hub ROI - 21-Day Payback")
    add_table_slide(
        prs,
        "Activity Hub ROI: 21-Day Payback Period",
        ["Metric", "Amount"],
        [
            ["Total Investment", "$184,000 (one-time)"],
            ["Annual Benefit", "$3,146,200"],
            ["Year 1 Net Benefit", "$2,962,200"],
            ["Year 1 ROI", "1,610%"],
            ["", ""],
            ["3-Year Total Benefit", "$9,919,200"],
            ["3-Year Net Benefit", "$9,735,200"],
            ["3-Year ROI", "5,292%"],
            ["", ""],
            ["Break-Even", "Month 1, Week 3"],
            ["Payback Period", "21 days"],
            ["NPV (8% discount)", "$8,960,000"]
        ],
        highlight_row=8
    )
    
    # Add callout
    slide = prs.slides[-1]
    add_callout_box(slide, "Every dollar invested returns $53 over 3 years",
                   2, 6.5, 6, 0.6, WALMART_YELLOW)
    
    # Slide 16: Activity Hub Timeline
    print("  Adding Slide 16: Activity Hub Timeline")
    add_table_slide(
        prs,
        "Activity Hub: 8 Months to Full Deployment",
        ["Phase", "Timeline", "Activities", "Milestone"],
        [
            ["Foundation", "Months 1-3", "Requirements, data model, API architecture", "Design complete"],
            ["Core Dev", "Months 3-5", "API dev, authentication, 6 core integrations", "Core platform functional"],
            ["Integration", "Months 5-7", "10 remaining integrations, data migration, testing", "All 16 platforms connected"],
            ["Launch", "Month 8", "Training, deployment, phased onboarding", "Activity Hub live"],
            ["Optimize", "Month 8+", "Monitor, gather feedback, continuous improvement", "Full suite operational"]
        ]
    )
    
    # COMBINED SUMMARY SECTION
    print("  Adding Slide 16: Section Title - Combined Summary")
    add_section_title(prs, "Combined Investment\n& ROI")
    
    # Slide 17: Combined Investment Summary
    print("  Adding Slide 17: Combined Investment Summary")
    add_table_slide(
        prs,
        "Total 3-Year Investment: $464,000",
        ["Component", "One-Time", "Year 1 Recurring", "Year 2 Recurring", "Year 3 Recurring", "3-Year Total"],
        [
            ["Path to Production", "$129,000", "$90,000", "$92,500", "$97,500", "$409,000"],
            ["Activity Hub", "$55,000", "$0", "$0", "$0", "$55,000"],
            ["TOTAL", "$184,000", "$90,000", "$92,500", "$97,500", "$464,000"]
        ],
        highlight_row=3
    )
    
    # Slide 19: Combined ROI
    print("  Adding Slide 19: Combined ROI - $25.6M Value")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shapes = slide.shapes
    
    title_box = shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_tf = title_box.text_frame
    title_tf.text = "$25.6M Total Value Over 3 Years"
    title_tf.paragraphs[0].font.size = Pt(40)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = WALMART_BLUE
    
    # Waterfall visualization (simplified as text)
    waterfall_data = [
        ("Investment Base", "-$464,000", RED, 1.8),
        ("Activity Hub Direct Benefits", "+$9,919,200", GREEN, 2.8),
        ("Future Products Enabled", "+$15,700,000", GREEN, 3.8),
        ("Total Value", "$25,619,200", WALMART_BLUE, 4.8),
        ("Net Benefit", "$25,155,200", WALMART_BLUE, 5.8),
        ("Combined ROI", "5,420%", WALMART_YELLOW, 6.8)
    ]
    
    for label, amount, color, top in waterfall_data:
        box = shapes.add_textbox(Inches(1.5), Inches(top), Inches(7), Inches(0.7))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(245, 245, 245) if color != WALMART_YELLOW else color
        tf = box.text_frame
        tf.text = f"{label}: {amount}"
        tf.paragraphs[0].font.size = Pt(22) if "ROI" in label or "Total" in label or "Net" in label else Pt(18)
        tf.paragraphs[0].font.bold = "ROI" in label or "Total" in label or "Net" in label
        tf.paragraphs[0].font.color.rgb = color if color != WALMART_YELLOW else DARK_GRAY
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Slide 20: Risk Assessment
    print("  Adding Slide 20: Risk Assessment")
    add_table_slide(
        prs,
        "Risks Are Manageable with Clear Mitigation",
        ["Risk", "Impact", "Probability", "Mitigation"],
        [
            ["User Adoption", "High", "Medium", "Phased rollout, training, exec sponsorship"],
            ["Integration Complexity", "Medium", "Medium", "API-first design, thorough testing"],
            ["Security/Compliance", "High", "Low", "Follow Production_Path, security audits"],
            ["Data Quality", "Medium", "Medium", "Validation rules, migration testing"],
            ["Timeline Delays", "Medium", "Medium", "15% contingency buffer, agile approach"]
        ]
    )
    
    # Slide 21: Success Metrics
    print("  Adding Slide 21: Success Metrics & KPIs")
    add_two_column_slide(
        prs,
        "How We'll Measure Success",
        {
            'title': 'Path to Production Metrics',
            'bullets': [
                '✅ APM Certified status achieved',
                '✅ SSP Approved status achieved',
                '✅ CI/CD pipeline operational',
                '✅ <2 week deployment time for new products',
                '✅ Zero security incidents'
            ]
        },
        {
            'title': 'Activity Hub Metrics',
            'bullets': [
                '✅ 16 platforms integrated',
                '✅ 100+ users onboarded',
                '✅ 275 → 58 hours/week (79% reduction)',
                '✅ 30% → 5% message overlap',
                '✅ $3.1M+ annual benefits achieved',
                '✅ 95%+ user satisfaction score'
            ]
        }
    )
    
    # Slide 22: The Ask
    print("  Adding Slide 22: The Ask - Two Approval Decisions")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shapes = slide.shapes
    
    title_box = shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_tf = title_box.text_frame
    title_tf.text = "What We Need from Leadership Today"
    title_tf.paragraphs[0].font.size = Pt(40)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = WALMART_BLUE
    
    # Two approval boxes
    approval1_box = shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.4), Inches(1.8))
    approval1_box.fill.solid()
    approval1_box.fill.fore_color.rgb = RGBColor(240, 245, 255)
    a1_tf = approval1_box.text_frame
    a1_tf.text = "☐ Approval #1: Path to Production Infrastructure"
    a1_tf.paragraphs[0].font.size = Pt(22)
    a1_tf.paragraphs[0].font.bold = True
    a1_tf.paragraphs[0].font.color.rgb = WALMART_BLUE
    
    approval1_details = [
        "Investment: $129,000 one-time + $90,000/year",
        "Deliverable: Enterprise development and deployment capability",
        "Timeline: 5 months to operational",
        "Value: Enables ALL future Store Support products"
    ]
    
    for detail in approval1_details:
        p = a1_tf.add_paragraph()
        p.text = f"  • {detail}"
        p.font.size = Pt(15)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(4)
    
    approval2_box = shapes.add_textbox(Inches(0.8), Inches(3.8), Inches(8.4), Inches(1.8))
    approval2_box.fill.solid()
    approval2_box.fill.fore_color.rgb = RGBColor(240, 255, 240)
    a2_tf = approval2_box.text_frame
    a2_tf.text = "☐ Approval #2: Activity Hub Platform"
    a2_tf.paragraphs[0].font.size = Pt(22)
    a2_tf.paragraphs[0].font.bold = True
    a2_tf.paragraphs[0].font.color.rgb = GREEN
    
    approval2_details = [
        "Investment: $55,000 development (one-time)",
        "Deliverable: Integrated platform unifying 16 tools",
        "Timeline: 8 months to production launch",
        "Value: $9.9M over 3 years, 5,292% ROI"
    ]
    
    for detail in approval2_details:
        p = a2_tf.add_paragraph()
        p.text = f"  • {detail}"
        p.font.size = Pt(15)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(4)
    
    # Combined package callout
    combined_box = shapes.add_textbox(Inches(1.5), Inches(5.8), Inches(7), Inches(1.3))
    combined_box.fill.solid()
    combined_box.fill.fore_color.rgb = WALMART_YELLOW
    c_tf = combined_box.text_frame
    c_tf.text = "☐ Combined Package:"
    c_tf.paragraphs[0].font.size = Pt(20)
    c_tf.paragraphs[0].font.bold = True
    c_tf.paragraphs[0].font.color.rgb = DARK_GRAY
    
    combined_details = [
        "Total Investment: $464,000 over 3 years | Total Value: $25.6M | ROI: 5,420%"
    ]
    
    for detail in combined_details:
        p = c_tf.add_paragraph()
        p.text = detail
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = DARK_GRAY
        p.alignment = PP_ALIGN.CENTER
    
    # Slide 23: Why This, Why Now?
    print("  Adding Slide 23: Why This, Why Now?")
    add_two_column_slide(
        prs,
        "The Opportunity Cost of Waiting",
        {
            'title': 'Scenario 1: Approve Today',
            'bullets': [
                'Launch: August 2026',
                'Year 1 benefit: $3.1M',
                '3-year benefit: $9.9M',
                'Future products: 15+ enabled'
            ]
        },
        {
            'title': 'Scenario 2: Delay 6 Months',
            'bullets': [
                'Launch: February 2027',
                'Year 1 benefit: $1.6M (partial year)',
                'Lost value: $1.5M',
                'Future products: Delayed roadmap'
            ]
        },
        GREEN,
        RED
    )
    
    # Add bottom callout
    slide = prs.slides[-1]
    add_callout_box(slide, "Every month of delay costs $262,000 in unrealized benefits",
                   2, 6.8, 6, 0.6, RED)
    
    # Slide 24: Call to Action
    print("  Adding Slide 24: Call to Action")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WALMART_BLUE
    shapes = slide.shapes
    
    title_box = shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1))
    title_tf = title_box.text_frame
    title_tf.text = "Ready to Transform Store Support?"
    title_tf.paragraphs[0].font.size = Pt(48)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = WHITE
    title_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    decision_box = shapes.add_textbox(Inches(1.5), Inches(3), Inches(7), Inches(2))
    decision_tf = decision_box.text_frame
    decision_tf.text = "Today's Decision:"
    decision_tf.paragraphs[0].font.size = Pt(32)
    decision_tf.paragraphs[0].font.bold = True
    decision_tf.paragraphs[0].font.color.rgb = WALMART_YELLOW
    decision_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    decisions = [
        "☑️ Approve Path to Production Infrastructure",
        "☑️ Approve Activity Hub Platform"
    ]
    
    for decision in decisions:
        p = decision_tf.add_paragraph()
        p.text = decision
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(12)
    
    stats_box = shapes.add_textbox(Inches(1.5), Inches(5.2), Inches(7), Inches(1.2))
    stats_tf = stats_box.text_frame
    stats_tf.text = "$464,000 investment | $25.6M value | 5,420% ROI"
    stats_tf.paragraphs[0].font.size = Pt(24)
    stats_tf.paragraphs[0].font.color.rgb = WALMART_YELLOW
    stats_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    p = stats_tf.add_paragraph()
    p.text = "Target Production Launch: August 2026"
    p.font.size = Pt(20)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    cta_box = shapes.add_textbox(Inches(2), Inches(6.5), Inches(6), Inches(0.8))
    cta_tf = cta_box.text_frame
    cta_tf.text = "Let's Build the Future of Store Support"
    cta_tf.paragraphs[0].font.size = Pt(36)
    cta_tf.paragraphs[0].font.bold = True
    cta_tf.paragraphs[0].font.color.rgb = WHITE
    cta_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    return prs


def main():
    """Main execution function."""
    print("\n" + "="*60)
    print("Store Support Path To Production PowerPoint Generator")
    print("Path to Production + Activity Hub")
    print("="*60 + "\n")
    
    prs = create_presentation()
    
    output_file = "Store_Support_Path_To_Production.pptx"
    prs.save(output_file)
    
    print("\n" + "="*60)
    print(f"✓ Presentation created successfully!")
    print(f"✓ File saved as: {output_file}")
    print(f"✓ Total slides: {len(prs.slides)}")
    print("="*60 + "\n")
    print("Next steps:")
    print("1. Open the presentation in PowerPoint")
    print("2. Add Walmart logo/branding as needed")
    print("3. Review and customize for your specific audience")
    print("4. Practice the delivery with speaker notes from outline")
    print("\n")


if __name__ == "__main__":
    main()
